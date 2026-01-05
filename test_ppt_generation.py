#!/usr/bin/env python3
"""Script de test pour la génération de PowerPoint"""

import pandas as pd
import os
import tempfile
from io import BytesIO
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from PIL import Image, ImageDraw, ImageFont
import base64

def generate_table_image():
    """Génère une image du tableau avec PIL"""
    print("📊 Génération de l'image du tableau...")
    
    # Créer une image
    width, height = 1200, 400
    img = Image.new('RGB', (width, height), 'white')
    draw = ImageDraw.Draw(img)
    
    # Dessiner un tableau simple
    # Header
    draw.rectangle([0, 0, width, 50], fill='#9C182F')
    headers = ['Entité', 'Postes avant', 'Nouveaux', 'Pourvus', 'En cours']
    col_width = width // len(headers)
    
    try:
        font = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf", 16)
        font_data = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", 14)
    except:
        font = ImageFont.load_default()
        font_data = ImageFont.load_default()
    
    # Headers
    for i, header in enumerate(headers):
        x = i * col_width + col_width // 2
        draw.text((x, 20), header, fill='white', font=font, anchor='mm')
    
    # Données
    data = [
        ('TGCC', '10', '5', '2', '13'),
        ('TGEM', '8', '3', '1', '10'),
        ('TG ALU', '5', '2', '0', '7'),
    ]
    
    y_offset = 70
    for row in data:
        for i, val in enumerate(row):
            x = i * col_width + col_width // 2
            draw.text((x, y_offset), val, fill='black', font=font_data, anchor='mm')
        draw.line([(0, y_offset + 25), (width, y_offset + 25)], fill='#ddd', width=1)
        y_offset += 50
    
    # Ligne TOTAL
    draw.rectangle([0, y_offset - 25, width, y_offset + 25], fill='#9C182F')
    total_data = ('TOTAL', '23', '10', '3', '30')
    for i, val in enumerate(total_data):
        x = i * col_width + col_width // 2
        draw.text((x, y_offset), val, fill='white', font=font, anchor='mm')
    
    # Sauvegarder
    output_path = os.path.join(tempfile.gettempdir(), 'table.png')
    img.save(output_path)
    print(f"✅ Image générée: {output_path}")
    return output_path

def generate_kanban_image():
    """Génère une image du Kanban avec PIL"""
    print("📊 Génération de l'image du Kanban...")
    
    # Créer une image
    width, height = 1400, 600
    img = Image.new('RGB', (width, height), 'white')
    draw = ImageDraw.Draw(img)
    
    try:
        font_header = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf", 16)
        font_card = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", 12)
    except:
        font_header = ImageFont.load_default()
        font_card = ImageFont.load_default()
    
    # Données du Kanban
    columns = {
        'Sourcing': ['Ingénieur BTP', 'Électricien', 'Maçon'],
        'Shortlisté': ['Chef de projet', 'Conducteur'],
        'Signature DRH': ['Ingénieur QSE'],
        'Clôture': ['Technicien', 'Assistant'],
        'Désistement': ['Chauffeur']
    }
    
    col_width = width // len(columns)
    x_offset = 10
    
    for col_name, cards in columns.items():
        # Header de colonne
        draw.rectangle([x_offset, 10, x_offset + col_width - 20, 50], fill='#9C182F', outline='#9C182F')
        draw.text((x_offset + col_width // 2 - 10, 30), f"{col_name} ({len(cards)})", 
                 fill='white', font=font_header, anchor='mm')
        
        # Cartes
        y_offset = 70
        for card in cards:
            draw.rectangle([x_offset + 5, y_offset, x_offset + col_width - 25, y_offset + 40], 
                          fill='#f0f0f0', outline='#ddd', width=2)
            draw.text((x_offset + 10, y_offset + 20), card, fill='black', font=font_card, anchor='lm')
            y_offset += 50
        
        x_offset += col_width
    
    # Sauvegarder
    output_path = os.path.join(tempfile.gettempdir(), 'kanban.png')
    img.save(output_path)
    print(f"✅ Image générée: {output_path}")
    return output_path

def generate_powerpoint():
    """Génère le PowerPoint avec les images"""
    print("📝 Génération du PowerPoint...")
    
    # Générer les images
    table_img = generate_table_image()
    kanban_img = generate_kanban_image()
    
    # Charger le template
    template_path = '/workspaces/TG_Hire/MASQUE PPT TGCC (2).pptx'
    prs = Presentation(template_path)
    
    print(f"✅ Template chargé: {len(prs.slides)} slides")
    
    # Parcourir les slides
    for slide_idx, slide in enumerate(prs.slides):
        shapes_to_remove = []
        
        for shape in slide.shapes:
            # Pylance-safe text extraction: some shapes don't expose `.text`
            shape_text = getattr(shape, "text", None)
            if shape_text is None or not str(shape_text).strip():
                # Fallback to text_frame when available
                has_tf = getattr(shape, "has_text_frame", False)
                if has_tf:
                    tf = getattr(shape, "text_frame", None)
                    if tf and getattr(tf, "text", None):
                        shape_text = tf.text
                    else:
                        continue
                else:
                    continue
            
            # Tableau des entités
            if "{{TABLEAU_BESOINS_ENTITES}}" in shape_text:
                print(f"  - Slide {slide_idx + 1}: Insertion du tableau")
                left = shape.left
                top = shape.top
                width = shape.width
                height = shape.height
                shapes_to_remove.append(shape)
                
                if os.path.exists(table_img):
                    slide.shapes.add_picture(table_img, left, top, width=width, height=height)
            
            # Kanban des postes
            elif "{{METRIC_TOTAL_POSTES}}" in shape_text:
                print(f"  - Slide {slide_idx + 1}: Insertion du Kanban")
                left = shape.left
                top = shape.top
                width = shape.width
                height = shape.height
                shapes_to_remove.append(shape)
                
                if os.path.exists(kanban_img):
                    slide.shapes.add_picture(kanban_img, left, top, width=width, height=height)
        
        # Supprimer les shapes marqués
        for shape in shapes_to_remove:
            sp = shape.element
            sp.getparent().remove(sp)
    
    # Sauvegarder
    output_path = '/tmp/test_report.pptx'
    prs.save(output_path)
    print(f"✅ PowerPoint sauvegardé: {output_path}")
    
    # Nettoyer les images temporaires
    for img_path in [table_img, kanban_img]:
        if os.path.exists(img_path):
            os.remove(img_path)
    
    return output_path

if __name__ == "__main__":
    print("🚀 Test de génération PowerPoint")
    print("=" * 50)
    
    try:
        ppt_path = generate_powerpoint()
        print(f"\n✅ Test réussi!")
        print(f"📄 Fichier généré: {ppt_path}")
    except Exception as e:
        print(f"\n❌ Erreur: {e}")
        import traceback
        traceback.print_exc()
