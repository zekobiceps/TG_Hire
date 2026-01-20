import streamlit as st
import pandas as pd
import plotly.express as px
import plotly.graph_objects as go
from plotly.subplots import make_subplots
import numpy as np
from datetime import datetime, timedelta
import warnings
import os
import base64
warnings.filterwarnings('ignore')
import re
import streamlit.components.v1 as components
import unicodedata
from io import BytesIO
import json
import gspread
from google.oauth2 import service_account
import unicodedata
from pptx import Presentation
def _format_long_title(title: str, max_line_length: int = 20) -> str:
    """Insert spaces in very long titles without natural breaks to force wrapping.
    - If the title already contains spaces, return as-is (CSS will wrap).
    - If no spaces and length > max_line_length, insert spaces every max_line_length chars.
    """
    if title is None:
        return ""
    s = str(title).strip()
    if len(s) <= max_line_length:
        return s
    if " " in s:
        return s
    chunks = [s[i:i+max_line_length] for i in range(0, len(s), max_line_length)]
    return " ".join(chunks)
def smart_wrap_title(title, max_line_length=25):
    """Retourne un titre avec des balises <br> aux bons endroits.
    - Coupe aux espaces si possible pour respecter max_line_length
    - Gère les titres sans espaces (souvent en MAJUSCULES) via regex
    - Dernier recours: coupe tous les max_line_length caractères
    """
    if not isinstance(title, str):
        return title

    s = title.strip()
    if len(s) <= max_line_length:
        return s

    # Si le titre contient des espaces, essayer de couper aux espaces
    if " " in s:
        words = s.split()
        lines = []
        current = ""
        for w in words:
            sep = (1 if current else 0)
            if len(current) + sep + len(w) <= max_line_length:
                current = (current + (" " if current else "") + w)
            else:
                if current:
                    lines.append(current)
                current = w
        if current:
            lines.append(current)
        return "<br>".join(lines)

    # Pour les titres sans espaces (majuscules, acronymes, chiffres)
    try:
        parts = re.findall(r"[A-Z][a-z]+|[A-Z]{2,}|[a-z]+|\d+", s)
    except Exception:
        parts = []

    if len(parts) > 1:
        spaced = " ".join(parts)
        return smart_wrap_title(spaced, max_line_length)

    # Dernier recours: couper tous les max_line_length caractères
    chunks = [s[i:i+max_line_length] for i in range(0, len(s), max_line_length)]
    return "<br>".join(chunks)
from pptx.util import Inches, Pt
from PIL import Image
import io

def _normalize_text(text):
    """A global function to safely normalize text, handling None and NaN values."""
    if text is None or (isinstance(text, float) and np.isnan(text)):
        return ''
    s = str(text)
    s = unicodedata.normalize('NFKD', s)
    s = ''.join(ch for ch in s if not unicodedata.combining(ch))
    return s.lower().strip()

def _norm(x):
    """Robust text normalization for status/keywords matching."""
    return _normalize_text(x)

st.set_page_config(
    page_title="📊 Reporting RH Complet",
    page_icon="📊",
    layout="wide"
)

# Vérification de la connexion
if not st.session_state.get("logged_in", False):
    st.stop()


def _truncate_label(label: str, max_len: int = 20) -> str:
    """Truncate long labels to a max length and append ellipsis.

    Returns a truncated string (if needed). Keep a plain truncation without
    changing accents/characters. Default max_len=20 (adjustable).
    """
    if not isinstance(label, str):
        return label
    if len(label) <= max_len:
        return label
    return label[: max_len - 4].rstrip() + '....'


def render_generic_metrics(metrics):
    """Render a horizontal row of metric cards via HTML.
    metrics: list of tuples (title, value, color_hex)
    """
    css = """
    <style>
    .gen-kpi-row{display:flex;gap:18px;justify-content:center;align-items:stretch;margin-bottom:8px}
    .gen-kpi{background:#fff;border-radius:8px;padding:14px 18px;min-width:220px;flex:0 1 auto;border:1px solid #e6eef6;box-shadow:0 2px 6px rgba(0,0,0,0.04);display:flex;flex-direction:column;align-items:center;justify-content:center;text-align:center}
    .gen-kpi .t{font-size:17px;color:#2c3e50;margin-bottom:8px;font-weight:700;text-align:center}
    .gen-kpi .v{font-size:36px;color:#172b4d;font-weight:800;text-align:center}
    </style>
    """
    cards = []
    for title, value, color in metrics:
        cards.append(f"<div class='gen-kpi'><div class='t'>{title}</div><div class='v' style='color:{color};'>{value}</div></div>")
    html = css + "<div class='gen-kpi-row'>" + "".join(cards) + "</div>"
    return html


# Shared title font used for all main charts so typography is consistent
TITLE_FONT = dict(family="Arial, sans-serif", size=18, color="#111111", )

# --- Centralisation des Logos et Affichage Entités ---
ENTITY_LOGO_MAP = {
    'TGCC': 'TGCC.PNG',
    'TGEM': 'TGEM.PNG',
    'TG ALU': 'TG ALU.PNG',
    'TG COVER': 'TG COVER.PNG',
    'TG STEEL': 'TG STEEL.PNG',
    'TG STONE': 'TG STONE.PNG',
    'TG WOOD': 'TG WOOD.PNG',
    'STAM': 'STAM.png',
    'BFO': 'BFO.png',
    'TGCC IMMOBILIER': 'tgcc-immobilier.png',
    'TGCC-IMMOBILIER': 'tgcc-immobilier.png'
}

@st.cache_data
def load_all_logos_b64():
    """Charge tous les logos disponibles en base64."""
    logos_dict = {}
    current_dir = os.path.dirname(os.path.abspath(__file__))
    root_dir = os.path.dirname(current_dir)
    logo_dir = os.path.join(root_dir, "LOGO")
    
    if os.path.exists(logo_dir):
        for filename in os.listdir(logo_dir):
            if filename.lower().endswith(('.png', '.jpg', '.jpeg')):
                try:
                    path = os.path.join(logo_dir, filename)
                    with open(path, "rb") as f:
                        logos_dict[filename] = base64.b64encode(f.read()).decode()
                except Exception:
                    pass
    return logos_dict

def get_entity_display_html_with_logo(name, logos_dict):
    """Retourne le HTML (Logo + Suffixe éventuel) pour l'affichage avec logo."""
    if not name or pd.isna(name): return ""
    name_str = str(name).strip()
    name_upper = name_str.upper()
    
    logo_file = None
    if name_upper in ENTITY_LOGO_MAP:
        logo_file = ENTITY_LOGO_MAP[name_upper]
    
    if not logo_file:
        for filename in logos_dict.keys():
            if os.path.splitext(filename)[0].upper() == name_upper:
                logo_file = filename
                break
                
    if not logo_file:
        for key in sorted(ENTITY_LOGO_MAP.keys(), key=len, reverse=True):
            if key in name_upper:
                logo_file = ENTITY_LOGO_MAP[key]
                break

    if logo_file and logo_file in logos_dict:
        logo_b64_str = logos_dict[logo_file]
        if name_upper in ['TG STEEL', 'TG STONE'] or 'IMMOBILIER' in name_upper:
            logo_height = 50
        elif name_upper == 'BFO':
            logo_height = 80
        else:
            logo_height = 63
            
        img_tag = f'<img src="data:image/png;base64,{logo_b64_str}" height="{logo_height}" style="vertical-align: middle; display: block; margin: 0 auto;">'
        
        matched_key = next((k for k in sorted(ENTITY_LOGO_MAP.keys(), key=len, reverse=True) if k in name_upper), None)
        if matched_key:
            pattern = re.escape(matched_key) + r'\s*-\s*(.*)'
            match = re.search(pattern, name_str, flags=re.IGNORECASE)
            if match and match.group(1):
                return f'{img_tag} <span style="vertical-align: middle; margin-top: 5px; font-weight: bold; display: block; text-align: center;">{match.group(1)}</span>'
        
        return img_tag
    
    return f'<div style="text-align: center; font-weight: 600;">{name_str}</div>'

def apply_title_style(fig):
    """Applique la police et le style de titre standardisé à une figure Plotly."""
    try:
        fig.update_layout(title_font=TITLE_FONT)
    except Exception:
        try:
            current = ''
            if hasattr(fig.layout, 'title') and getattr(fig.layout.title, 'text', None):
                current = fig.layout.title.text
            fig.update_layout(title=dict(text=current, x=0, xanchor='left', font=TITLE_FONT))
        except Exception:
            pass
    # Increase generic data-label fontsize and legend font for better readability
    try:
        fig.update_traces(textfont=dict(size=15))
    except Exception:
        pass
    try:
        if hasattr(fig.layout, 'legend'):
            fig.update_layout(legend=dict(font=dict(size=13)))
    except Exception:
        pass
    return fig


def get_current_commit_hash(short: bool = True) -> str:
    """Return the current git commit hash (short by default).

    Tries subprocess git first, then falls back to reading .git/HEAD or
    environment variable GIT_COMMIT. Returns 'unknown' if not available.
    """
    import subprocess
    try:
        if short:
            out = subprocess.check_output(["git", "rev-parse", "--short", "HEAD"], stderr=subprocess.DEVNULL)
        else:
            out = subprocess.check_output(["git", "rev-parse", "HEAD"], stderr=subprocess.DEVNULL)
        return out.decode().strip()
    except Exception:
        # Try environment variable
        import os
        env = os.environ.get('GIT_COMMIT') or os.environ.get('COMMIT_SHA') or os.environ.get('GITHUB_SHA')
        if env:
            return env[:7] if short else env
        # Try minimal .git reading
        try:
            git_head = None
            head_path = os.path.join(os.path.dirname(__file__), '..', '.git', 'HEAD')
            head_path = os.path.abspath(head_path)
            if os.path.exists(head_path):
                with open(head_path, 'r') as f:
                    ref = f.read().strip()
                if ref.startswith('ref:'):
                    ref_path = ref.split(' ', 1)[1]
                    ref_file = os.path.join(os.path.dirname(head_path), ref_path)
                    if os.path.exists(ref_file):
                        with open(ref_file, 'r') as rf:
                            val = rf.read().strip()
                            return val[:7] if short else val
                else:
                    return ref[:7] if short else ref
        except Exception:
            pass
    return 'unknown'


def _parse_mixed_dates(series):
    """Parse a pandas Series that may contain mixed date representations.

    Strategy:
      1. If values are numeric (Excel serial), convert using Excel epoch.
      2. Try pd.to_datetime(..., dayfirst=True) to favor dd/mm/YYYY formats.
      3. Fallback to pd.to_datetime(..., errors='coerce') for other formats.

    Returns a datetime64[ns] Series with NaT for unparseable values.
    """
    s = series.copy()
    try:
        # If the series is numeric (Excel serials), convert per-element where it looks numeric
        if pd.api.types.is_numeric_dtype(s):
            def _maybe_excel(x):
                try:
                    xf = float(x)
                    return pd.Timestamp('1899-12-30') + pd.Timedelta(days=xf)
                except Exception:
                    return pd.NaT

            return s.apply(lambda v: _maybe_excel(v) if pd.notna(v) and str(v).strip().replace('.', '', 1).isdigit() else pd.NaT).combine_first(pd.to_datetime(s, dayfirst=True, errors='coerce'))
    except Exception:
        # fall through to permissive parsing below
        pass

    # First try dayfirst parsing (dd/mm/YYYY common in French contexts)
    parsed = pd.to_datetime(s, dayfirst=True, errors='coerce')
    # If many values still NaT, try fallback parsing
    if parsed.isna().sum() > len(parsed) * 0.25:
        parsed_alt = pd.to_datetime(s, errors='coerce')
        parsed = parsed.combine_first(parsed_alt)

    return parsed


def render_kpi_cards(recrutements, postes, directions, delai_display, delai_help=None):
    """Render a single-row set of KPI cards (inline, bordered with colored left stripe).

    Cards: [Nombre de recrutements] [Postes concernés] [Directions concernées] [Délai moyen]
    Returns an HTML string ready to be inserted with st.markdown(..., unsafe_allow_html=True)
    """

    css = """
<style>
.kpi-row{display:flex;gap:16px;flex-wrap:nowrap;align-items:stretch;margin-bottom:14px;justify-content:center}
.kpi-card{flex:1 1 0;background:#fff;border-radius:8px;padding:14px;display:flex;flex-direction:column;align-items:center;justify-content:center;border:1px solid #e6eef6;min-width:180px;text-align:center}
.kpi-card .title{font-size:14px;color:#2c3e50;margin-bottom:8px;font-weight:700;text-align:center}
.kpi-card .value{font-size:26px;font-weight:700;color:#172b4d;text-align:center}
.kpi-accent{border-left:6px solid #1f77b4}
.kpi-green{border-left-color:#2ca02c}
.kpi-orange{border-left-color:#ff7f0e}
.kpi-purple{border-left-color:#6f42c1}
.kpi-help{font-size:12px;color:#555;margin-top:8px;text-align:center}
@media(max-width:800px){.kpi-row{flex-direction:column}}
</style>
"""

    html = f"""
{css}
<div class='kpi-row'>
    <div class='kpi-card kpi-accent' style='flex:2'>
        <div class='title'>Nombre de recrutements</div>
        <div class='value'>{recrutements:,}</div>
    </div>
    <div class='kpi-card kpi-green'>
        <div class='title'>Postes concernés</div>
        <div class='value'>{postes:,}</div>
    </div>
    <div class='kpi-card kpi-orange'>
        <div class='title'>Directions concernées</div>
        <div class='value'>{directions:,}</div>
    </div>
    <div class='kpi-card kpi-purple'>
        <div class='title'>Délai moyen (jours)</div>
        <div class='value'>{delai_display}</div>
        <div class='kpi-help'>{delai_help or ''}</div>
    </div>
</div>
"""

    return html

    df2 = df.copy()
    df2[start_col] = pd.to_datetime(df2[start_col], errors='coerce')
    df2[end_col] = pd.to_datetime(df2[end_col], errors='coerce')

    # Optionally filter to a specific status (e.g., 'Clôture') if the column exists
    if status_col in df2.columns and status_value is not None:
        df2 = df2[df2[status_col] == status_value].copy()

    # Compute delta in days
    df2['time_to_hire_days'] = (df2[end_col] - df2[start_col]).dt.days

    if drop_negative:
        df2 = df2[df2['time_to_hire_days'].notna() & (df2['time_to_hire_days'] >= 0)].copy()
    else:
        df2 = df2[df2['time_to_hire_days'].notna()].copy()

    if df2.empty:
        overall = None
    else:
        overall = {
            'mean': float(df2['time_to_hire_days'].mean()),
            'median': float(df2['time_to_hire_days'].median()),
            'std': float(df2['time_to_hire_days'].std()),
            'count': int(df2['time_to_hire_days'].count())
        }

    # Grouped aggregates
    by_direction = None
    by_poste = None
    if 'Direction concernée' in df2.columns:
        by_direction = df2.groupby('Direction concernée')['time_to_hire_days'].agg(['count', 'mean', 'median', 'std']).reset_index().sort_values('mean')
    if 'Poste demandé' in df2.columns:
        by_poste = df2.groupby('Poste demandé')['time_to_hire_days'].agg(['count', 'mean', 'median', 'std']).reset_index().sort_values('mean')

    return {
        'start_col': start_col,
        'end_col': end_col,
        'overall': overall,
        'by_direction': by_direction,
        'by_poste': by_poste,
        'df': df2
    }


# CSS pour styliser le bouton Google Sheets en rouge vif
st.markdown("""
<style>
.stButton > button[title="Synchroniser les données depuis Google Sheets"] {
    background-color: #FF4B4B !important;
    color: white !important;
    border: none !important;
    border-radius: 5px !important;
    font-weight: bold !important;
}
.stButton > button[title="Synchroniser les données depuis Google Sheets"]:hover {
    background-color: #FF3333 !important;
    color: white !important;
}
</style>
""", unsafe_allow_html=True)

# Global tweaks: enlarge tab labels, plot titles and legend fonts for better readability
st.markdown("""
<style>
/* Tabs: increase font-size and weight (stronger selectors) */
div[role="tablist"] > button, button[role="tab"], div[role="tablist"] button[role="tab"] {
    font-size: 20px !important;
    font-weight: 700 !important;
    padding: 8px 14px !important;
}

/* Plotly title and legend sizing (fallback selector) */
.plotly .gtitle, .plotly .gtitle text { font-family: Arial, sans-serif !important; font-size: 18px !important; fill: #111 !important; }
.plotly .legend { font-size: 14px !important; }

/* Streamlit metrics: aggressive selectors to increase label and value sizes */
/* Streamlit metrics: aggressive selectors to increase label and value sizes and center them */
div[data-testid="metric-container"] div[data-testid="stMetric"] span[data-testid], div[data-testid="metric-container"] .stMetricValue, .stMetricValue {
    font-size: 28px !important; font-weight: 800 !important; line-height:1 !important; display:block; text-align:center !important;
}
div[data-testid="metric-container"] div[data-testid="stMetric"] p, .stMetricLabel {
    font-size: 15px !important; color: #2c3e50 !important; margin:0 !important; display:block; text-align:center !important;
}

/* Fallback generic selectors for metrics */
span[data-testid="stMetricLabel"] { font-size:15px !important; }
span[data-testid="stMetricValue"] { font-size:28px !important; font-weight:700 !important; }

/* Slight increase for general section subtitles */
.stSubheader, .stMarkdown h3 { font-size: 1.06em !important; }
</style>
""", unsafe_allow_html=True)

# Données pour le Kanban
postes_data = [
    # Colonne Nouvelle demande
    {"titre": "Nouvelle Demande Exemple", "entite": "TGCC", "lieu": "SIEGE", "demandeur": "X.DEMANDE", "recruteur": "Jalal", "statut": "Nouvelle demande"},
    # Colonne Sourcing
    {"titre": "Ingénieur Achat", "entite": "TGCC", "lieu": "SIEGE", "demandeur": "A.BOUZOUBAA", "recruteur": "Zakaria", "statut": "Sourcing"},
    {"titre": "Directeur Achats Adjoint", "entite": "TGCC", "lieu": "Siège", "demandeur": "C.BENABDELLAH", "recruteur": "Zakaria", "statut": "Sourcing"},
    {"titre": "INGENIEUR TRAVAUX", "entite": "TGCC", "lieu": "YAMED LOT B", "demandeur": "M.TAZI", "recruteur": "Zakaria", "statut": "Sourcing"},

    # Colonne Shortlisté
    {"titre": "CHEF DE PROJETS", "entite": "TGCC", "lieu": "DESSALMENT JORF", "demandeur": "M.FENNAN", "recruteur": "ZAKARIA", "statut": "Shortlisté"},
    {"titre": "Planificateur", "entite": "TGCC", "lieu": "ASFI-B", "demandeur": "SOUFIANI", "recruteur": "Ghita", "statut": "Shortlisté"},
    {"titre": "RESPONSABLE TRANS INTERCH", "entite": "TG PREFA", "lieu": "OUED SALEH", "demandeur": "FBOUZOUBAA", "recruteur": "Ghita", "statut": "Shortlisté"},

    # Colonne Signature DRH
    {"titre": "PROJETEUR DESSINATEUR", "entite": "TG WOOD", "lieu": "OUED SALEH", "demandeur": "S.MENJRA", "recruteur": "Zakaria", "statut": "Signature DRH"},
    {"titre": "Projeteur", "entite": "TGCC", "lieu": "TSP Safi", "demandeur": "B.MORABET", "recruteur": "Zakaria", "statut": "Signature DRH"},
    {"titre": "Consultant SAP", "entite": "TGCC", "lieu": "Siège", "demandeur": "O.KETTA", "recruteur": "Zakaria", "statut": "Signature DRH"},

    # Colonne Clôture
    {"titre": "Doc Controller", "entite": "TGEM", "lieu": "SIEGE", "demandeur": "A.SANKARI", "recruteur": "Zakaria", "statut": "Clôture"},
    {"titre": "Ingénieur étude/qualité", "entite": "TGCC", "lieu": "SIEGE", "demandeur": "A.MOUTANABI", "recruteur": "Zakaria", "statut": "Clôture"},
    {"titre": "Responsable Cybersecurité", "entite": "TGCC", "lieu": "Siège", "demandeur": "Ghazi", "recruteur": "Zakaria", "statut": "Clôture"},
    {"titre": "CHEF DE CHANTIER", "entite": "TGCC", "demandeur": "M.FENNAN", "recruteur": "Zakaria", "statut": "Clôture"},
    {"titre": "Ing contrôle de la performance", "entite": "TGCC", "lieu": "Siège", "demandeur": "H.BARIGOU", "recruteur": "Ghita", "statut": "Clôture"},
    {"titre": "Ingénieur Systèmes Réseaux", "entite": "TGCC", "lieu": "Siège", "demandeur": "M.JADDOR", "recruteur": "Ghita", "statut": "Clôture"},
    {"titre": "Responsable étude de prix", "entite": "TGCC", "lieu": "SIEGE", "demandeur": "S.Bennani Zitani", "recruteur": "Ghita", "statut": "Clôture"},
    {"titre": "Responsable Travaux", "entite": "TGEM", "lieu": "Zone Rabat", "demandeur": "S.ACHIR", "recruteur": "Zakaria", "statut": "Clôture"},

    # Colonne Désistement
    {"titre": "Conducteur de Travaux", "entite": "TGCC", "lieu": "JORF LASFAR", "demandeur": "M.FENNAN", "recruteur": "Zakaria", "statut": "Désistement"},
    {"titre": "Chef de Chantier", "entite": "TGCC", "lieu": "TOARC", "demandeur": "M.FENNAN", "recruteur": "Zakaria", "statut": "Désistement"},
    {"titre": "Magasinier", "entite": "TG WOOD", "lieu": "Oulad Saleh", "demandeur": "K.TAZI", "recruteur": "Ghita", "statut": "Désistement", "commentaire": "Pas de retour du demandeur"}
]


def create_integration_filters(df_recrutement, prefix=""):
    """Créer des filtres spécifiques pour les intégrations (sans les périodes)"""
    if df_recrutement is None or len(df_recrutement) == 0:
        return {}

    # Organiser les filtres dans deux colonnes dans la sidebar
    filters = {}
    left_col, right_col = st.sidebar.columns(2)

    # Filtre par entité demandeuse (colonne gauche)
    entites = ['Toutes'] + sorted(df_recrutement['Entité demandeuse'].dropna().unique())
    with left_col:
        filters['entite'] = st.selectbox("Entité demandeuse", entites, key=f"{prefix}_entite")

    # Filtre par direction concernée (colonne droite)
    directions = ['Toutes'] + sorted(df_recrutement['Direction concernée'].dropna().unique())
    with right_col:
        filters['direction'] = st.selectbox("Direction concernée", directions, key=f"{prefix}_direction")

    # Pas de filtres de période pour les intégrations
    filters['periode_recrutement'] = 'Toutes'
    filters['periode_demande'] = 'Toutes'

    return filters

def create_global_filters(df_recrutement, prefix="", include_periode_recrutement=True, include_periode_demande=True):
    """Créer des filtres globaux réutilisables pour tous les onglets.

    include_periode_recrutement et include_periode_demande contrôlent si le sélecteur
    de période correspondant est affiché (utile pour n'affecter qu'une section).
    """
    if df_recrutement is None or len(df_recrutement) == 0:
        return {}

    # Organiser les filtres dans deux colonnes dans la sidebar
    filters = {}
    left_col, right_col = st.sidebar.columns(2)

    # Filtre par entité demandeuse (colonne gauche)
    entites = ['Toutes'] + sorted(df_recrutement['Entité demandeuse'].dropna().unique())
    with left_col:
        filters['entite'] = st.selectbox("Entité demandeuse", entites, key=f"{prefix}_entite")

    # Filtre par direction concernée (colonne droite)
    directions = ['Toutes'] + sorted(df_recrutement['Direction concernée'].dropna().unique())
    with right_col:
        filters['direction'] = st.selectbox("Direction concernée", directions, key=f"{prefix}_direction")

    # Ajouter les filtres de période (sans ligne de séparation)
    left_col2, right_col2 = st.sidebar.columns(2)

    # Filtre Période de recrutement (basé sur Date d'entrée effective)
    if include_periode_recrutement:
        with left_col2:
            if 'Date d\'entrée effective du candidat' in df_recrutement.columns:
                df_recrutement['Année_Recrutement'] = df_recrutement['Date d\'entrée effective du candidat'].dt.year
                annees_rec = sorted([y for y in df_recrutement['Année_Recrutement'].dropna().unique() if not pd.isna(y)])
                if annees_rec:
                    filters['periode_recrutement'] = st.selectbox(
                        "Période de recrutement", 
                        ['Toutes'] + [int(a) for a in annees_rec], 
                        index=len(annees_rec), 
                        key=f"{prefix}_periode_rec"
                    )
                else:
                    filters['periode_recrutement'] = 'Toutes'
            else:
                filters['periode_recrutement'] = 'Toutes'
    else:
        # Ne pas afficher le sélecteur, s'assurer que la valeur reste 'Toutes'
        filters['periode_recrutement'] = 'Toutes'

    # Filtre Période de la demande (basé sur Date de réception de la demande)
    date_demande_col = 'Date de réception de la demande aprés validation de la DRH'
    if include_periode_demande:
        with right_col2:
            if date_demande_col in df_recrutement.columns:
                df_recrutement['Année_Demande'] = df_recrutement[date_demande_col].dt.year
                annees_dem = sorted([y for y in df_recrutement['Année_Demande'].dropna().unique() if not pd.isna(y)])
                if annees_dem:
                    filters['periode_demande'] = st.selectbox(
                        "Période de la demande", 
                        ['Toutes'] + [int(a) for a in annees_dem], 
                        index=len(annees_dem), 
                        key=f"{prefix}_periode_dem"
                    )
                else:
                    filters['periode_demande'] = 'Toutes'
            else:
                filters['periode_demande'] = 'Toutes'
    else:
        # Ne pas afficher le sélecteur, s'assurer que la valeur reste 'Toutes'
        filters['periode_demande'] = 'Toutes'

    return filters

def apply_global_filters(df, filters):
    """Appliquer les filtres globaux aux données"""
    df_filtered = df.copy()
    
    if filters.get('entite') != 'Toutes':
        df_filtered = df_filtered[df_filtered['Entité demandeuse'] == filters['entite']]
    
    if filters.get('direction') != 'Toutes':
        df_filtered = df_filtered[df_filtered['Direction concernée'] == filters['direction']]
    
    # Appliquer le filtre période de recrutement
    if filters.get('periode_recrutement') != 'Toutes' and 'Année_Recrutement' in df_filtered.columns:
        df_filtered = df_filtered[df_filtered['Année_Recrutement'] == filters['periode_recrutement']]
    
    # Appliquer le filtre période de la demande
    if filters.get('periode_demande') != 'Toutes' and 'Année_Demande' in df_filtered.columns:
        df_filtered = df_filtered[df_filtered['Année_Demande'] == filters['periode_demande']]
    
    return df_filtered


@st.cache_resource
def get_gsheet_client():
    """Crée et retourne un client gspread authentifié en utilisant les secrets Streamlit."""
    try:
        service_account_info = {
            "type": st.secrets["GCP_TYPE"],
            "project_id": st.secrets["GCP_PROJECT_ID"],
            "private_key_id": st.secrets["GCP_PRIVATE_KEY_ID"],
            "private_key": st.secrets["GCP_PRIVATE_KEY"].replace('\\n', '\n').strip(),
            "client_email": st.secrets["GCP_CLIENT_EMAIL"],
            "client_id": st.secrets["GCP_CLIENT_ID"],
            "auth_uri": st.secrets["GCP_AUTH_URI"],
            "token_uri": st.secrets["GCP_TOKEN_URI"],
            "auth_provider_x509_cert_url": st.secrets["GCP_AUTH_PROVIDER_CERT_URL"],
            "client_x509_cert_url": st.secrets["GCP_CLIENT_CERT_URL"]
        }
        return gspread.service_account_from_dict(service_account_info)
    except Exception as e:
        st.error(f"❌ Erreur de connexion à Google Sheets. Vérifiez vos secrets. Détails: {e}")
        return None


def load_data_from_google_sheets(sheet_url):
    """
    Charger les données depuis Google Sheets avec authentification automatique via les secrets.
    """
    try:
        # Extraire l'ID de la feuille et le GID depuis l'URL
        sheet_id_match = re.search(r"/d/([a-zA-Z0-9-_]+)", sheet_url)
        gid_match = re.search(r"[?&]gid=(\d+)", sheet_url)
        
        if not sheet_id_match:
            raise ValueError("Impossible d'extraire l'ID du Google Sheet depuis l'URL")
        
        sheet_id = sheet_id_match.group(1)
        gid = gid_match.group(1) if gid_match else '0'
        
        # Utiliser le client authentifié avec les secrets
        gc = get_gsheet_client()
        if gc is None:
            # Fallback vers l'export CSV public si l'authentification échoue
            export_url = f"https://docs.google.com/spreadsheets/d/{sheet_id}/export?format=csv&gid={gid}"
            df = pd.read_csv(export_url)
            return df
        
        # Ouvrir la feuille par ID
        sh = gc.open_by_key(sheet_id)
        
        # Sélectionner la worksheet par GID
        try:
            worksheet = sh.get_worksheet_by_id(int(gid))
        except:
            worksheet = sh.sheet1  # Fallback vers la première feuille
        
        # Récupérer toutes les données
        data = worksheet.get_all_records()
        df = pd.DataFrame(data)
        
        return df
        
    except Exception as e:
        # Si l'authentification échoue, essayer l'export CSV public
        try:
            sheet_id_match = re.search(r"/d/([a-zA-Z0-9-_]+)", sheet_url)
            gid_match = re.search(r"[?&]gid=(\d+)", sheet_url)
            
            if sheet_id_match:
                sheet_id = sheet_id_match.group(1)
                gid = gid_match.group(1) if gid_match else '0'
                export_url = f"https://docs.google.com/spreadsheets/d/{sheet_id}/export?format=csv&gid={gid}"
                df = pd.read_csv(export_url)
                return df
        except:
            pass
        
        raise e


def load_data_from_files(csv_file=None, excel_file=None):
    """Charger et préparer les données depuis les fichiers uploadés ou locaux"""
    df_integration = None
    df_recrutement = None
    try:
        # Charger le CSV (données d'intégration)
        if csv_file is not None:
            df_integration = pd.read_csv(csv_file)
        else:
            # Fallback vers fichier local s'il existe
            local_csv = '2025-10-09T20-31_export.csv'
            if os.path.exists(local_csv):
                df_integration = pd.read_csv(local_csv)

        if df_integration is not None and 'Date Intégration' in df_integration.columns:
            df_integration['Date Intégration'] = pd.to_datetime(df_integration['Date Intégration'])

        # Charger l'Excel (données de recrutement)
        # Si une synchronisation Google Sheets a été réalisée, utiliser ce DataFrame
        try:
            if 'synced_recrutement_df' in st.session_state and st.session_state.synced_recrutement_df is not None:
                df_recrutement = st.session_state.synced_recrutement_df.copy()
            elif excel_file is not None:
                df_recrutement = pd.read_excel(excel_file, sheet_name=0)
            else:
                # Fallback vers fichier local s'il existe (recherche du plus récent)
                import glob
                # Pattern pour trouver les fichiers Excel de recrutement
                excel_files = glob.glob('Recrutement global PBI All*.xlsx')
                if excel_files:
                    # Trier par date de modification (le plus récent en dernier)
                    excel_files.sort(key=os.path.getmtime)
                    latest_excel = excel_files[-1]
                    # st.info(f"Chargement automatique du fichier local : {latest_excel}")
                    df_recrutement = pd.read_excel(latest_excel, sheet_name=0)
                else:
                    # Fallback legacy
                    local_excel = 'Recrutement global PBI All  google sheet (5).xlsx'
                    if os.path.exists(local_excel):
                        df_recrutement = pd.read_excel(local_excel, sheet_name=0)
        except Exception as e:
            st.error(f"Erreur lors du chargement des données de recrutement: {e}")

        if df_recrutement is not None:
            # Nettoyer et préparer les données de recrutement
            # Convertir les dates
            date_columns = ['Date de réception de la demande aprés validation de la DRH',
                           'Date d\'entrée effective du candidat',
                           'Date d\'annulation /dépriorisation de la demande',
                           'Date de la 1er réponse du demandeur à l\'équipe RH',
                           'Date du 1er retour equipe RH  au demandeur',
                           'Date de désistement',
                           'Date d\'acceptation du candidat',
                           'Date d\'entrée prévisionnelle']
            
            for col in date_columns:
                if col in df_recrutement.columns:
                    try:
                        df_recrutement[col] = _parse_mixed_dates(df_recrutement[col])
                    except Exception:
                        # fallback
                        df_recrutement[col] = pd.to_datetime(df_recrutement[col], errors='coerce')
            
            # Nettoyer les colonnes avec des espaces
            df_recrutement.columns = df_recrutement.columns.str.strip()
            
            # Nettoyer les colonnes numériques pour éviter les erreurs de type
            numeric_columns = ['Nb de candidats pré-selectionnés']
            for col in numeric_columns:
                if col in df_recrutement.columns:
                    df_recrutement[col] = pd.to_numeric(df_recrutement[col], errors='coerce')
                    df_recrutement[col] = df_recrutement[col].fillna(0)

            # Vérification basique des colonnes critiques et message dans les logs
            required_cols = [
                'Statut de la demande', 'Poste demandé', 'Direction concernée',
                'Entité demandeuse', 'Modalité de recrutement'
            ]
            missing = [c for c in required_cols if c not in df_recrutement.columns]
            if missing:
                # Log via st.warning but don't raise — keep app running
                st.warning(f"Colonnes attendues manquantes dans le fichier de recrutement: {missing}")
        
        return df_integration, df_recrutement
        
    except Exception as e:
        st.error(f"Erreur lors du chargement des données: {e}")
        return None, None

def create_integration_timeline(df):
    """Créer un graphique de timeline des intégrations"""
    # Grouper par mois
    df['Mois'] = df['Date Intégration'].dt.to_period('M')
    monthly_stats = df.groupby(['Mois', 'Statut']).size().reset_index(name='Count')
    monthly_stats['Mois_str'] = monthly_stats['Mois'].astype(str)
    
    fig = px.bar(
        monthly_stats, 
        x='Mois_str', 
        y='Count',
        color='Statut',
        title="📈 Évolution des Intégrations par Mois",
        color_discrete_map={'En cours': '#ff6b6b', 'Complet': '#51cf66'}
    )
    
    fig.update_traces(hovertemplate='%{y}<extra></extra>')
    fig.update_layout(
        xaxis_title="Mois",
        yaxis_title="Nombre d'intégrations",
        showlegend=True,
        height=400
    )
    
    return fig

def create_affectation_chart(df):
    """Créer un graphique par affectation"""
    # Prefer 'Entité demandeuse' if present, else fallback to 'Affectation'
    if 'Entité demandeuse' in df.columns and df['Entité demandeuse'].notna().sum() > 0:
        serie = df['Entité demandeuse'].value_counts()
        title = "🏢 Répartition par Entité (Top 10)"
    elif 'Affectation' in df.columns:
        serie = df['Affectation'].value_counts()
        title = "🏢 Répartition par Affectation (Top 10)"
    else:
        serie = pd.Series([], dtype=int)
        title = "Répartition"

    serie = serie.head(10)
    fig = px.pie(
        values=serie.values,
        names=serie.index,
        title=title
    )
    fig.update_traces(textposition='inside', textinfo='percent+label', textfont=dict(size=14))
    fig.update_layout(
        height=420,
        title=dict(text=title, x=0, xanchor='left', font=TITLE_FONT),
        legend=dict(
            orientation="v",
            yanchor="middle",
            y=0.5,
            xanchor="left",
            x=0.75,  # Rapprochement vers la gauche vers le centre du graphique
            font=dict(size=16) # Taille de légende augmentée
        ),
        margin=dict(l=20, r=20, t=60, b=20)
    )
    return fig


def render_plotly_scrollable(fig, max_height=500):
    """Renders a plotly figure inside a scrollable HTML div so the user can scroll when there are many bars."""
    try:
        # Use components.html so the plotly JS is executed correctly in Streamlit
        html = fig.to_html(full_html=False, include_plotlyjs='cdn')
        # Inject a small CSS block that targets Plotly's title classes inside the
        # generated HTML so we can enforce consistent font and left alignment
        injected_css = """
    <style>
    /* Ensure the plotly SVG title (gtitle) uses our TITLE_FONT and is left-aligned */
    .plotly .gtitle, .plotly .gtitle text { font-family: Arial, sans-serif !important; font-size: 18px !important; fill: #111111 !important; }
    .plotly .gtitle { text-anchor: start !important; }
    /* Force the plot container to align left inside the Streamlit component */
    .streamlit-plotly-wrapper{ display:flex; justify-content:flex-start; }
    </style>
    """

        # Wrap the plot HTML in a left-aligned container so the chart isn't centered
        wrapper = f"""
<div class='streamlit-plotly-wrapper' style='width:100%;'>
  {injected_css}
  {html}
</div>
"""
        # components.html supports scrolling and executes scripts
        components.html(wrapper, height=max_height, scrolling=True)
    except Exception:
        # Fallback to default renderer
                fig = apply_title_style(fig)
                st.plotly_chart(fig, width="stretch")

def create_recrutements_clotures_tab(df_recrutement, global_filters):
    """Onglet Recrutements Clôturés avec style carte"""
    
    # Filtrer seulement les recrutements clôturés
    df_cloture = df_recrutement[df_recrutement['Statut de la demande'] == 'Clôture'].copy()
    
    if len(df_cloture) == 0:
        st.warning("Aucune donnée de recrutement clôturé disponible")
        return
    
    # Appliquer les filtres globaux
    df_filtered = apply_global_filters(df_cloture, global_filters)

    # Use the HTML KPI cards renderer for a more attractive layout
    recrutements = len(df_filtered)
    postes_uniques = df_filtered['Poste demandé'].nunique()
    directions_uniques = df_filtered['Direction concernée'].nunique()

    # Compute delai display and help
    date_reception_col = 'Date de réception de la demande aprés validation de la DRH'
    date_retour_rh_col = 'Date du 1er retour equipe RH  au demandeur'
    delai_display = "N/A"
    delai_help = "Colonnes manquantes ou pas de durées valides"
    if date_reception_col in df_filtered.columns and date_retour_rh_col in df_filtered.columns:
        try:
            s = pd.to_datetime(df_filtered[date_reception_col], errors='coerce')
            e = pd.to_datetime(df_filtered[date_retour_rh_col], errors='coerce')
            mask = s.notna() & e.notna()
            if mask.sum() > 0:
                durees = (e[mask] - s[mask]).dt.days
                durees = durees[durees >= 0]
                if len(durees) > 0:
                    delai_moyen = round(durees.mean(), 1)
                    delai_display = f"{delai_moyen}"
                    delai_help = f"Moyenne calculée sur {len(durees)} recrutements clôturés"
        except Exception:
            pass

    # Render KPIs as HTML cards for larger/consistent styling
    metrics_html = render_generic_metrics([
        ("Nombre de recrutements", recrutements, "#1f77b4"),
        ("Postes concernés", postes_uniques, "#2ca02c"),
        ("Nombre de Directions concernées", directions_uniques, "#ff7f0e"),
        ("Délai moyen recrutement (jours)", delai_display, "#6f42c1")
    ])
    st.markdown(metrics_html, unsafe_allow_html=True)
    
    # Graphiques en ligne 1
    col1, col2 = st.columns([2,1])
    
    with col1:
        # Évolution des recrutements par mois (comme dans l'image 1)
        if 'Date d\'entrée effective du candidat' in df_filtered.columns:
            # Générer une série complète de mois entre la plus petite et la plus grande date
            df_filtered['Mois_Année'] = df_filtered['Date d\'entrée effective du candidat'].dt.to_period('M').dt.to_timestamp()
            monthly_data = df_filtered.groupby('Mois_Année').size().rename('Count')
            if not monthly_data.empty:
                all_months = pd.date_range(start=monthly_data.index.min(), end=monthly_data.index.max(), freq='MS')
                monthly_data = monthly_data.reindex(all_months, fill_value=0)
                monthly_data = monthly_data.reset_index().rename(columns={'index': 'Mois_Année'})
                monthly_data['Mois_Année'] = monthly_data['Mois_Année'].dt.strftime('%b %Y')

                fig_evolution = px.bar(
                    monthly_data,
                    x='Mois_Année',
                    y='Count',
                    title="Évolution des recrutements",
                    text='Count'
                )
                fig_evolution.update_traces(
                    marker_color='#1f77b4',
                    textposition='inside',
                    texttemplate='<b>%{y}</b>',
                    textfont=dict(size=15, color='white'),
                    hovertemplate='%{y}<extra></extra>'
                )
                fig_evolution.update_layout(
                    height=360,
                    margin=dict(t=48, b=30, l=20, r=20),
                    xaxis_title=None,
                    yaxis_title=None,
                    xaxis=dict(
                        tickmode='array',
                        tickvals=monthly_data['Mois_Année'],
                        ticktext=monthly_data['Mois_Année'],
                        tickangle=45
                    )
                )
                fig_evolution = apply_title_style(fig_evolution)
                st.plotly_chart(fig_evolution, width="stretch")
    
    with col2:
        # Répartition par modalité de recrutement (CORRECTION: légende déplacée à l'extérieur)
        if 'Modalité de recrutement' in df_filtered.columns:
            modalite_data = df_filtered['Modalité de recrutement'].value_counts()
            
            fig_modalite = go.Figure(data=[go.Pie(
                labels=modalite_data.index, 
                values=modalite_data.values,
                hole=.5,
                textposition='inside',
                textinfo='percent'
            )])
            fig_modalite.update_traces(textfont=dict(size=14))
            fig_modalite.update_layout(
                title=dict(text="Répartition par Modalité de recrutement", x=0, xanchor='left', font=TITLE_FONT),
                height=380,
                # Légende positionnée à droite pour éviter le chevauchement
                legend=dict(
                    orientation="v", 
                    yanchor="middle", 
                    y=0.5, 
                    xanchor="left", 
                    x=1.0,
                    font=dict(size=14)
                ),
                # Ajuster les marges pour faire de la place à la légende
                margin=dict(l=20, r=140, t=60, b=20)
            )
            fig_modalite = apply_title_style(fig_modalite)
            st.plotly_chart(fig_modalite, width="stretch")

    # Graphiques en ligne 2
    col3, col4 = st.columns(2)
    
    with col3:
        # Comparaison par direction
        direction_counts = df_filtered['Direction concernée'].value_counts()
        # Convert Series to DataFrame for plotly express compatibility
        df_direction = direction_counts.rename_axis('Direction').reset_index(name='Count')
        df_direction = df_direction.sort_values('Count', ascending=False)
        # Truncate long labels for readability, keep full label in customdata for hover
        df_direction['Label_trunc'] = df_direction['Direction'].apply(lambda s: _truncate_label(s, max_len=24))
        # Ensure a display label exists (truncated + small gap) to be used for axis and ordering
        if 'Label_display' not in df_direction.columns:
            df_direction['Label_display'] = df_direction['Label_trunc'] + '\u00A0\u00A0'
        # Add two non-breaking spaces to create visual gap between label and bar
        df_direction['Label_display'] = df_direction['Label_trunc'] + '\u00A0\u00A0'
        fig_direction = px.bar(
            df_direction,
            x='Count',
            y='Label_display',
            title="Comparaison par direction",
            text='Count',
            orientation='h',
            custom_data=['Direction']
        )
        fig_direction.update_traces(
            marker_color='grey',
            textposition='auto',
            texttemplate='<b>%{x}</b>',
            textfont=dict(size=15, color='white'),
            textangle=0,  # Forcer l'orientation horizontale des valeurs
            hovertemplate='<b>%{customdata[0]}</b><br>Nombre: %{x}<extra></extra>'
        )
        fig_direction.update_layout(
            height=300,
            xaxis_title=None,
            yaxis_title=None,
            margin=dict(l=160, t=48, b=30, r=20),
            xaxis=dict(tickangle=0),
            yaxis=dict(automargin=True, tickfont=dict(size=15), ticklabelposition='outside left', categoryorder='array', categoryarray=list(df_direction['Label_display'][::-1])),
            title=dict(text="<b>Comparaison par direction</b>", x=0, xanchor='left', font=TITLE_FONT)
        )
        fig_direction = apply_title_style(fig_direction)
        try:
            fig_direction.update_traces(textfont=dict(size=15, color='white'))
        except Exception:
            pass
        # Use a compact default visible area and allow scrolling when long
        render_plotly_scrollable(fig_direction, max_height=320)

    with col4:
        # Comparaison par poste
        poste_counts = df_filtered['Poste demandé'].value_counts()
        df_poste = poste_counts.rename_axis('Poste').reset_index(name='Count')
        df_poste = df_poste.sort_values('Count', ascending=False)
        df_poste['Label_trunc'] = df_poste['Poste'].apply(lambda s: _truncate_label(s, max_len=24))
        if 'Label_display' not in df_poste.columns:
            df_poste['Label_display'] = df_poste['Label_trunc'] + '\u00A0\u00A0'
        df_poste['Label_display'] = df_poste['Label_trunc'] + '\u00A0\u00A0'
        fig_poste = px.bar(
            df_poste,
            x='Count',
            y='Label_display',
            title="Comparaison par poste",
            text='Count',
            orientation='h',
            custom_data=['Poste']
        )
        fig_poste.update_traces(
            marker_color='grey',
            textposition='auto',
            texttemplate='<b>%{x}</b>',
            textfont=dict(size=15, color='white'),
            textangle=0,  # Forcer l'orientation horizontale des valeurs
            hovertemplate='<b>%{customdata[0]}</b><br>Nombre: %{x}<extra></extra>'
        )
        height_poste = max(300, 28 * len(df_poste))
        fig_poste.update_layout(
            height=300 if height_poste < 360 else height_poste,
            xaxis_title=None,
            yaxis_title=None,
            margin=dict(l=160, t=48, b=30, r=20),
            xaxis=dict(tickangle=0),
            yaxis=dict(automargin=True, tickfont=dict(size=15), ticklabelposition='outside left', categoryorder='array', categoryarray=list(df_poste['Label_display'][::-1])),
            title=dict(text="<b>Comparaison par poste</b>", x=0, xanchor='left', font=TITLE_FONT)
        )
        fig_poste = apply_title_style(fig_poste)
        try:
            fig_poste.update_traces(textfont=dict(size=15, color='white'))
        except Exception:
            pass
        render_plotly_scrollable(fig_poste, max_height=320)


    # Ligne 3 - KPIs de délai et candidats
    col5, col6 = st.columns(2)

    with col5:
        # Nombre de candidats présélectionnés - avec conversion sécurisée
        try:
            # Convertir les valeurs en numérique, remplacer les erreurs par 0
            candidats_series = pd.to_numeric(df_filtered['Nb de candidats pré-selectionnés'], errors='coerce').fillna(0)
            total_candidats = int(candidats_series.sum())
        except (KeyError, ValueError):
            total_candidats = 0
            
        # Titre stylé comme les autres graphiques
        st.markdown("<div style='font-family:Arial,sans-serif; font-size:18px; font-weight:700; color:#111111; text-align:left; margin:8px 0 4px 0;'>Nombre de candidats présélectionnés</div>", unsafe_allow_html=True)

        fig_candidats = go.Figure(go.Indicator(
            mode = "gauge+number",
            value = total_candidats,
            gauge = {'axis': {'range': [0, max(total_candidats * 2, 100)], 'visible': True},
                     'bar': {'color': "green"},
                    }))
        fig_candidats.update_layout(height=260, margin=dict(t=10, b=10, l=20, r=20))
        st.plotly_chart(fig_candidats, width="stretch")

    with col6:
        # Taux de refus = (Nb de refus aux promesses / Nb de promesses réalisées) * 100
        # Recherche robuste des colonnes dans le dataframe
        prom_col = next((c for c in df_filtered.columns if 'promess' in c.lower() and 'réalis' in c.lower()), None)
        refus_col = next((c for c in df_filtered.columns if 'refus' in c.lower()), None)
        prom_sum = 0
        refus_sum = 0
        try:
            if prom_col and prom_col in df_filtered.columns:
                prom_sum = pd.to_numeric(df_filtered[prom_col], errors='coerce').fillna(0).sum()
            if refus_col and refus_col in df_filtered.columns:
                refus_sum = pd.to_numeric(df_filtered[refus_col], errors='coerce').fillna(0).sum()
        except Exception:
            prom_sum = 0; refus_sum = 0

        taux_refus = 0.0
        if prom_sum and prom_sum > 0:
            taux_refus = float(refus_sum) / float(prom_sum) * 100.0

        st.markdown("<div style='font-family:Arial,sans-serif; font-size:18px; font-weight:700; color:#111111; text-align:left; margin:8px 0 4px 0;'>Taux de refus des promesses d'embauche (%)</div>", unsafe_allow_html=True)
        fig_refus = go.Figure(go.Indicator(
            mode='gauge+number',
            value=round(taux_refus, 1),
            number={'suffix':' %'},
            gauge={'axis': {'range':[0,100], 'visible': True}, 'bar': {'color':'#d62728'}}
        ))
        fig_refus.update_layout(height=280, margin=dict(t=20, b=20, l=20, r=20))
        st.plotly_chart(fig_refus, width='stretch')

    # Debug local pour l'onglet Recrutements Clôturés
    st.markdown("---")
    with st.expander("🔍 Debug - Détails des lignes (Recrutements Clôturés)", expanded=False):
        try:
            st.markdown("**Lignes contribuant aux graphiques avec calcul du délai:**")
            df_debug_clo = df_filtered.copy()
            # Colonnes source pour le délai
            date_reception_col = 'Date de réception de la demande aprés validation de la DRH'
            date_entree_col = "Date d'entrée effective du candidat"
            # Formater les dates sans heure
            if date_entree_col in df_debug_clo.columns:
                df_debug_clo[date_entree_col] = pd.to_datetime(df_debug_clo[date_entree_col], errors='coerce').dt.strftime('%d/%m/%Y')
            if date_reception_col in df_debug_clo.columns:
                df_debug_clo['Date Réception'] = pd.to_datetime(df_debug_clo[date_reception_col], errors='coerce').dt.strftime('%d/%m/%Y')
                # Calculer le délai en jours
                date_rec = pd.to_datetime(df_filtered[date_reception_col], errors='coerce')
                date_ent = pd.to_datetime(df_filtered[date_entree_col], errors='coerce')
                df_debug_clo['Délai (jours)'] = (date_ent - date_rec).dt.days
            cols_debug = ['Poste demandé', 'Entité demandeuse', 'Direction concernée', 'Date Réception', date_entree_col, 'Délai (jours)', 'Modalité de recrutement']
            cols_available = [c for c in cols_debug if c in df_debug_clo.columns]
            if cols_available:
                st.dataframe(df_debug_clo[cols_available].reset_index(drop=True), use_container_width=True, hide_index=True)
            else:
                st.dataframe(df_debug_clo.reset_index(drop=True), use_container_width=True, hide_index=True)
        except Exception:
            st.write("Aucune donnée disponible pour le debug.")

    # ... KPI row now includes Délai moyen de recrutement (moved up)

def create_demandes_recrutement_tab(df_recrutement, global_filters):
    """Onglet Demandes de Recrutement avec style carte"""
    
    # Appliquer les filtres globaux
    df_filtered = apply_global_filters(df_recrutement, global_filters)
    
    # Colonne de date pour les calculs
    date_col = 'Date de réception de la demande aprés validation de la DRH'
    
    # KPIs principaux - Indicateurs de demandes sur la même ligne
    # Rendre les KPI principaux via HTML pour contrôle précis des tailles
    total_demandes = len(df_filtered)
    today = datetime.now()
    start_of_month = today.replace(day=1)
    nouvelles_demandes = 0
    if date_col in df_filtered.columns:
        try:
            nouvelles_demandes = int((pd.to_datetime(df_filtered[date_col], errors='coerce') >= start_of_month).sum())
        except Exception:
            nouvelles_demandes = 0

    demandes_annulees = 0
    if 'Statut de la demande' in df_filtered.columns:
        demandes_annulees = int(df_filtered['Statut de la demande'].astype(str).str.contains('annul|déprioris', case=False, na=False).sum())
    taux_annulation = f"{round((demandes_annulees / total_demandes) * 100, 1)}%" if total_demandes > 0 else "N/A"

    metrics_html = render_generic_metrics([
        ("Nombre de demandes", total_demandes, "#1f77b4"),
        ("Nouvelles Demandes (ce mois-ci)", nouvelles_demandes, "#2ca02c"),
        ("Demandes Annulées/Dépriorisées", demandes_annulees, "#ff7f0e"),
        ("Taux d'annulation", taux_annulation, "#d62728")
    ])
    st.markdown(metrics_html, unsafe_allow_html=True)

    # Graphiques principaux
    col1, col2, col3 = st.columns([1,1,2])
    
    with col1:
        # Répartition par statut de la demande
        statut_counts = df_filtered['Statut de la demande'].value_counts()
        fig_statut = go.Figure(data=[go.Pie(labels=statut_counts.index, values=statut_counts.values, hole=.5)])
        fig_statut.update_traces(textfont=dict(size=14))
        fig_statut.update_layout(
            title=dict(text="Répartition par statut de la demande", x=0, xanchor='left', font=TITLE_FONT),
            height=380,
            legend=dict(orientation="h", yanchor="bottom", y=-0.45, xanchor="center", x=0.5, font=dict(size=13))
        )
        fig_statut = apply_title_style(fig_statut)
        st.plotly_chart(fig_statut, width="stretch")
    
    with col2:
        # Comparaison par raison du recrutement
        if 'Raison du recrutement' in df_filtered.columns:
            raison_counts = df_filtered['Raison du recrutement'].value_counts()
            df_raison = raison_counts.rename_axis('Raison').reset_index(name='Count')
            fig_raison = px.bar(
                df_raison,
                x='Raison',
                y='Count',
                title="<b>Comparaison par raison du recrutement</b>",
                text='Count',
                orientation='v'
            )
            fig_raison.update_traces(
                marker_color='grey', 
                textposition='auto',
                texttemplate='<b>%{y}</b>',
                textfont=dict(size=15, color='white'),
                hovertemplate='%{y}<extra></extra>'
            )
            fig_raison.update_layout(
                height=300, 
                xaxis_title=None, 
                yaxis_title=None,
                xaxis={'categoryorder':'total descending'},
                title=dict(text="<b>Comparaison par raison du recrutement</b>", x=0, xanchor='left', font=TITLE_FONT)
            )
            fig_raison = apply_title_style(fig_raison)
            st.plotly_chart(fig_raison, width="stretch")
    
    with col3:
        # Évolution des demandes
        if date_col in df_filtered.columns:
            # Générer une série complète de mois entre la plus petite et la plus grande date de demande
            df_filtered['Mois_Année_Demande'] = df_filtered[date_col].dt.to_period('M').dt.to_timestamp()
            monthly_demandes = df_filtered.groupby('Mois_Année_Demande').size().rename('Count')
            if not monthly_demandes.empty:
                all_months = pd.date_range(start=monthly_demandes.index.min(), end=monthly_demandes.index.max(), freq='MS')
                monthly_demandes = monthly_demandes.reindex(all_months, fill_value=0)
                monthly_demandes = monthly_demandes.reset_index().rename(columns={'index': 'Mois_Année_Demande'})
                monthly_demandes['Mois_Année_Demande'] = monthly_demandes['Mois_Année_Demande'].dt.strftime('%b %Y')

                fig_evolution_demandes = px.bar(
                    monthly_demandes,
                    x='Mois_Année_Demande',
                    y='Count',
                    title="Évolution des demandes",
                    text='Count'
                )
                fig_evolution_demandes.update_traces(
                    marker_color='#1f77b4',
                    textposition='inside',
                    texttemplate='<b>%{y}</b>',
                    textfont=dict(size=15, color='white'),
                    hovertemplate='%{y}<extra></extra>'
                )
                # Aligner la marge supérieure avec les autres titres (ex: pie statuts)
                fig_evolution_demandes.update_layout(height=320, margin=dict(t=48, b=30, l=20, r=20), xaxis_title=None, yaxis_title=None)
                fig_evolution_demandes = apply_title_style(fig_evolution_demandes)
                st.plotly_chart(fig_evolution_demandes, width="stretch")
    
    # Deuxième ligne de graphiques
    col4, col5 = st.columns(2)
    
    with col4:
        # Comparaison par direction
        direction_counts = df_filtered['Direction concernée'].value_counts()
        df_direction = direction_counts.rename_axis('Direction').reset_index(name='Count')
        df_direction = df_direction.sort_values('Count', ascending=False)
        # Truncate long labels for readability, keep full label in customdata for hover
        df_direction['Label_trunc'] = df_direction['Direction'].apply(lambda s: _truncate_label(s, max_len=24))
        # Ensure display label exists (truncated + small gap) to be used for axis and ordering
        if 'Label_display' not in df_direction.columns:
            df_direction['Label_display'] = df_direction['Label_trunc'] + '\u00A0\u00A0'
        fig_direction = px.bar(
            df_direction,
            x='Count',
            y='Label_display',
            title="Comparaison par direction",
            text='Count',
            orientation='h',
            custom_data=['Direction']
        )
        # Show values inside bars and full label on hover
        fig_direction.update_traces(
            marker_color='grey',
            textposition='auto',
            texttemplate='<b>%{x}</b>',
            textfont=dict(size=15, color='white'),
            textangle=0, # Orientation horizontale des valeurs
            hovertemplate='<b>%{customdata[0]}</b><br>Nombre: %{x}<extra></extra>'
        )
        # Dynamic height so long lists become scrollable on the page
        height_dir = max(300, 28 * len(df_direction))
        # Ensure largest values appear on top by reversing the truncated label array
        category_array_dir = list(df_direction['Label_display'][::-1])
        fig_direction.update_layout(
            height=height_dir,
            xaxis_title=None,
            yaxis_title=None,
            margin=dict(l=160, t=48, b=30, r=20),
            xaxis=dict(tickangle=0),
            yaxis=dict(automargin=True, tickfont=dict(size=15), ticklabelposition='outside left', categoryorder='array', categoryarray=category_array_dir),
            title=dict(text="<b>Comparaison par direction</b>", x=0, xanchor='left', font=TITLE_FONT)
        )
        # Render inside the column so the two charts are on the same row and the component width matches the column
        render_plotly_scrollable(fig_direction, max_height=320)
    
    with col5:
        # Comparaison par poste
        poste_counts = df_filtered['Poste demandé'].value_counts()
        df_poste = poste_counts.rename_axis('Poste').reset_index(name='Count')
        df_poste = df_poste.sort_values('Count', ascending=False)
        df_poste['Label_trunc'] = df_poste['Poste'].apply(lambda s: _truncate_label(s, max_len=24))
        if 'Label_display' not in df_poste.columns:
            df_poste['Label_display'] = df_poste['Label_trunc'] + '\u00A0\u00A0'
        df_poste['Label_display'] = df_poste['Label_trunc'] + '\u00A0\u00A0'
        fig_poste = px.bar(
            df_poste,
            x='Count',
            y='Label_display',
            title="Comparaison par poste",
            text='Count',
            orientation='h',
            custom_data=['Poste']
        )
        fig_poste.update_traces(
            marker_color='grey',
            textposition='auto',
            texttemplate='<b>%{x}</b>',
            textfont=dict(size=15, color='white'),
            textangle=0, # Orientation horizontale des valeurs
            hovertemplate='<b>%{customdata[0]}</b><br>Nombre: %{x}<extra></extra>'
        )
        height_poste = 320
        category_array_poste = list(df_poste['Label_display'][::-1])
        fig_poste.update_layout(
            height=320,
            xaxis_title=None,
            yaxis_title=None,
            margin=dict(l=160, t=48, b=30, r=20),
            xaxis=dict(tickangle=0),
            yaxis=dict(automargin=True, tickfont=dict(size=15), ticklabelposition='outside left', categoryorder='array', categoryarray=category_array_poste),
            title=dict(text="<b>Comparaison par poste</b>", x=0, xanchor='left', font=TITLE_FONT)
        )
        render_plotly_scrollable(fig_poste, max_height=320)

    # Debug local pour l'onglet Demandes de Recrutement (en dehors des colonnes, aligné à gauche)
    st.markdown("---")
    with st.expander("🔍 Debug - Détails des lignes (Demandes de Recrutement)", expanded=False):
        try:
            st.markdown("**Lignes contribuant aux graphiques (toutes les données filtrées):**")
            df_debug_dem = df_filtered.copy()
            # Formater la date sans heure
            date_col_dem = 'Date de réception de la demande aprés validation de la DRH'
            if date_col_dem in df_debug_dem.columns:
                df_debug_dem[date_col_dem] = pd.to_datetime(df_debug_dem[date_col_dem], errors='coerce').dt.strftime('%d/%m/%Y')
            cols_debug = ['Poste demandé', 'Entité demandeuse', 'Direction concernée', 'Raison du recrutement', 'Statut de la demande', date_col_dem]
            cols_available = [c for c in cols_debug if c in df_debug_dem.columns]
            if cols_available:
                st.dataframe(df_debug_dem[cols_available].reset_index(drop=True), use_container_width=True, hide_index=True)
            else:
                st.dataframe(df_debug_dem.reset_index(drop=True), use_container_width=True, hide_index=True)
        except Exception:
            st.write("Aucune donnée disponible pour le debug.")

def create_integrations_tab(df_recrutement, global_filters):
    """Onglet Intégrations basé sur les bonnes données"""
    st.header("📊 Intégrations")
    
    # Filtrer les données : Statut "En cours" ET candidat ayant accepté (nom présent)
    candidat_col = "Nom Prénom du candidat retenu yant accepté la promesse d'embauche"
    date_integration_col = "Date d'entrée prévisionnelle"
    plan_integration_col = "Plan d'intégration à préparer"
    
    # Diagnostic des données disponibles
    total_en_cours = len(df_recrutement[df_recrutement['Statut de la demande'] == 'En cours'])
    avec_candidat = len(df_recrutement[
        (df_recrutement['Statut de la demande'] == 'En cours') &
        (df_recrutement[candidat_col].notna()) &
        (df_recrutement[candidat_col].str.strip() != "")
    ])
    avec_date_prevue = len(df_recrutement[
        (df_recrutement['Statut de la demande'] == 'En cours') &
        (df_recrutement[candidat_col].notna()) &
        (df_recrutement[candidat_col].str.strip() != "") &
        (df_recrutement[date_integration_col].notna())
    ])
    
    # Critères : Statut "En cours" ET candidat avec nom
    df_integrations = df_recrutement[
        (df_recrutement['Statut de la demande'] == 'En cours') &
        (df_recrutement[candidat_col].notna()) &
        (df_recrutement[candidat_col].str.strip() != "")
    ].copy()
    
    # Message de diagnostic
    if total_en_cours > 0:
        st.info(f"📊 Diagnostic: {total_en_cours} demandes 'En cours' • {avec_candidat} avec candidat nommé • {avec_date_prevue} avec date d'entrée prévue")
    
    if len(df_integrations) == 0:
        st.warning("Aucune intégration en cours trouvée")
        st.info("Vérifiez que les demandes ont le statut 'En cours' ET un nom de candidat dans la colonne correspondante.")
        return
    
    # Appliquer les filtres globaux
    df_filtered = apply_global_filters(df_integrations, global_filters)
    
    # KPIs d'intégration (HTML cards pour contrôle précis des tailles)
    nb_int = len(df_filtered)
    a_preparer = 0
    if plan_integration_col in df_filtered.columns:
        try:
            a_preparer = int((df_filtered[plan_integration_col].astype(str).str.lower() == 'oui').sum())
        except Exception:
            a_preparer = 0

    en_retard = "N/A"
    if date_integration_col in df_filtered.columns:
        try:
            df_temp = df_filtered.copy()
            df_temp[date_integration_col] = pd.to_datetime(df_temp[date_integration_col], errors='coerce')
            reporting_date = st.session_state.get('reporting_date', None)
            if reporting_date is None:
                today = datetime.now()
            else:
                if isinstance(reporting_date, datetime):
                    today = reporting_date
                else:
                    today = datetime.combine(reporting_date, datetime.min.time())
            en_retard = int(((df_temp[date_integration_col].notna()) & (df_temp[date_integration_col] < today)).sum())
        except Exception:
            en_retard = "N/A"

    metrics_html = render_generic_metrics([
        ("👥 Intégrations en cours", nb_int, "#1f77b4"),
        ("📋 Plan d'intégration à préparer", a_preparer, "#ff7f0e"),
        ("⚠️ En retard", en_retard, "#d62728")
    ])
    st.markdown(metrics_html, unsafe_allow_html=True)
    
    # Graphiques
    col1, col2 = st.columns(2)

    # Graphique par affectation réactivé
    with col1:
        if 'Affectation' in df_filtered.columns:
            # Utiliser la fonction existante create_affectation_chart
            fig_affectation = create_affectation_chart(df_filtered)
            fig_affectation = apply_title_style(fig_affectation)
            st.plotly_chart(fig_affectation, width="stretch")
        else:
            st.warning("Colonne 'Affectation' non trouvée dans les données.")

    with col2:
        # Évolution des dates d'intégration prévues
        if date_integration_col in df_filtered.columns:
            df_filtered['Mois_Integration'] = df_filtered[date_integration_col].dt.to_period('M')
            monthly_integration = df_filtered.groupby('Mois_Integration').size().reset_index(name='Count')
            # Convertir en nom de mois seulement (ex: "Janvier", "Février")
            monthly_integration['Mois_str'] = monthly_integration['Mois_Integration'].dt.strftime('%B').str.capitalize()
            
            fig_evolution_int = px.bar(
                monthly_integration, 
                x='Mois_str', 
                y='Count',
                title="📈 Évolution des Intégrations Prévues",
                text='Count'
            )
            fig_evolution_int.update_traces(
                marker_color='#2ca02c', 
                textposition='inside',
                texttemplate='<b>%{y}</b>',
                textfont=dict(size=15, color='white'),
                hovertemplate='%{y}<extra></extra>'
            )
            fig_evolution_int.update_layout(height=400, xaxis_title="Mois", yaxis_title="Nombre")
            st.plotly_chart(fig_evolution_int, width="stretch")
    
    # Tableau détaillé des intégrations
    st.subheader("📋 Détail des Intégrations en Cours")
    
    # Détection automatique de la colonne Poste pour éviter les problèmes d'espaces
    poste_col_detected = 'Poste demandé ' if 'Poste demandé ' in df_filtered.columns else 'Poste demandé'
    
    colonnes_affichage = [
        candidat_col, 
        poste_col_detected,
        'Entité demandeuse',
        'Affectation',
        date_integration_col,
        plan_integration_col
    ]
    # Filtrer les colonnes qui existent
    colonnes_disponibles = [col for col in colonnes_affichage if col in df_filtered.columns]
    
    if colonnes_disponibles:
        # Travailler sur une copie des colonnes disponibles
        df_display = df_filtered[colonnes_disponibles].copy()

        # Exclure les lignes sans date d'intégration prévue AVANT le formatage
        if date_integration_col in df_display.columns:
            df_display = df_display[df_display[date_integration_col].notna() & df_display[date_integration_col].astype(str).str.strip().ne('')].copy()
            
            # Trier par date d'intégration croissante AVANT formatage
            try:
                df_display['_sort_date'] = pd.to_datetime(df_display[date_integration_col], errors='coerce')
                df_display = df_display.sort_values('_sort_date', ascending=True)
                df_display = df_display.drop(columns=['_sort_date'])
            except Exception:
                pass

            # Formater la date pour enlever l'heure et s'assurer du bon format DD/MM/YYYY
            def format_date_safely(date_str):
                if pd.isna(date_str) or date_str == '' or str(date_str).strip().upper() == 'N/A' or date_str is pd.NaT:
                    return 'N/A'
                parsed_date = pd.to_datetime(date_str, errors='coerce')
                if pd.notna(parsed_date):
                    return parsed_date.strftime('%d/%m/%Y')
                else:
                    return 'N/A'

            df_display[date_integration_col] = df_display[date_integration_col].apply(format_date_safely)

        # Renommer pour affichage plus propre
        rename_map = {
            candidat_col: "Candidat",
            date_integration_col: "Date d'Intégration Prévue"
        }
        if poste_col_detected in df_display.columns:
            rename_map[poste_col_detected] = "Poste demandé"
            
        df_display_table = df_display.rename(columns=rename_map)

        # Style CSS (version texte uniquement, plus aérée)
        st.markdown("""
        <style>
        .int-table-container {
            display: flex;
            justify-content: center;
            width: 100%;
            margin: 20px 0;
        }
        .int-custom-table {
            border-collapse: collapse;
            font-family: Arial, sans-serif;
            box-shadow: 0 1px 4px rgba(0,0,0,0.1);
            width: 85%;
            margin: 0 auto;
        }
        .int-custom-table th {
            background-color: #9C182F !important;
            color: white !important;
            font-weight: bold !important;
            text-align: center !important;
            padding: 10px 15px !important;
            border: 1px solid white !important;
            font-size: 1.05em;
            white-space: nowrap;
        }
        .int-custom-table td {
            text-align: center !important;
            padding: 8px 12px !important;
            border: 1px solid #ddd !important;
            background-color: white !important;
            font-size: 1.0em;
            font-weight: 500;
        }
        .int-custom-table .candidate-cell {
            font-weight: bold;
            color: #2c3e50;
            text-align: left !important;
            padding-left: 15px !important;
        }
        </style>
        """, unsafe_allow_html=True)

        html_table = '<div class="int-table-container"><table class="int-custom-table"><thead><tr>'
        for col in df_display_table.columns:
            html_table += f'<th>{col}</th>'
        html_table += '</tr></thead><tbody>'

        for _, row in df_display_table.iterrows():
            html_table += '<tr>'
            for col in df_display_table.columns:
                val = row[col]
                if col == 'Candidat':
                    html_table += f'<td class="candidate-cell">{val}</td>'
                else:
                    html_table += f'<td>{val}</td>'
            html_table += '</tr>'
        
        html_table += '</tbody></table></div>'
        st.markdown(html_table, unsafe_allow_html=True)
    else:
        st.warning("Colonnes d'affichage non disponibles")

    # Debug local pour l'onglet Intégrations avec colonnes de contribution
    st.markdown("---")
    with st.expander("🔍 Debug - Détails des lignes (Intégrations)", expanded=False):
        try:
            st.markdown("**Lignes contribuant aux indicateurs (en cours, en retard, plan d'intégration à préparer):**")
            # Construire les colonnes de contribution pour les indicateurs
            df_debug_int = df_filtered.copy()
            
            # Contribution: En cours (toutes les lignes filtrées)
            df_debug_int['contrib_en_cours'] = True
            
            # Contribution: En retard (date d'intégration prévue < aujourd'hui)
            today = datetime.now()
            if date_integration_col in df_debug_int.columns:
                df_debug_int['_date_check'] = pd.to_datetime(df_debug_int[date_integration_col], errors='coerce')
                df_debug_int['contrib_en_retard'] = df_debug_int['_date_check'] < today
                df_debug_int = df_debug_int.drop(columns=['_date_check'])
            else:
                df_debug_int['contrib_en_retard'] = False
            
            # Contribution: Plan d'intégration à préparer (plan = 'oui')
            if plan_integration_col in df_debug_int.columns:
                df_debug_int['contrib_plan_a_preparer'] = df_debug_int[plan_integration_col].astype(str).str.lower() == 'oui'
            else:
                df_debug_int['contrib_plan_a_preparer'] = False
            
            # Trier par date d'intégration croissante
            if date_integration_col in df_debug_int.columns:
                df_debug_int['_sort_date'] = pd.to_datetime(df_debug_int[date_integration_col], errors='coerce')
                df_debug_int = df_debug_int.sort_values('_sort_date', ascending=True)
                df_debug_int = df_debug_int.drop(columns=['_sort_date'])
            
            # Formater la date sans l'heure
            if date_integration_col in df_debug_int.columns:
                df_debug_int[date_integration_col] = pd.to_datetime(df_debug_int[date_integration_col], errors='coerce').dt.strftime('%d/%m/%Y')
            
            # Sélectionner colonnes pertinentes (sans Plan d'intégration)
            cols_display = [candidat_col, 'Poste demandé ', 'Entité demandeuse', date_integration_col, 'contrib_en_cours', 'contrib_en_retard', 'contrib_plan_a_preparer']
            cols_available = [c for c in cols_display if c in df_debug_int.columns]
            
            st.dataframe(df_debug_int[cols_available].reset_index(drop=True), use_container_width=True, hide_index=True)
        except Exception as e:
            st.write(f"Aucune donnée disponible pour le debug. Erreur: {e}")


def create_demandes_recrutement_combined_tab(df_recrutement):
    """Onglet combiné Demandes et Recrutement avec cartes expandables comme Home.py"""
    st.header("📊 Demandes & Recrutement")
    
    # CSS pour les cartes style Home.py
    st.markdown("""
    <style>
    .report-card {
        border-radius: 8px;
        background-color: #f8f9fa;
        padding: 15px;
        margin-bottom: 15px;
        border-left: 4px solid #007bff;
        box-shadow: 0 2px 4px rgba(0,0,0,0.1);
    }
    .report-card h4 {
        margin-top: 0;
        margin-bottom: 10px;
        color: #2c3e50;
        font-size: 1.25em;
        font-weight: 600;
    }
    .report-card p {
        margin-bottom: 8px;
        font-size: 1.02em;
        color: #5a6c7d;
    }
    .report-card .status-badge {
        display: inline-block;
        padding: 4px 8px;
        border-radius: 12px;
        font-size: 0.8em;
        font-weight: bold;
        color: white;
        background-color: #007bff;
    }
    </style>
    """, unsafe_allow_html=True)
    
    # Créer un seul jeu de filtres (4 contrôles): Entité, Direction, Période de la demande, Période de recrutement
    st.sidebar.subheader("🔧 Filtres Globaux")
    shared_filters = create_global_filters(df_recrutement, "combined", include_periode_recrutement=True, include_periode_demande=True)

    # Dériver deux jeux de filtres à partir des filtres partagés pour que chaque section
    # n'applique que la période qui lui est pertinente.
    filters_demandes = {
        'entite': shared_filters.get('entite', 'Toutes'),
        'direction': shared_filters.get('direction', 'Toutes'),
        'periode_demande': shared_filters.get('periode_demande', 'Toutes'),
        # Ne pas filtrer par période de recrutement dans la section Demandes
        'periode_recrutement': 'Toutes'
    }

    filters_clotures = {
        'entite': shared_filters.get('entite', 'Toutes'),
        'direction': shared_filters.get('direction', 'Toutes'),
        'periode_recrutement': shared_filters.get('periode_recrutement', 'Toutes'),
        # Ne pas filtrer par période de demande dans la section Clôtures
        'periode_demande': 'Toutes'
    }

    # Créer deux cartes expandables principales (comme dans Home.py)
    with st.expander("📋 **DEMANDES DE RECRUTEMENT**", expanded=False):
        create_demandes_recrutement_tab(df_recrutement, filters_demandes)
    
    with st.expander("🎯 **RECRUTEMENTS CLÔTURÉS**", expanded=False):
        create_recrutements_clotures_tab(df_recrutement, filters_clotures)


def calculate_weekly_metrics(df_recrutement):
    """Calcule les métriques hebdomadaires basées sur les vraies données"""
    if df_recrutement is None or len(df_recrutement) == 0:
        return {}
    
    # Obtenir la date de reporting (utiliser st.session_state si défini)
    reporting_date = st.session_state.get('reporting_date', None)
    if reporting_date is None:
        today = datetime.now()
    else:
        if isinstance(reporting_date, datetime):
            today = reporting_date
        else:
            today = datetime.combine(reporting_date, datetime.min.time())
    # start_of_week = Monday (00:00) of the reporting_date's week
    start_of_week = datetime(year=today.year, month=today.month, day=today.day) - timedelta(days=today.weekday())  # Lundi de cette semaine à 00:00
    # Définir la semaine précédente (Lundi -> Vendredi). Exemple: si reporting_date=2025-10-15 (mercredi),
    # start_of_week = 2025-10-13 (lundi), previous_monday = 2025-10-06, previous_friday = 2025-10-10.
    previous_monday = start_of_week - timedelta(days=7)   # Lundi de la semaine précédente (00:00)
    previous_friday = start_of_week - timedelta(days=3)   # Vendredi de la semaine précédente (00:00)
    # exclusive upper bound for inclusive-Friday semantics: < previous_friday + 1 day
    previous_friday_exclusive = previous_friday + timedelta(days=1)
    
    # Définir les colonnes attendues avec des alternatives possibles
    date_reception_col = "Date de réception de la demande après validation de la DRH"
    date_integration_col = "Date d'intégration prévisionnelle"
    candidat_col = "Nom Prénom du candidat retenu yant accepté la promesse d'embauche"
    statut_col = "Statut de la demande"
    entite_col = "Entité demandeuse"
    
    # Créer une copie pour les calculs
    df = df_recrutement.copy()

    # Toujours exclure les demandes clôturées/annulées du compteur 'avant'
    # (la case UI a été supprimée : les demandes clôturées ne sont pas comptées)
    include_closed = False
    
    # Vérifier les colonnes disponibles
    available_columns = df.columns.tolist()
    
    # Chercher les colonnes similaires si les noms exacts n'existent pas
    def find_similar_column(target_col, available_cols):
        """Trouve une colonne similaire dans la liste disponible"""
        target_lower = target_col.lower()
        for col in available_cols:
            if col.lower() == target_lower:
                return col
        # Chercher des mots-clés
        if "date" in target_lower and "réception" in target_lower:
            for col in available_cols:
                if "date" in col.lower() and ("réception" in col.lower() or "reception" in col.lower() or "demande" in col.lower()):
                    return col
        elif "date" in target_lower and "intégration" in target_lower:
            for col in available_cols:
                if "date" in col.lower() and ("intégration" in col.lower() or "integration" in col.lower() or "entrée" in col.lower()):
                    return col
        elif "candidat" in target_lower and "retenu" in target_lower:
            for col in available_cols:
                if ("candidat" in col.lower() and "retenu" in col.lower()) or ("nom" in col.lower() and "prénom" in col.lower()):
                    return col
        elif "statut" in target_lower:
            for col in available_cols:
                if "statut" in col.lower() or "status" in col.lower():
                    return col
        elif "entité" in target_lower:
            for col in available_cols:
                if "entité" in col.lower() or "entite" in col.lower():
                    return col
        return None
    
    # Trouver les colonnes réelles
    real_date_reception_col = find_similar_column(date_reception_col, available_columns)
    real_date_integration_col = find_similar_column(date_integration_col, available_columns)
    # chercher une colonne d'acceptation du candidat (date d'acceptation de la promesse)
    possible_accept_cols = [
        "Date d'acceptation",
        "Date d'acceptation du candidat",
        "Date d'acceptation de la promesse",
        "Date d'acceptation promesse",
        "Date d'accept",
        "Date d'acceptation de la promesse d'embauche",
        "Date d'acceptation de la promesse d embauche",
        "Date d'acceptation du candidat",
        'Date d\'acceptation',
    ]
    real_accept_col = find_similar_column('Date d\'acceptation du candidat', available_columns)
    # fallback to integration date if no explicit acceptance date
    if real_accept_col is None:
        # try some common alternatives
        for alt in ['Date d\'acceptation', 'Date d\'acceptation de la promesse', "Date d'accept"]:
            c = find_similar_column(alt, available_columns)
            if c:
                real_accept_col = c
                break
    real_candidat_col = find_similar_column(candidat_col, available_columns)
    real_statut_col = find_similar_column(statut_col, available_columns)
    real_entite_col = find_similar_column(entite_col, available_columns)
    # detect Poste demandé column (used for optional entity-specific title filters)
    real_poste_col = find_similar_column('Poste demandé', available_columns) or find_similar_column('Poste', available_columns)
    
    # Si les colonnes essentielles n'existent pas, retourner vide
    if not real_entite_col:
        st.warning(f"⚠️ Colonne 'Entité' non trouvée. Colonnes disponibles: {', '.join(available_columns[:5])}...")
        return {}
    
    if not real_statut_col:
        st.warning(f"⚠️ Colonne 'Statut' non trouvée. Colonnes disponibles: {', '.join(available_columns[:5])}...")
        return {}
    
    # Convertir les dates si les colonnes existent (gardes robustes)
    if real_date_reception_col and real_date_reception_col in df.columns:
        df[real_date_reception_col] = pd.to_datetime(df[real_date_reception_col], errors='coerce')
    else:
        real_date_reception_col = None
    if real_date_integration_col and real_date_integration_col in df.columns:
        df[real_date_integration_col] = pd.to_datetime(df[real_date_integration_col], errors='coerce')
    else:
        real_date_integration_col = None
    if real_accept_col and real_accept_col in df.columns:
        # Use robust mixed-date parser (handles Excel serials and mixed formats)
        try:
            df[real_accept_col] = _parse_mixed_dates(df[real_accept_col])
        except Exception:
            # Last-resort fallback
            df[real_accept_col] = pd.to_datetime(df[real_accept_col], errors='coerce')
    else:
        real_accept_col = None

    # Normalisation et mots-clés de statuts fermés (utiles dans plusieurs blocs)
    closed_keywords = ['cloture', 'clôture', 'annule', 'annulé', 'depriorise', 'dépriorisé', 'desistement', 'désistement', 'annul', 'reject', 'rejett']
    # Optional: entity-specific title inclusion lists. If an entity is present here,
    # only postes matching the provided list will be counted in 'en_cours' for that entity.
    SPECIAL_TITLE_FILTERS = {
        'TGCC': [
            'CHEF DE PROJETS',
            'INGENIEUR TRAVAUX',
            'CONDUCTEUR TRAVAUX SENIOR (ou Ingénieur Junior)',
            'CONDUCTEUR TRAVAUX',
            'INGENIEUR TRAVAUX JUNIOR',
            'RESPONSABLE QUALITE',
            'CHEF DE CHANTIER',
            'METREUR',
            'RESPONSABLE HSE',
            'SUPERVISEUR HSE',
            'ANIMATEUR HSE',
            'DIRECTEUR PROJETS',
            'RESPONSABLE ADMINISTRATIF ET FINANCIER',
            'RESPONSABLE MAINTENANCE',
            'RESPONSABLE ENERGIE INDUSTRIELLE',
            'RESPONSABLE CYBER SECURITE',
            'RESPONSABLE VRD',
            'RESPONSABLE ACCEUIL',
            'RESPONSABLE ETUDES',
            'TECHNICIEN SI',
            'RESPONSABLE GED & ARCHIVAGE',
            'ARCHIVISTE SENIOR',
            'ARCHIVISTE JUNIOR',
            'TOPOGRAPHE'
        ]
    }
    
    # Calculer les métriques par entité
    # Normaliser les noms d'entités pour éviter les doublons (ex: "TG WOOD" vs "TG WOOD ")
    if real_entite_col:
        df[real_entite_col] = df[real_entite_col].astype(str).str.strip().str.upper()
    
    entites = df[real_entite_col].dropna().unique()
    metrics_by_entity = {}
    
    for entite in entites:
        df_entite = df[df[real_entite_col] == entite]
        # Définitions temporelles par entité (basées sur la semaine précédente Lundi->Vendredi)
        # previous_monday / previous_friday_exclusive sont définis en haut de la fonction

        # 1. Nb postes ouverts avant début semaine
        # Doit compter toutes les lignes dont 'Date de réception ...' est STRICTEMENT antérieure
        # au vendredi de la semaine précédente (i.e. < previous_friday)
        postes_avant = 0
        if real_date_reception_col and real_date_reception_col in df_entite.columns:
            mask_date_avant = df_entite[real_date_reception_col] < previous_monday
            mask_not_closed = None
            if not include_closed and real_statut_col and real_statut_col in df_entite.columns:
                mask_not_closed = ~df_entite[real_statut_col].fillna("").astype(str).apply(lambda s: any(k in _norm(s) for k in closed_keywords))
            # Exclure les lignes avec une date d'acceptation du candidat antérieure à la semaine précédente
            mask_old_accept = None
            if real_accept_col and real_accept_col in df_entite.columns:
                mask_old_accept = (df_entite[real_accept_col] < previous_monday)
            mask_avant = mask_date_avant
            if mask_not_closed is not None:
                mask_avant = mask_avant & mask_not_closed
            if mask_old_accept is not None:
                mask_avant = mask_avant & (~mask_old_accept)
            postes_avant = int(df_entite[mask_avant].shape[0])
        else:
            postes_avant = 0

        # 2. Nb nouveaux postes ouverts cette semaine = demandes validées par la DRH
        # dans la semaine précédente (du lundi au vendredi inclus).
        nouveaux_postes = 0
        if real_date_reception_col and real_date_reception_col in df_entite.columns:
            mask_nouveaux = (df_entite[real_date_reception_col] >= previous_monday) & (df_entite[real_date_reception_col] < previous_friday_exclusive)
            nouveaux_postes = int(df_entite[mask_nouveaux].shape[0])
        else:
            nouveaux_postes = 0

        # 3. Nb postes pourvus cette semaine: compter les acceptations du candidat
        # dont la date d'acceptation est dans la même fenêtre (previous_monday..previous_friday)
        postes_pourvus = 0
        mask_has_name = None
        if real_candidat_col and real_candidat_col in df_entite.columns:
            mask_has_name = df_entite[real_candidat_col].notna() & (df_entite[real_candidat_col].astype(str).str.strip() != '')

        # préparer un masque de statut 'En cours' si la colonne statut existe (réutilisé plus bas)
        mask_status_en_cours = None
        if real_statut_col and real_statut_col in df_entite.columns:
            mask_status_en_cours = df_entite[real_statut_col].fillna("").astype(str).apply(lambda s: 'en cours' in _norm(s) or 'encours' in _norm(s))

        if real_accept_col and real_accept_col in df_entite.columns:
            mask_accept_prev_week = (df_entite[real_accept_col] >= previous_monday) & (df_entite[real_accept_col] < previous_friday_exclusive)
            if mask_has_name is not None:
                postes_pourvus = int(df_entite[mask_accept_prev_week & mask_has_name].shape[0])
            else:
                postes_pourvus = int(df_entite[mask_accept_prev_week].shape[0])
        else:
            # fallback: try using integration date as proxy for pourvus in the same window
            if real_date_integration_col and real_date_integration_col in df_entite.columns:
                mask_integ_prev_week = (df_entite[real_date_integration_col] >= previous_monday) & (df_entite[real_date_integration_col] < previous_friday_exclusive)
                if mask_has_name is not None:
                    postes_pourvus = int(df_entite[mask_integ_prev_week & mask_has_name].shape[0])
                else:
                    postes_pourvus = int(df_entite[mask_integ_prev_week].shape[0])
            else:
                postes_pourvus = 0

        # 4. Nb postes en cours cette semaine
        # Par défaut on calcule une formule de secours (nouveaux + avant - pourvus)
        postes_en_cours_formula = nouveaux_postes + postes_avant - postes_pourvus
        if postes_en_cours_formula is None:
            postes_en_cours_formula = 0

        # Si la colonne Statut existe, on suit la règle métier stricte demandée :
        # "Postes en cours" = lignes avec statut 'En cours' ET sans valeur dans
        # la colonne 'Nom Prénom du candidat retenu...' (donc pas d'acceptation).
        # Si la colonne Statut est absente, on retombe sur la formule de secours.
        postes_en_cours = int(postes_en_cours_formula)
        postes_en_cours_status = 0
        if mask_status_en_cours is not None:
            mask_has_name_local = None
            if real_candidat_col and real_candidat_col in df_entite.columns:
                mask_has_name_local = df_entite[real_candidat_col].notna() & (df_entite[real_candidat_col].astype(str).str.strip() != '')

            # Exclure les lignes où la date d'acceptation du candidat est antérieure à la semaine précédente
            mask_old_accept = None
            if real_accept_col and real_accept_col in df_entite.columns:
                mask_old_accept = (df_entite[real_accept_col] < previous_monday)
            # Règle stricte : "Postes en cours" = statut 'En cours' ET sans candidat ET pas d'acceptation ancienne
            if mask_has_name_local is not None:
                mask_sourcing = mask_status_en_cours & (~mask_has_name_local)
            else:
                mask_sourcing = mask_status_en_cours
            if mask_old_accept is not None:
                mask_sourcing = mask_sourcing & (~mask_old_accept)

            postes_en_cours_status = int(df_entite[mask_sourcing].shape[0])
            postes_en_cours = int(postes_en_cours_status)

        metrics_by_entity[entite] = {
            'avant': postes_avant,
            'nouveaux': nouveaux_postes,
            'pourvus': postes_pourvus,
            'en_cours': postes_en_cours,
            # ajouter info additionnelle utile pour debug/UI
            'en_cours_status_count': postes_en_cours_status
        }
    
    # Créer aussi table_data pour l'export PowerPoint
    excluded_entities = {'BESIX-TGCC', 'DECO EXCELL', 'TG PREFA'}
    metrics_included = {k: v for k, v in metrics_by_entity.items() if k not in excluded_entities}
    
    # Calculer les totaux
    total_avant = sum(data['avant'] for data in metrics_included.values())
    total_nouveaux = sum(data['nouveaux'] for data in metrics_included.values())
    total_pourvus = sum(data['pourvus'] for data in metrics_included.values())
    total_en_cours = sum(data['en_cours'] for data in metrics_included.values())
    total_en_cours_status = sum(data.get('en_cours_status_count', 0) for data in metrics_included.values())
    
    # Préparer les données pour le tableau
    table_data = []
    for entite, data in metrics_included.items():
        table_data.append({
            'Entité': entite,  # Sans les logos pour le PowerPoint
            'Nb postes ouverts avant début semaine': data['avant'] if data['avant'] > 0 else 0,
            'Nb nouveaux postes ouverts cette semaine': data['nouveaux'] if data['nouveaux'] > 0 else 0,
            'Nb postes pourvus cette semaine': data['pourvus'] if data['pourvus'] > 0 else 0,
            'Nb postes en cours cette semaine (sourcing)': data['en_cours'] if data['en_cours'] > 0 else 0
        })
    
    # Ajouter la ligne TOTAL
    table_data.append({
        'Entité': 'TOTAL',
        'Nb postes ouverts avant début semaine': total_avant,
        'Nb nouveaux postes ouverts cette semaine': total_nouveaux,
        'Nb postes pourvus cette semaine': total_pourvus,
        'Nb postes en cours cette semaine (sourcing)': total_en_cours
    })
    
    return {
        'metrics_by_entity': metrics_by_entity,
        'table_data': table_data,
        'totals': {
            'avant': total_avant,
            'nouveaux': total_nouveaux,
            'pourvus': total_pourvus,
            'en_cours': total_en_cours
        }
    }

def create_weekly_report_tab(df_recrutement=None):
    """Onglet Reporting Hebdomadaire (simplifié)

    Cette version affiche les KPI calculés par calculate_weekly_metrics
    en respectant la `st.session_state['reporting_date']` si fournie.
    """
    st.header("📅 Reporting Hebdomadaire : Chiffres Clés de la semaine")

    # Note: les demandes clôturées sont exclues du compteur 'avant' par défaut.
    # La case pour inclure les demandes clôturées a été supprimée.

    # Calculer les métriques
    if df_recrutement is not None and len(df_recrutement) > 0:
        try:
            metrics_result = calculate_weekly_metrics(df_recrutement)
            # Extraire metrics_by_entity de la nouvelle structure
            metrics = metrics_result.get('metrics_by_entity', {}) if isinstance(metrics_result, dict) else {}
        except Exception as e:
            st.error(f"⚠️ Erreur lors du calcul des métriques: {e}")
            metrics = {}
    else:
        metrics = {}

    # Exclure certaines entités (Besix, DECO EXCELL et TG PREFA) de l'affichage et des totaux
    excluded_entities = set(['BESIX-TGCC', 'DECO EXCELL', 'TG PREFA'])
    metrics_included = {e: m for e, m in metrics.items() if e not in excluded_entities}

    total_avant = sum(m.get('avant', 0) for m in metrics_included.values())
    total_nouveaux = sum(m.get('nouveaux', 0) for m in metrics_included.values())
    total_pourvus = sum(m.get('pourvus', 0) for m in metrics_included.values())
    total_en_cours = sum(m.get('en_cours', 0) for m in metrics_included.values())
    # total lines with statut 'En cours' (may include those with candidate)
    total_en_cours_status = sum(m.get('en_cours_status_count', 0) for m in metrics_included.values())

    # KPI cards (styled and centered like other KPI cards)
    kpi_cards_html = render_generic_metrics([
        ("Postes en cours (sourcing)", total_en_cours, "#1f77b4"),
        ("Postes pourvus cette semaine", total_pourvus, "#2ca02c"),
        ("Nouveaux postes ouverts", total_nouveaux, "#ff7f0e"),
        ("Total postes ouverts avant la semaine", total_avant, "#6f42c1")
    ])
    st.markdown(kpi_cards_html, unsafe_allow_html=True)

    # Tableau récapitulatif par entité (HTML personnalisé, rendu centralisé)
        st.markdown('Avant de générer le reporting, assurez-vous que tous les recrutements sont saisis, que la colonne «Status» de la demande est renseignée et que la colonne «TG Hire» reflète l\'état du recrutement (Nouvelle demande, Sourcing, Signature DRH, ...).<div style="color:red; font-weight:700; margin-top:6px;">NB : au fur et à mesure de l\'évolution des recrutements, mettez à jour leur état, sinon le reporting ne sera pas fiable.</div>', unsafe_allow_html=True)

        st.markdown("""
1.  **Chargement des Données** : 

    *   Allez dans l'onglet **"📂 Upload & Téléchargement"**.
    *   **Option A** : Cliquez sur le bouton rouge **"🔁 Synchroniser depuis Google Sheets"** pour récupérer les données les plus récentes.
    *   **Option B** : Glissez-déposez votre fichier Excel de recrutement dans la zone de chargement.

        # Préparer les données pour le HTML
        table_data = []
        for entite, data in metrics_included.items():
            table_data.append({
                'Entité': get_entity_display(entite),
                'Nb postes ouverts avant début semaine': data['avant'] if data['avant'] > 0 else '-',
                'Nb nouveaux postes ouverts cette semaine': data['nouveaux'] if data['nouveaux'] > 0 else '-',
                'Nb postes pourvus cette semaine': data['pourvus'] if data['pourvus'] > 0 else '-',
                "Nb postes statut 'En cours' (total)": data.get('en_cours_status_count', 0) if data.get('en_cours_status_count', 0) > 0 else '-',
                'Nb postes en cours cette semaine (sourcing)': data['en_cours'] if data['en_cours'] > 0 else '-'
            })

        # Ajouter la ligne de total
        table_data.append({
            'Entité': '<div style="text-align: center; font-weight: bold;">TOTAL</div>',
            'Nb postes ouverts avant début semaine': f'**{total_avant}**',
            'Nb nouveaux postes ouverts cette semaine': f'**{total_nouveaux}**',
            'Nb postes pourvus cette semaine': f'**{total_pourvus}**',
            "Nb postes statut 'En cours' (total)": f'**{total_en_cours_status}**',
            'Nb postes en cours cette semaine (sourcing)': f'**{total_en_cours}**'
        })

        # HTML + CSS (RETOUR À LA VERSION ORIGINALE AVEC LOGOS)
        st.markdown("""
        <style>
        .table-container {
            display: flex;
            justify-content: center;
            width: 100%;
            margin: 15px 0;
        }
        .custom-table {
            border-collapse: collapse;
            font-family: Arial, sans-serif;
            box-shadow: 0 1px 4px rgba(0,0,0,0.1);
            width: 75%;
            margin: 0 auto;
        }
        .custom-table th {
            background-color: #9C182F !important;
            color: white !important;
            font-weight: bold !important;
            text-align: center !important;
            padding: 6px 4px !important;
            border: 1px solid white !important;
            font-size: 1.1em;
            line-height: 1.3;
            white-space: normal !important;
            word-wrap: break-word !important;
            max-width: 150px;
        }
        .custom-table td {
            text-align: center !important;
            padding: 6px 4px !important;
            border: 1px solid #ddd !important;
            background-color: white !important;
            font-size: 1.1em;
            line-height: 1.2;
            font-weight: 500;
        }
        .custom-table .entity-cell {
            text-align: center !important;
            padding: 4px 2px !important;
            font-weight: 600;
            min-width: 80px;
            max-width: 100px;
        }
        .custom-table .total-row td {
            background-color: #9C182F !important;
            color: white !important;
            font-weight: bold !important;
            border: 1px solid #9C182F !important;
        }
        </style>
        """, unsafe_allow_html=True)

        # Construire le tableau HTML
        html_table = '<div class="table-container">'
        html_table += '<table class="custom-table">'
        html_table += '<thead><tr>'
        html_table += '<th>Entité</th>'
        html_table += '<th>Nb postes ouverts avant début semaine</th>'
        html_table += '<th>Nb nouveaux postes ouverts cette semaine</th>'
        html_table += '<th>Nb postes pourvus cette semaine</th>'
        html_table += '<th>Nb postes en cours cette semaine (sourcing)</th>'
        html_table += '</tr></thead>'
        html_table += '<tbody>'

        # Ajouter les lignes de données (toutes sauf la dernière qui est TOTAL)
        data_rows = table_data[:-1]
        for row in data_rows:
            html_table += '<tr>'
            html_table += f'<td class="entity-cell">{row["Entité"]}</td>'
            html_table += f'<td>{row["Nb postes ouverts avant début semaine"]}</td>'
            html_table += f'<td>{row["Nb nouveaux postes ouverts cette semaine"]}</td>'
            html_table += f'<td>{row["Nb postes pourvus cette semaine"]}</td>'
            html_table += f'<td>{row["Nb postes en cours cette semaine (sourcing)"]}</td>'
            html_table += '</tr>'

        # Ligne TOTAL (la dernière)
        total_row = table_data[-1]
        html_table += '<tr class="total-row">'
        html_table += f'<td class="entity-cell">TOTAL</td>'
        html_table += f'<td>{total_row["Nb postes ouverts avant début semaine"].replace("**", "")}</td>'
        html_table += f'<td>{total_row["Nb nouveaux postes ouverts cette semaine"].replace("**", "")}</td>'
        html_table += f'<td>{total_row["Nb postes pourvus cette semaine"].replace("**", "")}</td>'
        html_table += f'<td>{total_row["Nb postes en cours cette semaine (sourcing)"].replace("**", "")}</td>'
        html_table += '</tr>'
        html_table += '</tbody></table></div>'

        st.markdown(html_table, unsafe_allow_html=True)

        # --- Tableau : Recrutements en cours par recruteur (juste après 'Besoins en Cours par Entité')
        try:
            if df_recrutement is not None and 'Colonne TG Hire' in df_recrutement.columns:
                # On veut un tableau avec une colonne par statut: Nouvelle demande, Sourcing, Shortlisté, Signature DRH
                wanted_statuses = ['Nouvelle demande', 'Sourcing', 'Shortlisté', 'Signature DRH']
                # Construire un pivot complet par recruteur (colonnes = Colonne TG Hire demandées)
                # Le TOTAL affiché est le nombre de lignes pour lesquelles 'Statut de la demande' == 'En cours'
                # Identifier la colonne recruteur de façon robuste (sur le DF complet)
                recruteur_col = next((c for c in df_recrutement.columns if 'responsable' in c.lower() and 'traitement' in c.lower()), None)
                if not recruteur_col:
                    recruteur_col = next((c for c in df_recrutement.columns if 'recruteur' in c.lower() or 'responsable' in c.lower()), None)

                if recruteur_col:
                    # Pivot sur l'ensemble des données pour récupérer les colonnes d'intérêt
                    pivot = pd.crosstab(df_recrutement[recruteur_col].fillna('').astype(str).str.strip(), df_recrutement['Colonne TG Hire'])
                    # S'assurer de l'ordre des colonnes et présence de toutes (ajoute 0 si manquante)
                    for s in wanted_statuses:
                        if s not in pivot.columns:
                            pivot[s] = 0
                    pivot = pivot[wanted_statuses]

                    # Supprimer les lignes sans recruteur explicite et exclure certains recruteurs indésirables
                    pivot.index = pivot.index.astype(str)
                    pivot = pivot[~pivot.index.str.strip().str.lower().isin(['', '(sans)', 'nan', 'none'])]
                    exclude_list = {n.lower() for n in [
                        'Bouchra AJBILOU','Bouchra AOUISSE','Ghita LAKHDAR',
                        'Reda Berrada','Reda Mohamed BERRADA','Saad FATI'
                    ]}
                    pivot = pivot[~pivot.index.str.lower().isin(exclude_list)]

                    # Calculer la colonne demandée: Total = Nouvelle demande + Sourcing + Shortlisté + Signature DRH
                    cols_sum = [c for c in wanted_statuses if c in pivot.columns]
                    if cols_sum:
                        pivot['Total'] = pivot[cols_sum].sum(axis=1).clip(lower=0).astype(int)
                    else:
                        pivot['Total'] = 0

                    # Construire le HTML du tableau
                    html_rec = '<div class="table-container" style="margin-top:8px;">'
                    html_rec += '<table class="custom-table" style="width:60%; margin:0;">'
                    # Header
                    html_rec += '<thead><tr><th>Recruteur</th>'
                    for s in wanted_statuses:
                        html_rec += f'<th>{s}</th>'
                    html_rec += '<th>Total</th></tr></thead><tbody>'

                    # Lignes par recruteur
                    for rec, row in pivot.iterrows():
                        html_rec += '<tr>'
                        html_rec += f'<td class="entity-cell">{rec}</td>'
                        for s in wanted_statuses:
                            html_rec += f'<td>{int(row[s])}</td>'
                        html_rec += f'<td>{int(row["Total"])}</td>'
                        html_rec += '</tr>'

                    html_rec += '</tbody></table></div>'

                    st.markdown('<div style="font-family:Arial,sans-serif; font-size:1.15em; font-weight:700; margin-top:12px;">📋 Recrutements en cours par recruteur</div>', unsafe_allow_html=True)
                    st.markdown(html_rec, unsafe_allow_html=True)
                    
                    # Debug pour le tableau recruteur
                    with st.expander("🔍 Debug - Détails des lignes (Recrutements par recruteur)", expanded=False):
                        try:
                            st.markdown("**Lignes contribuant au tableau (statut 'En cours') par statut:**")
                            # Filtrer uniquement les lignes avec statut 'En cours'
                            df_rec_debug = df_recrutement[df_recrutement['Statut de la demande'] == 'En cours'].copy()
                            # Ajouter colonnes de contribution
                            df_rec_debug['contrib_Nouvelle_demande'] = df_rec_debug['Colonne TG Hire'] == 'Nouvelle demande'
                            df_rec_debug['contrib_Sourcing'] = df_rec_debug['Colonne TG Hire'] == 'Sourcing'
                            df_rec_debug['contrib_Shortlisté'] = df_rec_debug['Colonne TG Hire'] == 'Shortlisté'
                            df_rec_debug['contrib_Signature_DRH'] = df_rec_debug['Colonne TG Hire'] == 'Signature DRH'
                            cols_show = ['Entité demandeuse', recruteur_col, 'Poste demandé', 'Colonne TG Hire', 'contrib_Nouvelle_demande', 'contrib_Sourcing', 'contrib_Shortlisté', 'contrib_Signature_DRH']
                            cols_avail = [c for c in cols_show if c in df_rec_debug.columns]
                            st.dataframe(df_rec_debug[cols_avail].reset_index(drop=True), use_container_width=True, hide_index=True)
                        except Exception as e:
                            st.write(f"Erreur debug: {e}")
        except Exception:
            pass
    else:
        # Affichage par défaut si pas de metrics (avec logos pour démo)
        logos_dict = load_all_logos_b64()
        tgcc_logo = get_entity_display_html_with_logo('TGCC', logos_dict)
        tgem_logo = get_entity_display_html_with_logo('TGEM', logos_dict)
        
        default_html = f"""
<div class="table-container">
    <table class="custom-table">
        <thead>
            <tr>
                <th>Entité</th>
                <th>Nb postes ouverts avant début semaine</th>
                <th>Nb nouveaux postes ouverts cette semaine</th>
                <th>Nb postes pourvus cette semaine</th>
                <th>Nb postes en cours cette semaine</th>
            </tr>
        </thead>
        <tbody>
            <tr>
                <td class="entity-cell">{tgcc_logo}</td>
                <td>19</td>
                <td>12</td>
                <td>5</td>
                <td>26</td>
            </tr>
            <tr>
                <td class="entity-cell">{tgem_logo}</td>
                <td>2</td>
                <td>2</td>
                <td>0</td>
                <td>4</td>
            </tr>
            <tr class="total-row">
                <td class="entity-cell">TOTAL</td>
                <td>21</td>
                <td>14</td>
                <td>5</td>
                <td>30</td>
            </tr>
        </tbody>
    </table>
</div>
"""
        st.markdown(default_html, unsafe_allow_html=True)

    # Section Debug (expandable): montrer les lignes et pourquoi elles sont comptées
    with st.expander("🔍 Debug - Détails des lignes", expanded=False):
        st.markdown("""
        **Explication des KPI du tableau 'Besoins en Cours par Entité'**
        - **Nb postes ouverts avant début semaine** : demandes dont la date de réception est antérieure au lundi de la semaine de reporting, et qui ne sont pas clôturées/annulées. Les lignes avec une date d'acceptation du candidat antérieure à la semaine précédente sont également exclues.
        - **Nb nouveaux postes ouverts cette semaine** : demandes validées par la DRH entre le lundi et le vendredi de la semaine précédente.
        - **Nb postes pourvus cette semaine** : postes pour lesquels un candidat a accepté (ou date d'intégration) dans la même fenêtre temporelle.
        - **Nb postes en cours cette semaine (sourcing)** : calculé comme (nouveaux + avant - pourvus), ou plus précisément comme les lignes avec statut "En cours" et sans candidat retenu. Les lignes avec une date d'acceptation du candidat antérieure à la semaine précédente sont également exclues, même si le statut est "En cours".
        - **Nb postes statut 'En cours' (total)** : nombre de lignes avec statut "En cours" (peut inclure celles avec candidat).
        """)
        try:
            df_debug = df_recrutement.copy() if df_recrutement is not None else pd.DataFrame()
            if not df_debug.empty:
                cols = df_debug.columns.tolist()

                def find_similar_column(target_col, available_cols):
                    target_lower = target_col.lower()
                    for col in available_cols:
                        if col.lower() == target_lower:
                            return col
                    if "date" in target_lower and "réception" in target_lower:
                        for col in available_cols:
                            if "date" in col.lower() and ("réception" in col.lower() or "reception" in col.lower() or "demande" in col.lower()):
                                return col
                    if "date" in target_lower and "intégration" in target_lower:
                        for col in available_cols:
                            if "date" in col.lower() and ("intégration" in col.lower() or "integration" in col.lower() or "entrée" in col.lower()):
                                return col
                    # Prefer columns explicitly mentioning the candidate / promesse
                    if "candidat" in target_lower or "promesse" in target_lower or "accept" in target_lower:
                        for col in available_cols:
                            lc = col.lower()
                            if ("candidat" in lc and "retenu" in lc) or ("accept" in lc and "candidat" in lc) or ("promesse" in lc and "candidat" in lc):
                                return col
                        # fallback to any column containing 'candidat'
                        for col in available_cols:
                            if 'candidat' in col.lower():
                                return col
                    # As a last resort, match generic 'nom'/'prénom' columns
                    if "candidat" in target_lower and "retenu" in target_lower:
                        for col in available_cols:
                            if ("nom" in col.lower() and "pr" in col.lower()) or ("nom" in col.lower() and "prenom" in col.lower()):
                                return col
                    if "statut" in target_lower:
                        for col in available_cols:
                            if "statut" in col.lower() or "status" in col.lower():
                                return col
                    if "entité" in target_lower or "entite" in target_lower:
                        for col in available_cols:
                            if "entité" in col.lower() or "entite" in col.lower():
                                return col
                    return None

                date_reception_col = "Date de réception de la demande après validation de la DRH"
                date_integration_col = "Date d'intégration prévisionnelle"
                candidat_col = "Nom Prénom du candidat retenu yant accepté la promesse d'embauche"

                real_date_reception_col = find_similar_column(date_reception_col, cols)
                real_date_integration_col = find_similar_column(date_integration_col, cols)
                real_accept_col = None
                for alt in ['Date d\'acceptation du candidat','Date d\'acceptation','Date d\'acceptation de la promesse',"Date d'accept"]:
                    c = find_similar_column(alt, cols)
                    if c:
                        real_accept_col = c
                        break
                real_candidat_col = candidat_col if candidat_col in cols else find_similar_column(candidat_col, cols)
                real_statut_col = find_similar_column('Statut de la demande', cols)
                real_entite_col = find_similar_column('Entité demandeuse', cols) or find_similar_column('Entité', cols)

                rd = st.session_state.get('reporting_date', None)
                if rd is None:
                    today = datetime.now()
                else:
                    today = rd if isinstance(rd, datetime) else datetime.combine(rd, datetime.min.time())
                
                # Re-create the same date ranges as in `calculate_weekly_metrics` for consistency
                start_of_week = datetime(year=today.year, month=today.month, day=today.day) - timedelta(days=today.weekday())
                previous_monday = start_of_week - timedelta(days=7)
                previous_friday_exclusive = (start_of_week - timedelta(days=3))
                previous_friday_exclusive = previous_friday_exclusive + timedelta(days=1)  # Vendredi + 1 jour

                if real_date_reception_col and real_date_reception_col in df_debug.columns:
                    df_debug[real_date_reception_col] = pd.to_datetime(df_debug[real_date_reception_col], errors='coerce')
                if real_date_integration_col and real_date_integration_col in df_debug.columns:
                    df_debug[real_date_integration_col] = pd.to_datetime(df_debug[real_date_integration_col], errors='coerce')
                if real_accept_col and real_accept_col in df_debug.columns:
                    try:
                        df_debug[real_accept_col] = _parse_mixed_dates(df_debug[real_accept_col])
                    except Exception:
                        df_debug[real_accept_col] = pd.to_datetime(df_debug[real_accept_col], errors='coerce')

                def _local_norm(x):
                    if pd.isna(x) or x is None:
                        return ''
                    s = str(x)
                    s = unicodedata.normalize('NFKD', s)
                    s = ''.join(ch for ch in s if not unicodedata.combining(ch))
                    return s.lower().strip()

                # Masks for contribution flags, aligned with `calculate_weekly_metrics`
                mask_avant_semaine = pd.Series(False, index=df_debug.index)
                mask_nouveaux_semaine = pd.Series(False, index=df_debug.index)
                mask_pourvus_semaine = pd.Series(False, index=df_debug.index)
                mask_status_en_cours = pd.Series(False, index=df_debug.index)
                mask_has_name = pd.Series(False, index=df_debug.index)
                closed_keywords = ['cloture', 'clôture', 'annule', 'annulé', 'depriorise', 'dépriorisé', 'desistement', 'désistement', 'annul', 'reject', 'rejett']

                if real_date_reception_col and real_date_reception_col in df_debug.columns:
                    # "Avant": reception date is before the start of the previous week
                    mask_avant_semaine = df_debug[real_date_reception_col] < previous_monday
                    # "Nouveaux": reception date is within the previous week (Mon-Fri)
                    mask_nouveaux_semaine = (df_debug[real_date_reception_col] >= previous_monday) & (df_debug[real_date_reception_col] < previous_friday_exclusive)

                if real_statut_col and real_statut_col in df_debug.columns:
                    mask_status_en_cours = df_debug[real_statut_col].fillna("").astype(str).apply(lambda s: ('en cours' in _local_norm(s)) or ('encours' in _local_norm(s)))
                    mask_not_closed = ~df_debug[real_statut_col].fillna("").astype(str).apply(lambda s: any(k in _local_norm(s) for k in closed_keywords))
                    mask_avant_semaine = mask_avant_semaine & mask_not_closed


                if real_candidat_col and real_candidat_col in df_debug.columns:
                    mask_has_name = df_debug[real_candidat_col].notna() & (df_debug[real_candidat_col].astype(str).str.strip() != '')

                if real_accept_col and real_accept_col in df_debug.columns:
                    # "Pourvus": acceptance date is within the previous week (Mon-Fri)
                    mask_pourvus_semaine = (df_debug[real_accept_col] >= previous_monday) & (df_debug[real_accept_col] < previous_friday_exclusive)
                
                contributes_avant = mask_avant_semaine
                contributes_nouveaux = mask_nouveaux_semaine
                contributes_pourvus = mask_pourvus_semaine & mask_has_name
                # contrib_en_cours: statut 'En cours'
                contributes_en_cours = mask_status_en_cours

                # For the debug view we want the contributors to exactly reflect the
                # 'Besoins en Cours' table. We will show all rows that contribute to *any* of the KPIs.
                any_contrib = contributes_avant | contributes_nouveaux | contributes_pourvus | contributes_en_cours

                df_selected = df_debug[any_contrib].copy()

                display_cols = []
                if real_entite_col and real_entite_col in df_selected.columns:
                    display_cols.append(real_entite_col)
                if real_candidat_col and real_candidat_col in df_selected.columns:
                    display_cols.append(real_candidat_col)
                if real_statut_col and real_statut_col in df_selected.columns:
                    display_cols.append(real_statut_col)
                if real_date_reception_col and real_date_reception_col in df_selected.columns:
                    display_cols.append(real_date_reception_col)
                if real_accept_col and real_accept_col in df_selected.columns:
                    display_cols.append(real_accept_col)

                df_out = df_selected[display_cols].copy() if display_cols else df_selected.copy()
                # mark if the row matches any SPECIAL_TITLE_FILTERS for its entity (useful for TGCC overrides)
                # compute a boolean flag indicating whether the row's title matches
                # any entry in the SPECIAL_TITLE_FILTERS for that entity.
                # Use safe fallbacks if some variables are not present in this scope.
                try:
                    # try to access SPECIAL_TITLE_FILTERS from the module scope
                    st_special_filters = globals().get('SPECIAL_TITLE_FILTERS', None)
                    # if real_poste_col isn't defined in this block, try to discover a 'Poste' column
                    rp = globals().get('real_poste_col', None)
                    if not rp:
                        # conservative search for a poste-like column name in df_selected
                        for c in df_selected.columns:
                            if 'poste' in c.lower() or 'title' in c.lower():
                                rp = c
                                break

                    def _matches_special_title(row):
                        ent = row.get(real_entite_col, '') if real_entite_col in row.index else (row.get('Entité demandeuse', '') or row.get('Entité', ''))
                        titre = ''
                        if rp and rp in row.index:
                            titre = row.get(rp, '')
                        else:
                            titre = row.get('Poste demandé', '') or row.get('Poste demandé ', '') or ''

                        if not ent or not st_special_filters:
                            return False
                        filter_list = st_special_filters.get(ent, [])
                        if not filter_list:
                            return False
                        tnorm = str(titre).strip().upper()
                        return any(tnorm == s for s in filter_list)

                    df_out['special_title_match'] = df_selected.apply(_matches_special_title, axis=1)
                except Exception:
                    df_out['special_title_match'] = False
                df_out['contrib_avant'] = pd.Series(contributes_avant.loc[df_out.index], dtype='bool')
                df_out['contrib_nouveaux'] = pd.Series(contributes_nouveaux.loc[df_out.index], dtype='bool')
                df_out['contrib_pourvus'] = pd.Series(contributes_pourvus.loc[df_out.index], dtype='bool')
                df_out['contrib_en_cours'] = pd.Series(contributes_en_cours.loc[df_out.index], dtype='bool')

                # Ensure the candidate value is present in a predictable column named 'candidate_value'
                try:
                    if real_candidat_col and real_candidat_col in df_selected.columns:
                        df_out['candidate_value'] = df_selected[real_candidat_col].fillna('').astype(str)
                    else:
                        df_out['candidate_value'] = ''
                except Exception:
                    df_out['candidate_value'] = ''

                for dc in [real_date_reception_col, real_accept_col, real_date_integration_col]:
                    if dc and dc in df_out.columns:
                        try:
                            # Safely format each date value in the column using .map
                            df_out[dc] = pd.to_datetime(df_out[dc], errors='coerce')
                            df_out[dc] = df_out[dc].map(lambda x: x.strftime('%d/%m/%Y') if pd.notna(x) and hasattr(x, 'strftime') else 'N/A')
                        except Exception:
                            pass

                # Allow user to choose display mode: KPI contributors or all rows with statut 'En cours'
                display_mode = st.radio("Mode d'affichage:", ["Contributeurs 'Postes en cours'", "Toutes lignes statut 'En cours'"], index=0)

                if display_mode == "Toutes lignes statut 'En cours'":
                    if real_statut_col and real_statut_col in df_debug.columns:
                        df_status = df_debug[mask_status_en_cours].copy()
                        
                        # Créer la colonne candidate_value avec la même logique que le mode contributeurs
                        if real_candidat_col and real_candidat_col in df_status.columns:
                            # Supprimer la colonne existante pour éviter les doublons
                            if real_candidat_col in df_status.columns:
                                df_status = df_status.drop(columns=[real_candidat_col])
                            df_status['candidate_value'] = df_debug.loc[df_status.index, real_candidat_col].fillna('').astype(str)
                        else:
                            df_status['candidate_value'] = ''
                        
                        # Formater les dates en jj/mm/aaaa
                        if real_date_reception_col and real_date_reception_col in df_status.columns:
                            df_status[real_date_reception_col] = pd.to_datetime(df_status[real_date_reception_col], errors='coerce')
                            df_status[real_date_reception_col] = df_status[real_date_reception_col].apply(
                                lambda x: x.strftime('%d/%m/%Y') if pd.notna(x) else 'N/A'
                            )
                        if real_accept_col and real_accept_col in df_status.columns:
                            df_status[real_accept_col] = pd.to_datetime(df_status[real_accept_col], errors='coerce')
                            df_status[real_accept_col] = df_status[real_accept_col].apply(
                                lambda x: x.strftime('%d/%m/%Y') if pd.notna(x) else 'N/A'
                            )
                        
                        # option: masquer par défaut les lignes où un candidat est déjà renseigné
                        show_with_candidate = st.checkbox("Afficher aussi les lignes avec candidat renseigné", value=False)
                        if not show_with_candidate:
                            df_status = df_status[~(df_status['candidate_value'].str.strip() != '')].copy()

                        # Colonnes à afficher (sans demandeur, Direction concernée, Raison du recrutement)
                        desired_cols = [
                            'Poste demandé', 'Entité demandeuse', 'Affectation',
                            'candidate_value',  # Colonne du candidat
                            real_date_reception_col,  # Date de réception de la demande
                            real_accept_col  # Date d'acceptation du candidat
                        ]
                        # Filtrer les colonnes qui sont réellement disponibles et non nulles
                        available_show = [c for c in desired_cols if c and c in df_status.columns]
                        
                        # Renommer candidate_value pour l'affichage
                        if 'candidate_value' in available_show and real_candidat_col:
                            df_status = df_status.rename(columns={'candidate_value': real_candidat_col})  # type: ignore
                            available_show = [real_candidat_col if c == 'candidate_value' else c for c in available_show]

                        df_out_status = df_status[available_show].copy() if available_show else df_status.copy()
                        st.info(f"Lignes avec statut 'En cours' détectées (après filtre candidat): {len(df_out_status)}")
                        st.dataframe(df_out_status.reset_index(drop=True), width="stretch")
                    else:
                        st.warning("Colonne de statut introuvable — impossible de lister les lignes 'En cours'.")
                else:
                    st.info(f"Lignes contribuant au KPI 'Postes en cours' : {len(df_out)} lignes")
                    # Colonnes à afficher dans le dataframe final
                    display_cols_final = [
                        '_orig_index', 'contrib_avant', 'contrib_nouveaux', 'contrib_pourvus', 'contrib_en_cours'
                    ]
                    # Ajouter les colonnes de base si elles existent
                    if real_entite_col in df_out.columns: display_cols_final.insert(1, real_entite_col)
                    if real_statut_col in df_out.columns: display_cols_final.insert(2, real_statut_col)
                    # Utiliser 'candidate_value' au lieu de real_candidat_col directement
                    if 'candidate_value' in df_out.columns: 
                        display_cols_final.insert(3, 'candidate_value')
                    
                    df_out_display = df_out.reset_index().rename(columns={'index': '_orig_index'})
                    
                    # Renommer candidate_value pour l'affichage si besoin
                    if 'candidate_value' in df_out_display.columns and real_candidat_col:
                        # Si la colonne réelle existe déjà, la supprimer d'abord pour éviter les doublons
                        if real_candidat_col in df_out_display.columns:
                            df_out_display = df_out_display.drop(columns=[real_candidat_col])
                        df_out_display = df_out_display.rename(columns={'candidate_value': real_candidat_col})
                        # Mettre à jour display_cols_final avec le nouveau nom
                        display_cols_final = [real_candidat_col if col == 'candidate_value' else col for col in display_cols_final]
                    
                    # S'assurer que toutes les colonnes de contribution sont booléennes
                    for col in ['contrib_avant', 'contrib_nouveaux', 'contrib_pourvus', 'contrib_en_cours']:
                        if col in df_out_display.columns:
                            df_out_display[col] = df_out_display[col].astype(bool)

                    # Ajouter les deux colonnes de date si elles existent
                    extra_date_cols = ["Date d'acceptation du candidat", "Date de réception de la demande"]
                    for col in extra_date_cols:
                        if col in df_out_display.columns and col not in display_cols_final:
                            display_cols_final.append(col)
                    # Filtrer pour n'afficher que les colonnes désirées
                    existing_display_cols = [c for c in display_cols_final if c in df_out_display.columns]
                    st.dataframe(df_out_display[existing_display_cols], width="stretch")
            else:
                st.info('Aucune donnée pour le debug.')
        except Exception as e:
            st.error(f"Erreur lors de la génération du debug: {e}")

    st.markdown("---")

    # 3. Section "Pipeline de Recrutement (Kanban)"
    # (Le nombre total sera ajouté après avoir collecté les données)

    # Construire les données du Kanban à partir du fichier importé (préférence aux données réelles)
    # Détection heuristique des colonnes utiles
    def _find_col(cols, keywords):
        for k in keywords:
            for c in cols:
                if k in c.lower():
                    return c
        return None

    cols = df_recrutement.columns.tolist() if df_recrutement is not None else []
    kanban_col = 'Colonne TG Hire' if 'Colonne TG Hire' in cols else _find_col(cols, ['colonne tg hire'])
    poste_col = _find_col(cols, ['poste', 'title', 'post'])
    entite_col = _find_col(cols, ['entité', 'entite', 'entité demandeuse', 'entite demandeuse', 'entité'])
    lieu_col = _find_col(cols, ['lieu', 'affectation', 'site'])
    demandeur_col = _find_col(cols, ['demandeur', 'requester'])
    recruteur_col = _find_col(cols, ['responsable de traitement', 'recruteur', 'recruiter'])
    commentaire_col = _find_col(cols, ['commentaire', 'comment'])
    accept_date_col = _find_col(cols, ["date d'acceptation du candidat", "date d'acceptation", "date d'acceptation de la promesse", "date d'acceptation promesse", "date d'accept"])  # robust search
    desistement_date_col = _find_col(cols, ["date de désistement", "date désistement", "date desistement", "date desistement"])

    def _normalize_kanban(text):
        if text is None or (isinstance(text, float) and np.isnan(text)):
            return ''
        s = str(text)
        s = unicodedata.normalize('NFKD', s)
        s = ''.join(ch for ch in s if not unicodedata.combining(ch))
        return s.lower().strip()

    # Statuts Kanban canoniques (ordre d'affichage)
    statuts_kanban_display = ["Nouvelle demande", "Sourcing", "Shortlisté", "Signature DRH", "Clôture", "Désistement"]

    # Mapping de formes possibles -> statut canonique
    status_map = {
        'nouvelle demande': 'Nouvelle demande',
        'nouvelle': 'Nouvelle demande',
        'desistement': 'Désistement',
        'desisté': 'Désistement',
        'sourcing': 'Sourcing',
        'shortlist': 'Shortlisté',
        'shortlisté': 'Shortlisté',
        'shortliste': 'Shortlisté',
        'signature drh': 'Signature DRH',
        'signature': 'Signature DRH',
        'cloture': 'Clôture',
        'clôture': 'Clôture',
        'dépriorisé': 'Dépriorisé',
        'depriorise': 'Dépriorisé',
    }

    postes_data = []
    if df_recrutement is not None and kanban_col:
        # Utiliser uniquement la colonne "Colonne TG Hire" pour le statut Kanban
        df_kanban = df_recrutement[df_recrutement[kanban_col].notna()].copy()
        df_kanban[kanban_col] = df_kanban[kanban_col].astype(str)
        for _, r in df_kanban.iterrows():
            raw_kanban = r.get(kanban_col)
            if pd.isna(raw_kanban):
                continue
            norm = _normalize_kanban(raw_kanban)
            canon = None
            for key, val in status_map.items():
                if key in norm:
                    canon = val
                    break
            if canon is None:
                for tgt in statuts_kanban_display:
                    if _normalize_kanban(tgt) == norm:
                        canon = tgt
                        break
            if canon is None:
                continue

            titre = r.get(poste_col, '') if poste_col else r.get('Poste demandé', '')
            # Ajout de la date d'acceptation du candidat pour le filtrage "Clôture"
            accept_date = r.get(accept_date_col) if accept_date_col else None
            # Ajout de la date de désistement pour le filtrage "Désistement"
            desistement_date = r.get(desistement_date_col) if desistement_date_col else None
            postes_data.append({
                'statut': canon,
                'titre': titre or '',
                'entite': r.get(entite_col, '') if entite_col else r.get('Entité demandeuse', ''),
                'lieu': r.get(lieu_col, '') if lieu_col else '',
                'demandeur': r.get(demandeur_col, '') if demandeur_col else '',
                'recruteur': str(r.get(recruteur_col, '')).replace('nan', '') if recruteur_col else '',
                'commentaire': r.get(commentaire_col, '') if commentaire_col else '',
                'date_acceptation': accept_date,
                'date_desistement': desistement_date
            })
    
    # Afficher le titre avec le nombre total de cartes
    # total_cartes = len(postes_data)
    st.subheader("Pipeline de Recrutement (Kanban)")
    
    # CSS pour styliser les cartes (flexbox pour affichage fluide)
    st.markdown("""
    <style>
    .kanban-container {
        display: flex;
        flex-wrap: wrap;
        gap: 10px;
        margin-bottom: 10px;
    }
    .kanban-card {
        flex: 0 0 auto;
        width: 250px;
        border-radius: 5px;
        background-color: #f0f2f6;
        padding: 8px;
        border-left: 4px solid #1f77b4;
        box-shadow: 0 1px 3px rgba(0,0,0,0.1);
        min-height: 80px;
    }
    .kanban-card h4 {
        margin-top: 0;
        margin-bottom: 4px;
        font-size: 1.1em !important;
        color: #2c3e50;
        line-height: 1.2;
        white-space: normal;
        word-wrap: break-word;
        overflow-wrap: break-word;
        hyphens: auto;
        word-break: break-all;
        display: -webkit-box;
        -webkit-line-clamp: 3;
        -webkit-box-orient: vertical;
        overflow: hidden;
    }
    .kanban-card p {
        margin-bottom: 2px;
        font-size: 1.0em !important;
        color: #555;
        line-height: 1.1;
        white-space: normal;
    }
    .kanban-header {
        text-align: center !important;
        font-weight: bold;
        font-size: 1.3em !important;
        color: #FFFFFF !important;
        padding: 8px;
        background-color: #9C182F !important;
        border-radius: 6px;
        margin-bottom: 10px;
        margin-top: 20px;
        border: 1px solid #B01030;
    }
    </style>
    """, unsafe_allow_html=True)
    
    # Remplir chaque section (ligne) avec les postes correspondants
    for i, statut in enumerate(statuts_kanban_display):
        # ...filtrage des postes_in_col (inchangé)...
        if statut == "Clôture":
            reporting_date = st.session_state.get('reporting_date', None)
            if reporting_date is None:
                today = datetime.now()
            else:
                today = reporting_date if isinstance(reporting_date, datetime) else datetime.combine(reporting_date, datetime.min.time())
            current_week_monday = datetime(year=today.year, month=today.month, day=today.day) - timedelta(days=today.weekday())
            start_filter = current_week_monday - timedelta(days=7)
            end_filter = current_week_monday + timedelta(days=6)
            def in_reporting_week(poste):
                accept_date = poste.get('date_acceptation')
                if not accept_date:
                    return False
                if isinstance(accept_date, str):
                    try:
                        accept_date = pd.to_datetime(accept_date, errors='coerce')
                    except Exception:
                        return False
                if not pd.notna(accept_date):
                    return False
                return start_filter <= accept_date <= end_filter + timedelta(days=1)
            postes_in_col = [p for p in postes_data if p["statut"] == statut and in_reporting_week(p)]
        elif statut == "Désistement":
            reporting_date = st.session_state.get('reporting_date', None)
            if reporting_date is None:
                today = datetime.now()
            else:
                today = reporting_date if isinstance(reporting_date, datetime) else datetime.combine(reporting_date, datetime.min.time())
            current_week_monday = datetime(year=today.year, month=today.month, day=today.day) - timedelta(days=today.weekday())
            start_filter = current_week_monday - timedelta(days=7)
            end_filter = current_week_monday + timedelta(days=6)
            def in_reporting_week_desistement(poste):
                desist_date = poste.get('date_desistement')
                if not desist_date:
                    return False
                if isinstance(desist_date, str):
                    try:
                        desist_date = pd.to_datetime(desist_date, errors='coerce')
                    except Exception:
                        return False
                if not pd.notna(desist_date):
                    return False
                return start_filter <= desist_date <= end_filter + timedelta(days=1)
            postes_in_col = [p for p in postes_data if p["statut"] == statut and in_reporting_week_desistement(p)]
        else:
            postes_in_col = [p for p in postes_data if p["statut"] == statut]
        nb_postes = len(postes_in_col)
        st.markdown(f'<div class="kanban-header">{statut} ({nb_postes})</div>', unsafe_allow_html=True)

        # Limiter à 8 cartes par ligne
        max_cards_per_row = 8
        for row_start in range(0, len(postes_in_col), max_cards_per_row):
            cards_html = '<div class="kanban-container">'
            for poste in postes_in_col[row_start:row_start+max_cards_per_row]:
                commentaire = poste.get('commentaire', '')
                commentaire_html = f"<p style='margin-top: 4px; font-style: italic; color: #666;'>💬 {commentaire}</p>" if commentaire and str(commentaire).strip() else ""
                titre_fmt = smart_wrap_title(poste.get('titre', ''))
                card_div = f"""<div class="kanban-card">
<h4><b>{titre_fmt}</b></h4>
<p>📍 {poste.get('entite', 'N/A')} - {poste.get('lieu', 'N/A')}</p>
<p>👤 {poste.get('demandeur', 'N/A')}</p>
<p>✍️ {poste.get('recruteur', 'N/A')}</p>
{commentaire_html}
</div>"""
                cards_html += card_div
            cards_html += '</div>'
            st.markdown(cards_html, unsafe_allow_html=True)

    # Debug local pour Pipeline de Recrutement (Kanban)
    st.markdown("---")
    with st.expander("🔍 Debug - Détails des lignes (Pipeline Kanban)", expanded=False):
        try:
            st.markdown("**Lignes contribuant au pipeline par statut:**")
            if postes_data:
                # Convertir postes_data en DataFrame pour affichage
                df_kanban_debug = pd.DataFrame(postes_data)
                # Ajouter colonnes de contribution par statut
                for s in statuts_kanban_display:
                    df_kanban_debug[f'contrib_{s}'] = df_kanban_debug['statut'] == s
                # Afficher colonnes pertinentes
                cols_show = ['titre', 'entite', 'lieu', 'recruteur', 'statut'] + [f'contrib_{s}' for s in statuts_kanban_display]
                cols_available = [c for c in cols_show if c in df_kanban_debug.columns]
                st.dataframe(df_kanban_debug[cols_available].reset_index(drop=True), use_container_width=True, hide_index=True)
            else:
                st.write("Aucune donnée de pipeline disponible.")
        except Exception as e:
            st.write(f"Erreur debug pipeline: {e}")


def generate_table_image_simple(weekly_metrics):
    """Génère une image simple du tableau avec PIL incluant les LOGOS"""
    from PIL import Image, ImageDraw, ImageFont
    import tempfile
    
    try:
        metrics_by_entity = weekly_metrics.get('metrics_by_entity', {})
        excluded_entities = {'BESIX-TGCC', 'DECO EXCELL', 'TG PREFA'}
        metrics_included = {k: v for k, v in metrics_by_entity.items() if k not in excluded_entities}
        
        # Calculer les totaux
        total_avant = sum(data['avant'] for data in metrics_included.values())
        total_nouveaux = sum(data['nouveaux'] for data in metrics_included.values())
        total_pourvus = sum(data['pourvus'] for data in metrics_included.values())
        total_en_cours = sum(data['en_cours'] for data in metrics_included.values())
        
        # --- CHARGEMENT DES LOGOS ---
        current_dir = os.path.dirname(os.path.abspath(__file__))
        logo_folder = os.path.join(os.path.dirname(current_dir), "LOGO")
        loaded_logos = {}
        
        # Mapping des noms d'entités vers les fichiers (ordre prioritaire)
        entity_logo_map = {
            'TGCC IMMOBILIER': 'tgcc-immobilier.png',
            'TGCC-IMMOBILIER': 'tgcc-immobilier.png',
            'TGCC Immobilier': 'tgcc-immobilier.png',
            'TG STEEL': 'TG STEEL.PNG',
            'TG STONE': 'TG STONE.PNG',
            'TG ALU': 'TG ALU.PNG',
            'TG COVER': 'TG COVER.PNG',
            'TG WOOD': 'TG WOOD.PNG',
            'STAM': 'STAM.png',
            'BFO': 'BFO.png',
            'TGEM': 'TGEM.PNG',
            'TGCC': 'TGCC.PNG'
        }

        if os.path.exists(logo_folder):
            for filename in os.listdir(logo_folder):
                if filename.lower().endswith(('.png', '.jpg', '.jpeg', '.jfif')):
                    try:
                        img_path = os.path.join(logo_folder, filename)
                        loaded_logos[filename.upper()] = Image.open(img_path).convert("RGBA")
                    except Exception as e:
                        print(f"Erreur chargement logo {filename}: {e}")

        def get_logo_image(entity_name):
            """Récupère l'image du logo pour une entité (priorité à IMMOBILIER)"""
            name_upper = str(entity_name).upper().strip()
            
            # Chercher dans le mapping avec ordre de priorité (déjà ordonné dans entity_logo_map)
            filename = None
            for key in entity_logo_map:
                if key.upper() in name_upper:
                    filename = entity_logo_map[key]
                    break
            
            # Si pas trouvé, chercher directement
            if not filename:
                for fname in loaded_logos.keys():
                    base_name = os.path.splitext(fname)[0]
                    if base_name in name_upper or name_upper in base_name:
                        filename = fname
                        break
            
            if filename and filename.upper() in loaded_logos:
                return loaded_logos[filename.upper()]
            return None

        # --- DESSIN DU TABLEAU ---
        num_rows = len(metrics_included) + 2  # +1 pour header, +1 pour total
        row_height = 80  # Hauteur augmentée pour les logos
        width = 1920
        height = num_rows * row_height + 60
        
        img = Image.new('RGB', (width, height), 'white')
        draw = ImageDraw.Draw(img)
        
        # Police
        try:
            font_header = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf", 24)
            font_data = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", 20)
        except:
            font_header = ImageFont.load_default()
            font_data = ImageFont.load_default()
        
        # Header
        headers = ['Entité', 'Postes avant', 'Nouveaux postes', 'Postes pourvus', 'Postes en cours']
        col_widths = [350, 300, 300, 300, 300]  # Première colonne plus large pour logos
        x_positions = [0] + [sum(col_widths[:i+1]) for i in range(len(col_widths))]
        
        # Dessiner le header
        draw.rectangle([0, 0, width, row_height], fill='#9C182F')
        for i, header in enumerate(headers):
            x = x_positions[i] + col_widths[i] // 2
            draw.text((x, row_height // 2), header, fill='white', font=font_header, anchor='mm')
        
        # Données
        y_offset = row_height
        
        # Custom sort and filter
        sorted_items = []
        for entite, data in metrics_included.items():
            # Filter out empty TG WOOD (duplicate) or any empty entity
            # User specifically mentioned duplicate TG WOOD being empty
            entite_str = str(entite).upper().strip()
            if 'TG WOOD' in entite_str and data['avant'] == 0 and data['nouveaux'] == 0 and data['pourvus'] == 0 and data['en_cours'] == 0:
                continue
            # Also filter out literal "NAN" or "NONE" strings if they appear as entities
            if entite_str in ['NAN', 'NONE', '']:
                continue
                
            sorted_items.append((entite, data))
            
        # Sort: TGCC first, then others
        def sort_key(item):
            name = str(item[0]).upper().strip()
            if name == 'TGCC':
                return '000_TGCC' # Force first
            if 'TGCC' in name and 'IMMOBILIER' not in name:
                 return '001_TGCC_OTHER'
            return name
            
        sorted_items.sort(key=sort_key)

        for entite, data in sorted_items:
            # Lignes alternées
            if ((y_offset - row_height) // row_height) % 2 == 0:
                draw.rectangle([0, y_offset, width, y_offset + row_height], fill='#f9f9f9')
            
            # Logo ou texte de l'entité
            logo_img = get_logo_image(entite)
            cell_center_x = x_positions[0] + col_widths[0] // 2
            cell_center_y = y_offset + row_height // 2
            
            if logo_img:
                # Redimensionner le logo avec tailles personnalisées
                entite_upper = str(entite).upper().strip()
                if 'STEEL' in entite_upper or 'STONE' in entite_upper:
                    new_h = 45  # Plus petit pour STEEL et STONE
                elif 'BFO' in entite_upper:
                    new_h = 70  # Plus grand pour BFO
                elif 'ALU' in entite_upper or 'WOOD' in entite_upper:
                    new_h = 65  # Plus grand pour ALU et WOOD
                elif 'COVER' in entite_upper:
                    new_h = 60  # Plus grand pour COVER
                else:
                    new_h = 55  # Taille standard
                
                aspect_ratio = logo_img.width / logo_img.height
                new_w = int(new_h * aspect_ratio)
                if new_w > col_widths[0] - 20:
                    new_w = col_widths[0] - 20
                    new_h = int(new_w / aspect_ratio)
                
                logo_resized = logo_img.resize((new_w, new_h), Image.Resampling.LANCZOS)
                paste_x = int(cell_center_x - new_w / 2)
                paste_y = int(cell_center_y - new_h / 2)
                
                # Coller avec transparence
                if logo_resized.mode == 'RGBA':
                    img.paste(logo_resized, (paste_x, paste_y), logo_resized)
                else:
                    img.paste(logo_resized, (paste_x, paste_y))
            else:
                # Fallback texte si pas de logo
                draw.text((cell_center_x, cell_center_y), entite[:20], fill='black', font=font_data, anchor='mm')
            draw.text((x_positions[1] + col_widths[1] // 2, y_offset + row_height // 2), 
                     str(data['avant'] if data['avant'] > 0 else '-'), fill='black', font=font_data, anchor='mm')
            draw.text((x_positions[2] + col_widths[2] // 2, y_offset + row_height // 2), 
                     str(data['nouveaux'] if data['nouveaux'] > 0 else '-'), fill='black', font=font_data, anchor='mm')
            draw.text((x_positions[3] + col_widths[3] // 2, y_offset + row_height // 2), 
                     str(data['pourvus'] if data['pourvus'] > 0 else '-'), fill='black', font=font_data, anchor='mm')
            draw.text((x_positions[4] + col_widths[4] // 2, y_offset + row_height // 2), 
                     str(data['en_cours'] if data['en_cours'] > 0 else '-'), fill='black', font=font_data, anchor='mm')
            
            # Bordures
            draw.line([(0, y_offset + row_height), (width, y_offset + row_height)], fill='#ddd', width=1)
            y_offset += row_height
        
        # Ligne TOTAL
        draw.rectangle([0, y_offset, width, y_offset + row_height], fill='#9C182F')
        draw.text((x_positions[0] + col_widths[0] // 2, y_offset + row_height // 2), 
                 'TOTAL', fill='white', font=font_header, anchor='mm')
        draw.text((x_positions[1] + col_widths[1] // 2, y_offset + row_height // 2), 
                 str(total_avant), fill='white', font=font_header, anchor='mm')
        draw.text((x_positions[2] + col_widths[2] // 2, y_offset + row_height // 2), 
                 str(total_nouveaux), fill='white', font=font_header, anchor='mm')
        draw.text((x_positions[3] + col_widths[3] // 2, y_offset + row_height // 2), 
                 str(total_pourvus), fill='white', font=font_header, anchor='mm')
        draw.text((x_positions[4] + col_widths[4] // 2, y_offset + row_height // 2), 
                 str(total_en_cours), fill='white', font=font_header, anchor='mm')
        
        # Sauvegarder
        output_path = os.path.join(tempfile.gettempdir(), 'table_reporting.png')
        img.save(output_path)
        return output_path
        
    except Exception as e:
        st.error(f"Erreur génération image tableau: {e}")
        return None


def generate_kanban_image_simple(df_recrutement):
    """Génère une image simple du Kanban avec PIL"""
    from PIL import Image, ImageDraw, ImageFont
    import tempfile
    import pandas as pd
    
    try:
        # Si pas de dataframe, utiliser des données par défaut
        if df_recrutement is None or len(df_recrutement) == 0:
            colonnes = {
                'Sourcing': pd.DataFrame(),
                'Shortlisté': pd.DataFrame(),
                'Signature DRH': pd.DataFrame(),
                'Clôture': pd.DataFrame(),
                'Désistement': pd.DataFrame(),
            }
        else:
            # Charger les données du kanban
            colonnes = {
                'Sourcing': df_recrutement[df_recrutement['Colonne TG Hire'] == 'Sourcing'],
                'Shortlisté': df_recrutement[df_recrutement['Colonne TG Hire'] == 'Shortlisté'],
                'Signature DRH': df_recrutement[df_recrutement['Colonne TG Hire'] == 'Signature DRH'],
                'Clôture': df_recrutement[df_recrutement['Colonne TG Hire'] == 'Clôture'],
                'Désistement': df_recrutement[df_recrutement['Colonne TG Hire'] == 'Désistement'],
            }
        
        # Dimensions (haute résolution)
        col_width = 370  # Plus large
        num_cols = 5
        width = col_width * num_cols + 60
        height = 1080  # Haute résolution
        
        img = Image.new('RGB', (width, height), 'white')
        draw = ImageDraw.Draw(img)
        
        # Polices (plus grandes)
        try:
            font_header = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf", 20)
            font_card_title = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf", 14)
            font_card = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", 12)
        except:
            font_header = ImageFont.load_default()
            font_card_title = ImageFont.load_default()
            font_card = ImageFont.load_default()
        
        x_offset = 20
        for statut, df_col in colonnes.items():
            # En-tête de colonne
            draw.rectangle([x_offset, 20, x_offset + col_width - 10, 60], fill='#9C182F', outline='#9C182F')
            count = len(df_col)
            draw.text((x_offset + col_width // 2 - 5, 40), f"{statut} ({count})", 
                     fill='white', font=font_header, anchor='mm')
            
            # Cartes
            y_offset = 80
            for idx, row in df_col.iterrows():
                if y_offset > height - 100:
                    break
                
                # Carte
                draw.rectangle([x_offset + 5, y_offset, x_offset + col_width - 15, y_offset + 90], 
                              fill='#f9f9f9', outline='#ddd', width=2)
                
                # Titre du poste (wrap jusqu'à 3 lignes)
                titre_raw = str(row.get('Intitulé du poste', 'N/A')).strip()
                titre_formatted = smart_wrap_title(titre_raw, max_line_length=25)
                title_lines = [seg for seg in titre_formatted.replace('<br>', '\n').split('\n') if seg][:3]
                for i, line in enumerate(title_lines):
                    draw.text((x_offset + 10, y_offset + 10 + (i * 16)), line, fill='#9C182F', font=font_card_title)
                title_block_height = 10 + (len(title_lines) * 16)
                
                # Entité et lieu
                entite = str(row.get('Entité demandeuse', '')).strip()
                lieu = str(row.get('Ville', '')).strip()
                draw.text((x_offset + 10, y_offset + title_block_height + 18), f"{entite[:40]} - {lieu[:40]}", fill='black', font=font_card)
                
                # Demandeur
                demandeur = str(row.get('Nom du Demandeur', '')).strip()[:50]
                draw.text((x_offset + 10, y_offset + title_block_height + 36), f"👤 {demandeur}", fill='black', font=font_card)
                
                # Recruteur
                recruteur = str(row.get('RH en charge du recrutement', '')).strip()[:50]
                draw.text((x_offset + 10, y_offset + title_block_height + 54), f"✍️ {recruteur}", fill='black', font=font_card)

                # Commentaire (si présent)
                comment_col = next((c for c in row.index if 'comment' in c.lower()), None)
                if comment_col:
                    commentaire_val = str(row.get(comment_col, '')).strip()
                    if commentaire_val:
                        draw.text((x_offset + 10, y_offset + title_block_height + 72), f"[!] {commentaire_val[:80]}", fill='#666666', font=font_card)
                
                y_offset += max(110, title_block_height + 88)
            
            x_offset += col_width
        
        # Sauvegarder
        output_path = os.path.join(tempfile.gettempdir(), 'kanban_reporting.png')
        img.save(output_path)
        return output_path
        
    except Exception as e:
        st.error(f"Erreur génération image kanban: {e}")
        return None


def check_logos():
    """Vérifie que les logos sont bien accessibles"""
    current_dir = os.path.dirname(os.path.abspath(__file__))
    logo_folder = os.path.join(os.path.dirname(current_dir), "LOGO")
    
    if not os.path.exists(logo_folder):
        st.error(f"❌ Dossier LOGO non trouvé: {logo_folder}")
        return False
    
    logos = os.listdir(logo_folder)
    if not logos:
        st.error("❌ Aucun logo trouvé dans le dossier LOGO")
        return False
    
    st.success(f"✅ {len(logos)} logos trouvés dans: {logo_folder}")
    st.info(f"📁 Fichiers: {', '.join(logos[:5])}{'...' if len(logos) > 5 else ''}")
    return True


def find_chromium_executable():
    """Trouve le chemin de l'exécutable Chromium en testant plusieurs emplacements."""
    import shutil
    
    # Liste des chemins à tester, dans l'ordre de priorité
    # IMPORTANT: html2image a besoin du vrai binaire, pas du script wrapper
    possible_paths = [
        '/usr/lib/chromium/chromium',  # Debian/Ubuntu - VRAI BINAIRE (priorité 1)
        '/usr/bin/google-chrome',       # Google Chrome
        '/usr/bin/google-chrome-stable',
        '/usr/bin/chromium',            # Script wrapper (peut ne pas fonctionner avec html2image)
        '/usr/bin/chromium-browser',
        'chromium',                     # PATH
        'chromium-browser',
        'google-chrome'
    ]
    
    for path in possible_paths:
        # Vérifier si le fichier existe
        if os.path.isfile(path):
            st.info(f"🔍 Test chemin: {path} - ✅ TROUVÉ")
            return path
        elif shutil.which(path):
            found_path = shutil.which(path)
            st.info(f"🔍 Test chemin: {path} - ✅ TROUVÉ via PATH: {found_path}")
            return found_path
        else:
            st.warning(f"🔍 Test chemin: {path} - ❌ Non trouvé")
    
    st.error("❌ Aucun exécutable Chromium trouvé dans tous les chemins testés")
    return None


def generate_table_html_image(weekly_metrics):
    """Génère une image du tableau des besoins avec logos"""
    import tempfile
    
    st.info("🔍 Début génération tableau HTML avec logos...")
    
    try:
        # Essayer d'abord avec html2image + Chromium
        from html2image import Html2Image
        st.info("✅ Module Html2Image importé")
        
        # Trouver Chromium automatiquement
        chromium_path = find_chromium_executable()
        if not chromium_path:
            st.error("❌ Chromium non trouvé - les images seront générées avec PIL (sans logos)")
            raise Exception("Chromium non trouvé")
        
        st.success(f"✅ Chromium détecté: {chromium_path}")
        
        # Configurer html2image pour utiliser Chromium avec flags no-sandbox
        hti = Html2Image(
            output_path=tempfile.gettempdir(),
            browser_executable=chromium_path,
            custom_flags=['--no-sandbox', '--disable-dev-shm-usage', '--disable-gpu', '--headless']
        )
        st.success("✅ Html2Image configuré")
    except Exception as e:
        st.warning(f"⚠️ html2image non disponible ({e}), utilisation de PIL à la place")
        # Fallback vers PIL
        return generate_table_image_simple(weekly_metrics)
    
    try:
        # Créer le HTML du tableau exactement comme dans Streamlit
        metrics_by_entity = weekly_metrics.get('metrics_by_entity', {})
        excluded_entities = {'BESIX-TGCC', 'DECO EXCELL', 'TG PREFA'}
        metrics_included = {k: v for k, v in metrics_by_entity.items() if k not in excluded_entities}
        
        # Charger les logos
        logo_folder = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))), "LOGO")
        entity_logo_map = {
            'TGCC': 'TGCC.PNG',
            'TGEM': 'TGEM.PNG',
            'TG ALU': 'TG ALU.PNG',
            'TG COVER': 'TG COVER.PNG',
            'TG STEEL': 'TG STEEL.PNG',
            'TG STONE': 'TG STONE.PNG',
            'TG WOOD': 'TG WOOD.PNG',
            'STAM': 'STAM.png',
            'BFO': 'BFO.png',
            'TGCC IMMOBILIER': 'tgcc-immobilier.png',
            'TGCC Immobilier': 'tgcc-immobilier.png'
        }
        
        logos_dict = {}
        for entity, logo_file in entity_logo_map.items():
            logo_path = os.path.join(logo_folder, logo_file)
            if os.path.exists(logo_path):
                with open(logo_path, "rb") as f:
                    logos_dict[entity] = base64.b64encode(f.read()).decode()
        
        # Fonction pour obtenir l'affichage avec logo
        def get_entity_display(name_str):
            if not name_str or pd.isna(name_str):
                return name_str
            
            name_upper = str(name_str).upper().strip()
            
            # Ordre de priorité: vérifier IMMOBILIER en premier pour éviter confusion avec TGCC simple
            priority_order = [
                'TGCC IMMOBILIER', 'TGCC-IMMOBILIER', 'TGCC Immobilier',
                'TG STEEL', 'TG STONE', 'TG ALU', 'TG COVER', 'TG WOOD',
                'STAM', 'BFO', 'TGEM', 'TGCC'
            ]
            
            for entity_key in priority_order:
                if entity_key.upper() in name_upper:
                    if entity_key in logos_dict:
                        logo_b64 = logos_dict[entity_key]
                        # Log pour débogage
                        print(f"Match: '{name_str}' -> '{entity_key}'")
                        # Ajuster les tailles selon les besoins
                        size_map = {
                            'TG STEEL': 45,      # Réduit
                            'TG STONE': 45,      # Réduit
                            'TGCC IMMOBILIER': 50,
                            'TGCC Immobilier': 50,
                            'BFO': 90,           # Augmenté
                            'TG ALU': 75,        # Augmenté
                            'TG COVER': 70,      # Augmenté
                            'TG WOOD': 75,       # Augmenté
                            'STAM': 63,
                            'TGEM': 63,
                            'TGCC': 63
                        }
                        logo_size = size_map.get(entity_key, 63)
                        img_tag = f'<img src="data:image/png;base64,{logo_b64}" style="height: {logo_size}px; vertical-align: middle;" />'
                        return img_tag
            print(f"No match for: '{name_str}'")
            return name_str
        
        # Calculer les totaux
        total_avant = sum(data['avant'] for data in metrics_included.values())
        total_nouveaux = sum(data['nouveaux'] for data in metrics_included.values())
        total_pourvus = sum(data['pourvus'] for data in metrics_included.values())
        total_en_cours = sum(data['en_cours'] for data in metrics_included.values())
        
        # Créer le HTML
        html_table = f"""
        <html>
        <head>
            <style>
                body {{
                    margin: 0;
                    padding: 20px;
                    font-family: Arial, sans-serif;
                    background: white;
                }}
                .table-container {{
                    width: 100%;
                }}
                .custom-table {{
                    width: 90%;
                    margin: 0 auto;
                    border-collapse: collapse;
                    background-color: white;
                    box-shadow: 0 2px 4px rgba(0,0,0,0.1);
                }}
                .custom-table th {{
                    background-color: #9C182F;
                    color: white;
                    font-weight: bold;
                    text-align: center;
                    padding: 8px 6px;
                    border: 1px solid white;
                    font-size: 1.0em;
                    line-height: 1.2;
                    white-space: normal;
                    word-wrap: break-word;
                    overflow-wrap: break-word;
                    word-break: break-word;
                    min-width: 120px;
                }}
                .custom-table td {{
                    text-align: center;
                    padding: 6px 4px;
                    border: 1px solid #ddd;
                    background-color: white;
                    font-size: 1.1em;
                    line-height: 1.2;
                    font-weight: 500;
                }}
                .custom-table .entity-cell {{
                    text-align: center;
                    padding: 4px 2px;
                    font-weight: 600;
                    min-width: 80px;
                }}
                .custom-table .total-row {{
                    background-color: #9C182F;
                    color: white;
                    font-weight: bold;
                    border-top: 2px solid #9C182F;
                }}
                .custom-table .total-row td {{
                    background-color: #9C182F;
                    color: white;
                    font-size: 1.1em;
                    font-weight: bold;
                    border: 1px solid #9C182F;
                }}
                .custom-table .total-row .entity-cell {{
                    text-align: center;
                    padding: 4px 2px;
                    font-weight: bold;
                    background-color: #9C182F;
                    color: white;
                }}
            </style>
        </head>
        <body>
            <div class="table-container">
                <table class="custom-table">
                    <thead>
                        <tr>
                            <th>Entité</th>
                            <th>Nb postes ouverts avant début semaine</th>
                            <th>Nb nouveaux postes ouverts cette semaine</th>
                            <th>Nb postes pourvus cette semaine</th>
                            <th>Nb postes en cours cette semaine (sourcing)</th>
                        </tr>
                    </thead>
                    <tbody>
        """
        
        # Ajouter les lignes de données
        for entite, data in metrics_included.items():
            html_table += f'''<tr>
                <td class="entity-cell">{get_entity_display(entite)}</td>
                <td>{data['avant'] if data['avant'] > 0 else '-'}</td>
                <td>{data['nouveaux'] if data['nouveaux'] > 0 else '-'}</td>
                <td>{data['pourvus'] if data['pourvus'] > 0 else '-'}</td>
                <td>{data['en_cours'] if data['en_cours'] > 0 else '-'}</td>
            </tr>'''
        
        # Ligne TOTAL
        html_table += f'''<tr class="total-row">
            <td class="entity-cell">TOTAL</td>
            <td>{total_avant}</td>
            <td>{total_nouveaux}</td>
            <td>{total_pourvus}</td>
            <td>{total_en_cours}</td>
        </tr>
        '''
        
        html_table += """
                    </tbody>
                </table>
            </div>
        </body>
        </html>
        """
        
        # Convertir en image avec html2image (haute résolution)
        # Supprimer l'ancien fichier s'il existe
        old_file = os.path.join(tempfile.gettempdir(), 'table.png')
        if os.path.exists(old_file):
            os.remove(old_file)
            st.info(f"🗑️ Ancien fichier table.png supprimé")
        
        image_path = hti.screenshot(html_str=html_table, save_as='table.png', size=(1920, 1080))[0]
        st.success(f"✅ Image tableau générée: {image_path}")
        
        return image_path
    except Exception as e:
        st.error(f"Erreur lors de la génération de l'image du tableau avec html2image: {e}")
        st.info("Tentative avec PIL...")
        # Fallback vers PIL
        return generate_table_image_simple(weekly_metrics)


def generate_kanban_statut_image(df_recrutement, statut, max_cards=10):
    """Génère une image pour un statut Kanban spécifique.
    
    Args:
        df_recrutement: DataFrame avec les données
        statut: Le statut à afficher ('Sourcing', 'Shortlisté', etc.)
        max_cards: Nombre maximum de cartes à afficher
    
    Returns:
        Chemin de l'image générée
    """
    import tempfile
    
    # Vérifier si le DataFrame est vide
    if df_recrutement is None or len(df_recrutement) == 0:
        st.warning(f"Aucune donnée pour le statut {statut}")
        return None
    
    st.info(f"🔍 Début génération Kanban pour {statut}...")
    
    try:
        # Essayer d'abord avec html2image + Chromium
        from html2image import Html2Image
        st.info(f"✅ Module Html2Image importé pour {statut}")
        
        # Trouver Chromium automatiquement
        chromium_path = find_chromium_executable()
        if not chromium_path:
            st.error(f"❌ Chromium non trouvé pour Kanban {statut} - utilisation de PIL (limité)")
            raise Exception("Chromium non trouvé")
        
        st.success(f"✅ Chromium détecté pour Kanban {statut}: {chromium_path}")
        
        # Configurer html2image
        hti = Html2Image(
            output_path=tempfile.gettempdir(),
            browser_executable=chromium_path,
            custom_flags=['--no-sandbox', '--disable-dev-shm-usage', '--disable-gpu', '--headless']
        )
    except Exception as e:
        st.warning(f"⚠️ html2image non disponible pour {statut} ({e}), utilisation de PIL à la place")
        # Fallback vers PIL
        return generate_kanban_statut_image_simple(df_recrutement, statut, max_cards)
    
    try:
        # Charger les données pour ce statut
        postes_data = []
        
        for index, row in df_recrutement.iterrows():
            if row.get('Colonne TG Hire') == statut:
                # Utiliser les noms de colonnes originaux du fichier Excel
                postes_data.append({
                    "titre": str(row.get('Poste demandé ', 'N/A')).strip(),
                    "entite": str(row.get('Entité demandeuse', 'N/A')).strip(),
                    "lieu": str(row.get('Affectation', 'N/A')).strip(),
                    "demandeur": str(row.get('Nom Prénom du demandeur', 'N/A')).strip(),
                    "recruteur": str(row.get('Responsable de traitement de  la demande ', 'N/A')).strip(),
                    "commentaire": str(row.get('Commentaire', '')) if pd.notna(row.get('Commentaire')) else None
                })
                # Limiter le nombre de cartes
                if len(postes_data) >= max_cards:
                    break
        
        # Créer HTML
        html_content = f"""
        <html>
        <head>
            <meta charset="UTF-8">
            <style>
                body {{
                    margin: 0;
                    padding: 20px;
                    font-family: Arial, sans-serif;
                    background: white;
                }}
                .statut-header {{
                    background-color: #9C182F;
                    color: white;
                    padding: 22px;
                    border-radius: 8px;
                    font-size: 2.3em;
                    font-weight: bold;
                    text-align: center;
                    margin-bottom: 28px;
                }}
                .cards-container {{
                    display: grid;
                    grid-template-columns: repeat(auto-fill, minmax(260px, 1fr));
                    grid-auto-flow: dense;
                    gap: 22px;
                    padding: 20px;
                    justify-items: center;
                    justify-content: center;
                    align-items: start;
                    width: 100%;
                    margin: 0 auto;
                }}
                .kanban-card {{
                    background-color: #f0f2f6;
                    border-radius: 8px;
                    padding: 14px;
                    box-shadow: 0 2px 6px rgba(0,0,0,0.08);
                    border-left: 5px solid #1f77b4;
                    min-height: 140px;
                    width: 100%;
                    box-sizing: border-box;
                }}
                .kanban-card h4 {{
                    margin: 0 0 8px 0;
                    color: #9C182F;
                    font-size: 1.05em;
                    line-height: 1.25;
                    font-weight: 700;
                    white-space: normal !important;
                    word-break: break-word !important;
                    overflow-wrap: break-word !important;
                    hyphens: auto !important;
                    display: -webkit-box;
                    -webkit-line-clamp: 3;
                    -webkit-box-orient: vertical;
                    overflow: hidden;
                }}
                .kanban-card p {{
                    margin: 4px 0;
                    font-size: 0.96em;
                    color: #444;
                    line-height: 1.3;
                    word-wrap: break-word;
                    overflow-wrap: break-word;
                    word-break: break-word;
                }}
            </style>
        </head>
        <body>
            <div class="statut-header">{statut} ({len(postes_data)} postes)</div>
            <div class="cards-container">
        """
        
        for poste in postes_data:
                titre_fmt = smart_wrap_title(poste.get('titre', ''))
                commentaire_html = f"<p style='color:#666;font-style:italic;'>[!] {poste['commentaire']}</p>" if poste.get('commentaire') and str(poste.get('commentaire')).strip() not in ['nan', 'None', ''] else ""
                html_content += f"""
                <div class="kanban-card">
                    <h4><b>{titre_fmt}</b></h4>
                    <p><b>&gt;</b> {poste['entite']} - {poste['lieu']}</p>
                    <p><b>*</b> {poste['demandeur']}</p>
                    <p><b>#</b> {poste['recruteur']}</p>
                    {commentaire_html}
                </div>
                """
        
        html_content += """
            </div>
        </body>
        </html>
        """
        
        # Générer l'image
        image_filename = f'{statut.lower().replace(" ", "_")}.png'
        image_path = hti.screenshot(html_str=html_content, save_as=image_filename, size=(1920, 1080))[0]
        
        return image_path
    except Exception as e:
        st.error(f"Erreur lors de la génération de l'image {statut}: {e}")
        return generate_kanban_statut_image_simple(df_recrutement, statut, max_cards)


def generate_kanban_statut_image_simple(df_recrutement, statut, max_cards=10):
    """Génère une image simple pour un statut avec PIL (fallback)."""
    from PIL import Image, ImageDraw, ImageFont
    import tempfile
    
    width, height = 1920, 1080
    img = Image.new('RGB', (width, height), 'white')
    draw = ImageDraw.Draw(img)
    
    try:
        font_header = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf", 40)
        font_card_title = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf", 18)
        font_card = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", 16)
    except:
        font_header = ImageFont.load_default()
        font_card_title = ImageFont.load_default()
        font_card = ImageFont.load_default()
    
    # Header
    draw.rectangle([0, 0, width, 100], fill='#9C182F')
    
    # Filtrer les données
    if df_recrutement is not None and len(df_recrutement) > 0:
        # Trouver la colonne 'Colonne TG Hire' (avec gestion des espaces)
        col_tg_hire = next((c for c in df_recrutement.columns if 'colonne' in c.lower() and 'tg' in c.lower() and 'hire' in c.lower()), None)
        
        if col_tg_hire:
            df_statut = df_recrutement[df_recrutement[col_tg_hire] == statut]
            count = len(df_statut)
        else:
            # Fallback: essayer de trouver une colonne 'Statut'
            col_statut = next((c for c in df_recrutement.columns if 'statut' in c.lower()), None)
            if col_statut:
                df_statut = df_recrutement[df_recrutement[col_statut] == statut]
                count = len(df_statut)
            else:
                df_statut = None
                count = 0
    else:
        df_statut = None
        count = 0
    
    # Ajouter un timestamp de génération pour prouver que l'image est fraîche
    from datetime import datetime
    timestamp = datetime.now().strftime("%H:%M:%S")
    draw.text((width - 150, 50), f"Généré à {timestamp}", fill='white', font=font_card, anchor='mm')

    draw.text((width // 2, 50), f"{statut} ({count} postes)", fill='white', font=font_header, anchor='mm')
    
    # Cartes en grille 6 colonnes
    if df_statut is not None and count > 0:
        card_width = (width - 140) // 6  # 6 colonnes
        card_height = 160
        x_positions = [20 + i * (card_width + 20) for i in range(6)]
        y_offset = 120
        col = 0
        
        for idx, row in df_statut.iterrows():
            if y_offset > height - 200:
                break
            
            x = x_positions[col]
            
            # Carte
            draw.rectangle([x, y_offset, x + card_width, y_offset + card_height], 
                          fill='#f0f2f6', outline='#ddd', width=2)
            
            # Contenu - utiliser noms colonnes originaux (avec gestion des espaces potentiels)
            # Titre
            titre_col = 'Poste demandé' if 'Poste demandé' in row else 'Poste demandé '
            titre = str(row.get(titre_col, 'N/A')).strip()
            # Wrap title robustly (use <br> segments then draw lines)
            formatted = smart_wrap_title(titre, max_line_length=25)
            wrapped_lines = [seg for seg in formatted.replace('<br>', '\n').split('\n') if seg][:3]
            for i, line in enumerate(wrapped_lines):
                draw.text((x + 10, y_offset + 10 + (i * 20)), line, fill='#9C182F', font=font_card_title)
            title_block_height = 10 + (len(wrapped_lines) * 20)
            
            # Entité
            entite_val = str(row.get('Entité demandeuse', 'N/A')).strip()
            affectation_val = str(row.get('Affectation', row.get('Direction concernée', ''))).strip()
            entite_txt = f"> {entite_val}"
            if affectation_val:
                entite_txt += f" - {affectation_val}"
            draw.text((x + 10, y_offset + title_block_height + 20), entite_txt[:50], fill='#555', font=font_card)
            
            # Demandeur
            demandeur_col = 'Nom Prénom du demandeur'
            if demandeur_col not in row:
                # Chercher une colonne ressemblante
                demandeur_col = next((c for c in row.index if 'demandeur' in c.lower() and 'nom' in c.lower()), 'Demandeur')
            
            demandeur_val = str(row.get(demandeur_col, 'N/A')).strip()
            demandeur = f"* {demandeur_val}"
            draw.text((x + 10, y_offset + title_block_height + 45), demandeur[:50], fill='#555', font=font_card)
            
            # Recruteur
            recruteur_col = 'Responsable de traitement de la demande'
            if recruteur_col not in row:
                # Chercher avec double espace ou autre
                recruteur_col = next((c for c in row.index if 'responsable' in c.lower() and 'traitement' in c.lower()), 'RH')
                
            recruteur_val = str(row.get(recruteur_col, 'N/A')).strip()
            recruteur = f"# {recruteur_val}"
            draw.text((x + 10, y_offset + title_block_height + 70), recruteur[:50], fill='#555', font=font_card)

            # Commentaire (si présent)
            comment_col = next((c for c in row.index if 'comment' in c.lower()), None)
            if comment_col:
                commentaire_val = str(row.get(comment_col, '')).strip()
                if commentaire_val:
                    draw.text((x + 10, y_offset + title_block_height + 88), f"[!] {commentaire_val[:80]}", fill='#666666', font=font_card)
            
            col += 1
            if col >= 6:
                col = 0
                y_offset += card_height + 30
    
    # Sauvegarder
    output_path = os.path.join(tempfile.gettempdir(), f'{statut.lower().replace(" ", "_")}_reporting.png')
    img.save(output_path)
    return output_path


def generate_kanban_html_image(df_recrutement):
    """Génère une image du Kanban avec les données réelles"""
    import tempfile
    
    try:
        # Essayer d'abord avec html2image + Chromium
        from html2image import Html2Image
        
        # Trouver Chromium automatiquement
        chromium_path = find_chromium_executable()
        if not chromium_path:
            raise Exception("Chromium non trouvé")
        
        # Configurer html2image pour utiliser Chromium avec flags no-sandbox
        hti = Html2Image(
            output_path=tempfile.gettempdir(),
            browser_executable=chromium_path,
            custom_flags=['--no-sandbox', '--disable-dev-shm-usage', '--disable-gpu', '--headless']
        )
        st.info(f"✅ Chromium trouvé: {chromium_path}")
    except Exception as e:
        st.warning(f"⚠️ html2image non disponible ({e}), utilisation de PIL à la place")
        return generate_kanban_image_simple(df_recrutement)
    
    try:
        # Charger les données réelles du dataframe
        statuts = ["Nouvelle demande", "Sourcing", "Shortlisté", "Signature DRH", "Clôture", "Désistement"]
        postes_data = []
        
        for index, row in df_recrutement.iterrows():
            if row.get('Colonne TG Hire') in statuts:
                postes_data.append({
                    "titre": str(row.get('Poste demandé ', 'N/A')).strip(),
                    "entite": str(row.get('Entité demandeuse', 'N/A')).strip(),
                    "lieu": str(row.get('Affectation', 'N/A')).strip(),
                    "demandeur": str(row.get('Nom Prénom du demandeur', 'N/A')).strip(),
                    "recruteur": str(row.get('Responsable de traitement de  la demande ', 'N/A')).strip(),
                    "statut": str(row.get('Colonne TG Hire', 'N/A')),
                    "commentaire": str(row.get('Commentaire', '')) if pd.notna(row.get('Commentaire')) else None
                })
        
        # Si pas de données, utiliser un exemple
        if len(postes_data) == 0:
            postes_data = [
            {"titre": "Ingénieur Achat", "entite": "TGCC", "lieu": "SIEGE", "demandeur": "A.BOUZOUBAA", "recruteur": "Zakaria", "statut": "Sourcing"},
            {"titre": "Directeur Achats Adjoint", "entite": "TGCC", "lieu": "Siège", "demandeur": "C.BENABDELLAH", "recruteur": "Zakaria", "statut": "Sourcing"},
            {"titre": "INGENIEUR TRAVAUX", "entite": "TGCC", "lieu": "YAMED LOT B", "demandeur": "M.TAZI", "recruteur": "Zakaria", "statut": "Sourcing"},
            {"titre": "CHEF DE PROJETS", "entite": "TGCC", "lieu": "DESSALEMENT JORF", "demandeur": "M.FENNAN", "recruteur": "ZAKARIA", "statut": "Shortlisté"},
            {"titre": "Planificateur", "entite": "TGCC", "lieu": "ASFI-B", "demandeur": "SOUFIANI", "recruteur": "Ghita", "statut": "Shortlisté"},
            {"titre": "RESPONSABLE TRANS INTERCH", "entite": "TG PREFA", "lieu": "OUED SALEH", "demandeur": "FBOUZOUBAA", "recruteur": "Ghita", "statut": "Shortlisté"},
            {"titre": "PROJETEUR DESSINATEUR", "entite": "TG WOOD", "lieu": "OUED SALEH", "demandeur": "S.MENJRA", "recruteur": "Zakaria", "statut": "Signature DRH"},
            {"titre": "Projeteur", "entite": "TGCC", "lieu": "TSP Safi", "demandeur": "B.MORABET", "recruteur": "Zakaria", "statut": "Signature DRH"},
            {"titre": "Consultant SAP", "entite": "TGCC", "lieu": "Siège", "demandeur": "O.KETTA", "recruteur": "Zakaria", "statut": "Signature DRH"},
            {"titre": "Doc Controller", "entite": "TGEM", "lieu": "SIEGE", "demandeur": "A.SANKARI", "recruteur": "Zakaria", "statut": "Clôture"},
            {"titre": "Ingénieur étude/qualité", "entite": "TGCC", "lieu": "SIEGE", "demandeur": "A.MOUTANABI", "recruteur": "Zakaria", "statut": "Clôture"},
            {"titre": "Responsable Cybersecurité", "entite": "TGCC", "lieu": "Siège", "demandeur": "Ghazi", "recruteur": "Zakaria", "statut": "Clôture"},
            {"titre": "Conducteur de Travaux", "entite": "TGCC", "lieu": "JORF LASFAR", "demandeur": "M.FENNAN", "recruteur": "Zakaria", "statut": "Désistement"},
            {"titre": "Chef de Chantier", "entite": "TGCC", "lieu": "TOARC", "demandeur": "M.FENNAN", "recruteur": "Zakaria", "statut": "Désistement"},
            {"titre": "Magasinier", "entite": "TG WOOD", "lieu": "Oulad Saleh", "demandeur": "K.TAZI", "recruteur": "Ghita", "statut": "Désistement", "commentaire": "Pas de retour du demandeur"}
        ]
        
        statuts = ["Nouvelle demande", "Sourcing", "Shortlisté", "Signature DRH", "Clôture", "Désistement"]
        
        html_kanban = """
        <html>
        <head>
            <style>
                body {
                    margin: 0;
                    padding: 20px;
                    font-family: Arial, sans-serif;
                    background: #f5f5f5;
                }
                .kanban-container {
                    display: flex;
                    gap: 15px;
                    overflow-x: auto;
                    padding: 10px;
                }
                .kanban-column {
                    min-width: 260px;
                    background-color: #f0f0f0;
                    border-radius: 10px;
                    padding: 14px;
                    box-sizing: border-box;
                }
                .kanban-header {
                    background-color: #9C182F;
                    color: white;
                    padding: 12px;
                    border-radius: 6px;
                    font-weight: 800;
                    text-align: center;
                    margin-bottom: 12px;
                    font-size: 1.45em;
                }
                .kanban-card {
                    background-color: white;
                    border-radius: 6px;
                    padding: 14px;
                    margin-bottom: 12px;
                    box-shadow: 0 3px 6px rgba(0,0,0,0.1);
                    min-height: 120px;
                }
                .kanban-card h4 {
                    margin: 0 0 8px 0;
                    color: #9C182F;
                    white-space: normal !important;
                    word-break: break-all !important;
                    overflow-wrap: anywhere !important;
                    hyphens: auto !important;
                    display: -webkit-box;
                    -webkit-line-clamp: 3;
                    -webkit-box-orient: vertical;
                    overflow: hidden;
                }
                .kanban-card p {
                    margin: 4px 0;
                    font-size: 0.9em;
                }
            </style>
        </head>
        <body>
            <div class="kanban-container">
        """
        
        for statut in statuts:
            postes_in_col = [p for p in postes_data if p["statut"] == statut]
            html_kanban += f'<div class="kanban-column">'
            html_kanban += f'<div class="kanban-header">{statut} ({len(postes_in_col)})</div>'
            
            for poste in postes_in_col:
                titre_fmt = smart_wrap_title(poste.get('titre', ''))
                commentaire_html = f"<p>💬 {poste['commentaire']}</p>" if poste.get('commentaire') else ""
                html_kanban += f'''
                <div class="kanban-card">
                    <h4><b>{titre_fmt}</b></h4>
                    <p>📍 {poste.get('entite', 'N/A')} - {poste.get('lieu', 'N/A')}</p>
                    <p>👤 {poste.get('demandeur', 'N/A')}</p>
                    <p>✍️ {poste.get('recruteur', 'N/A')}</p>
                    {commentaire_html}
                </div>
                '''
            
            html_kanban += '</div>'
        
        html_kanban += """
            </div>
        </body>
        </html>
        """
        
        # Convertir en image avec html2image (haute résolution)
        image_path = hti.screenshot(html_str=html_kanban, save_as='kanban.png', size=(1920, 1080))[0]
        
        return image_path
    except Exception as e:
        st.error(f"Erreur lors de la génération de l'image du Kanban avec html2image: {e}")
        st.info("Tentative avec PIL...")
        # Fallback vers PIL
        return generate_kanban_image_simple(df_recrutement)


def _html_kpi_card(title, value, color="#1f77b4"):
        return f"""
        <div class='kpi-card'>
            <div class='kpi-title'>{title}</div>
            <div class='kpi-value' style='color:{color};'>{value}</div>
        </div>
        """

# --- Plotly image helpers for PPT composition ---
def _plotly_fig_to_pil(fig, width=900, height=520):
        """Render a Plotly figure to a PIL Image using kaleido. Falls back to PNG bytes in memory."""
        try:
            import plotly.io as pio
            import io
            img_bytes = pio.to_image(fig, format="png", width=width, height=height, engine="kaleido")
            from PIL import Image
            return Image.open(io.BytesIO(img_bytes)).convert('RGB')
        except Exception as e:
            st.warning(f"⚠️ Impossible d'exporter le graphique (kaleido manquant ?): {e}")
            return None

def _compose_dashboard_image(title_text, kpi_items, chart_rows, output_filename):
        """Compose a 1920x1080 dashboard image with a title, KPI row and chart rows.
        - title_text: string
        - kpi_items: list of tuples (label, value, color)
        - chart_rows: list of lists of PIL.Image objects (each inner list is a row)
        - output_filename: name to save in temp dir
        Returns absolute path to saved PNG.
        """
        import tempfile
        from PIL import Image, ImageDraw, ImageFont

        W, H = 1920, 1080
        M = 36  # outer margin
        img = Image.new('RGB', (W, H), 'white')
        draw = ImageDraw.Draw(img)
        try:
            font_title = ImageFont.truetype('/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf', 42)
            font_kpi_label = ImageFont.truetype('/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf', 22)
            font_kpi_value = ImageFont.truetype('/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf', 30)
        except Exception:
            font_title = ImageFont.load_default()
            font_kpi_label = ImageFont.load_default()
            font_kpi_value = ImageFont.load_default()

        # Title
        draw.text((M, M), title_text, fill="#9C182F", font=font_title)

        # KPI row layout
        kpi_top = M + 70
        kpi_height = 100
        kpi_gap = 16
        kpi_cols = max(1, len(kpi_items))
        kpi_col_width = int((W - 2*M - (kpi_gap * (kpi_cols - 1))) / kpi_cols)
        for i, (lab, val, col) in enumerate(kpi_items):
            x = M + i * (kpi_col_width + kpi_gap)
            y = kpi_top
            # card background
            draw.rounded_rectangle([x, y, x + kpi_col_width, y + kpi_height], radius=12, fill="#f8f9fa", outline="#e1e5ea")
            draw.line([(x+6, y), (x+6, y + kpi_height)], fill=col, width=6)
            draw.text((x + 16, y + 16), lab, fill="#334", font=font_kpi_label)
            draw.text((x + 16, y + 52), str(val), fill=col, font=font_kpi_value)

        # Charts layout
        current_y = kpi_top + kpi_height + 24
        row_gap = 24
        for row in chart_rows:
            if not row:
                continue
            # Determine widths per chart in row
            charts = [c for c in row if c is not None]
            if not charts:
                continue
            cols = len(charts)
            # leave margins left/right
            avail_w = W - 2*M
            gap = 16
            col_w = int((avail_w - gap * (cols - 1)) / cols)
            # uniform height per row
            row_h = 360
            x = M
            for cimg in charts:
                try:
                    # resize maintaining aspect
                    ci = cimg.copy()
                    ci.thumbnail((col_w, row_h))
                    img.paste(ci, (x, current_y))
                except Exception:
                    pass
                x += col_w + gap
            current_y += row_h + row_gap

        out_path = os.path.join(tempfile.gettempdir(), output_filename)
        img.save(out_path)
        return out_path

def generate_demandes_recrutement_html_image(df_recrutement):
        """Compose an image with the same charts as the Streamlit 'Demandes' tab."""
        df = df_recrutement.copy()
        # KPIs
        total = len(df)
        date_col = next((c for c in df.columns if "réception" in c.lower() and "date" in c.lower()), None)
        nouvelles = 0
        if date_col is not None:
            now = datetime.now(); start_of_month = now.replace(day=1)
            try:
                s = pd.to_datetime(df[date_col], errors='coerce')
                nouvelles = int((s >= start_of_month).sum())
            except Exception:
                nouvelles = 0
        annulees = 0
        if 'Statut de la demande' in df.columns:
            annulees = int(df['Statut de la demande'].astype(str).str.contains('annul|déprioris', case=False, na=False).sum())
        taux_annulation = f"{round((annulees/total)*100,1)}%" if total > 0 else "N/A"

        # Figures
        import plotly.express as px
        import plotly.graph_objects as go
        figs_row1 = []
        # pie statut
        try:
            if 'Statut de la demande' in df.columns:
                statut_counts = df['Statut de la demande'].value_counts()
                fig_statut = go.Figure(data=[go.Pie(labels=statut_counts.index, values=statut_counts.values, hole=.5)])
                fig_statut.update_traces(textfont=dict(size=14))
                fig_statut.update_layout(title=dict(text="Répartition par statut", x=0, xanchor='left', font=TITLE_FONT), height=360, margin=dict(l=20,r=20,t=48,b=12), legend=dict(font=dict(size=13)))
                figs_row1.append(fig_statut)
        except Exception:
            figs_row1.append(None)
        # bar raison
        try:
            if 'Raison du recrutement' in df.columns:
                raison_counts = df['Raison du recrutement'].value_counts()
                df_raison = raison_counts.rename_axis('Raison').reset_index(name='Count')
                fig_raison = px.bar(df_raison, x='Raison', y='Count', title="Raison du recrutement", text='Count')
                fig_raison.update_traces(marker_color='grey', textposition='outside')
                fig_raison.update_layout(height=320, margin=dict(l=20,r=20,t=40,b=10), xaxis_title=None, yaxis_title=None)
                figs_row1.append(fig_raison)
        except Exception:
            figs_row1.append(None)
        # evolution demandes monthly
        figs_row1.append(None)  # placeholder; we'll compute next
        if date_col in df.columns:
            try:
                df['Mois_Année_Demande'] = pd.to_datetime(df[date_col], errors='coerce').dt.to_period('M').dt.to_timestamp()
                monthly = df.groupby('Mois_Année_Demande').size().rename('Count')
                if not monthly.empty:
                    all_months = pd.date_range(start=monthly.index.min(), end=monthly.index.max(), freq='MS')
                    monthly = monthly.reindex(all_months, fill_value=0).reset_index().rename(columns={'index':'Mois'})
                    monthly['Label'] = monthly['Mois'].dt.strftime('%b %Y')
                    fig_evo = px.bar(monthly, x='Label', y='Count', title="Évolution des demandes", text='Count')
                    fig_evo.update_traces(marker_color='#1f77b4', textposition='outside')
                    fig_evo.update_layout(height=320, margin=dict(l=20,r=20,t=40,b=10), xaxis_title=None, yaxis_title=None)
                    figs_row1[-1] = fig_evo
            except Exception:
                pass

        # row2: direction and poste bars (horizontal)
        figs_row2 = []
        try:
            if 'Direction concernée' in df.columns:
                direction_counts = df['Direction concernée'].value_counts()
                dfd = direction_counts.rename_axis('Direction').reset_index(name='Count').sort_values('Count', ascending=False)
                dfd['Label'] = dfd['Direction']
                fig_dir = px.bar(dfd, x='Count', y='Label', title="Comparaison par direction", text='Count', orientation='h')
                fig_dir.update_traces(marker_color='grey', textposition='auto', texttemplate='%{x}', textfont=dict(size=15), hovertemplate='%{y}<extra></extra>')
                fig_dir.update_layout(height=320, margin=dict(l=160,t=48,b=30,r=20), xaxis_title=None, yaxis_title=None, yaxis=dict(tickfont=dict(size=13)))
                figs_row2.append(fig_dir)
        except Exception:
            figs_row2.append(None)
        try:
            if 'Poste demandé' in df.columns:
                poste_counts = df['Poste demandé'].value_counts()
                dfp = poste_counts.rename_axis('Poste').reset_index(name='Count').sort_values('Count', ascending=False)
                dfp['Label'] = dfp['Poste']
                fig_poste = px.bar(dfp, x='Count', y='Label', title="Comparaison par poste", text='Count', orientation='h')
                fig_poste.update_traces(marker_color='grey', textposition='auto', texttemplate='%{x}', textfont=dict(size=15), hovertemplate='%{y}<extra></extra>')
                fig_poste.update_layout(height=320, margin=dict(l=160,t=48,b=30,r=20), xaxis_title=None, yaxis_title=None, yaxis=dict(tickfont=dict(size=13)))
                figs_row2.append(fig_poste)
        except Exception:
            figs_row2.append(None)

        # Convert figures to PIL images
        row1_imgs = [(_plotly_fig_to_pil(f, width=600, height=320) if f else None) for f in figs_row1]
        row2_imgs = [(_plotly_fig_to_pil(f, width=900, height=360) if f else None) for f in figs_row2]

        # Compose final image
        kpis = [("Nombre de demandes", total, "#1f77b4"), ("Nouvelles (mois en cours)", nouvelles, "#2ca02c"), ("Annulées/Dépriorisées", annulees, "#ff7f0e"), ("Taux d'annulation", taux_annulation, "#d62728")]
        chart_rows = [row1_imgs, row2_imgs]
        return _compose_dashboard_image("📋 DEMANDES DE RECRUTEMENT", kpis, chart_rows, 'demandes_recrutement.png')

def generate_recrutements_clotures_html_image(df_recrutement):
    """Compose an image with charts matching the Streamlit 'Clôturés' tab."""
    df = df_recrutement.copy()
    df_cl = df[df['Statut de la demande'] == 'Clôture'] if 'Statut de la demande' in df.columns else df.iloc[0:0]
    nb = len(df_cl)
    postes_uniques = df_cl['Poste demandé'].nunique() if 'Poste demandé' in df_cl.columns else 0
    directions_uniques = df_cl['Direction concernée'].nunique() if 'Direction concernée' in df_cl.columns else 0
    delai_display = "N/A"
    rec_col = next((c for c in df_cl.columns if "réception" in c.lower() and "date" in c.lower()), None)
    ret_col = next((c for c in df_cl.columns if "retour" in c.lower() and "date" in c.lower()), None)
    if rec_col and ret_col:
        try:
            s = pd.to_datetime(df_cl[rec_col], errors='coerce'); e = pd.to_datetime(df_cl[ret_col], errors='coerce')
            mask = s.notna() & e.notna(); durees = (e[mask]-s[mask]).dt.days; durees = durees[durees >= 0]
            delai_display = str(round(durees.mean(),1)) if len(durees)>0 else "N/A"
        except Exception:
            delai_display = "N/A"

    import plotly.express as px
    import plotly.graph_objects as go
    # Row1: evolution monthly + modalité pie
    figs_row1 = []
    # Evolution
    try:
        mois_col = next((c for c in df_cl.columns if "entrée effective" in c.lower() or "date d'entrée" in c.lower()), None)
        if mois_col:
            s = pd.to_datetime(df_cl[mois_col], errors='coerce')
            series = s.dt.to_period('M').value_counts().sort_index()
            monthly = series.rename_axis('Mois').reset_index(name='Count')
            monthly['Label'] = monthly['Mois'].astype(str)
            fig_evo = px.bar(monthly, x='Label', y='Count', title="Évolution des recrutements", text='Count')
            fig_evo.update_traces(marker_color='#1f77b4', textposition='outside')
            fig_evo.update_layout(height=320, margin=dict(l=20,r=20,t=40,b=10), xaxis_title=None, yaxis_title=None)
            figs_row1.append(fig_evo)
    except Exception:
        figs_row1.append(None)
    # Modalité pie
    try:
        if 'Modalité de recrutement' in df_cl.columns:
            modalite_data = df_cl['Modalité de recrutement'].value_counts()
            fig_mod = go.Figure(data=[go.Pie(labels=modalite_data.index, values=modalite_data.values, hole=.5, textposition='inside', textinfo='percent')])
            fig_mod.update_traces(textfont=dict(size=14))
            fig_mod.update_layout(title=dict(text="Répartition par Modalité", x=0, xanchor='left', font=TITLE_FONT), height=360, margin=dict(l=20,r=20,t=48,b=12), legend=dict(font=dict(size=13)))
            figs_row1.append(fig_mod)
    except Exception:
        figs_row1.append(None)

    # Row2: direction and poste
    figs_row2 = []
    try:
        if 'Direction concernée' in df_cl.columns:
            direction_counts = df_cl['Direction concernée'].value_counts()
            dfd = direction_counts.rename_axis('Direction').reset_index(name='Count').sort_values('Count', ascending=False)
            dfd['Label'] = dfd['Direction']
            fig_dir = px.bar(dfd, x='Count', y='Label', title="Comparaison par direction", text='Count', orientation='h')
            fig_dir.update_traces(marker_color='#ff7f0e', textposition='inside', texttemplate='%{x}')
            fig_dir.update_layout(height=max(320, 24*len(dfd)), margin=dict(l=160,t=40,b=30,r=20), xaxis_title=None, yaxis_title=None)
            figs_row2.append(fig_dir)
    except Exception:
        figs_row2.append(None)
    try:
        if 'Poste demandé' in df_cl.columns:
            poste_counts = df_cl['Poste demandé'].value_counts()
            dfp = poste_counts.rename_axis('Poste').reset_index(name='Count').sort_values('Count', ascending=False)
            dfp['Label'] = dfp['Poste']
            fig_poste = px.bar(dfp, x='Count', y='Label', title="Comparaison par poste", text='Count', orientation='h')
            fig_poste.update_traces(marker_color='#2ca02c', textposition='inside', texttemplate='%{x}')
            fig_poste.update_layout(height=max(320, 24*len(dfp)), margin=dict(l=160,t=40,b=30,r=20), xaxis_title=None, yaxis_title=None)
            figs_row2.append(fig_poste)
    except Exception:
        figs_row2.append(None)

    row1_imgs = [(_plotly_fig_to_pil(f, width=900, height=320) if f else None) for f in figs_row1]
    row2_imgs = [(_plotly_fig_to_pil(f, width=900, height=360) if f else None) for f in figs_row2]

    kpis = [("Nombre de recrutements", nb, "#1f77b4"), ("Postes concernés", postes_uniques, "#2ca02c"), ("Directions concernées", directions_uniques, "#ff7f0e"), ("Délai moyen (jours)", delai_display, "#d62728")]
    chart_rows = [row1_imgs, row2_imgs]
    return _compose_dashboard_image("🎯 RECRUTEMENTS CLÔTURÉS", kpis, chart_rows, 'recrutements_clotures.png')

def generate_integrations_html_image(df_recrutement):
    """Compose an image with charts matching the Streamlit 'Intégrations' tab."""
    df = df_recrutement.copy()
    candidat_col = "Nom Prénom du candidat retenu yant accepté la promesse d'embauche"
    date_integration_col = "Date d'entrée prévisionnelle"
    plan_integration_col = "Plan d'intégration à préparer"
    if 'Statut de la demande' in df.columns:
        df = df[df['Statut de la demande'] == 'En cours']
    if candidat_col in df.columns:
        df = df[(df[candidat_col].notna()) & (df[candidat_col].astype(str).str.strip() != "")]

    nb_int = len(df)
    a_preparer = 0
    if plan_integration_col in df.columns:
        try:
            a_preparer = int((df[plan_integration_col].astype(str).str.lower() == 'oui').sum())
        except Exception:
            a_preparer = 0
    en_retard = 0
    if date_integration_col in df.columns:
        try:
            s = pd.to_datetime(df[date_integration_col], errors='coerce')
            reporting_date = st.session_state.get('reporting_date', datetime.now().date())
            today = reporting_date if isinstance(reporting_date, datetime) else datetime.combine(reporting_date, datetime.min.time())
            en_retard = int(((s.notna()) & (s < today)).sum())
        except Exception:
            en_retard = 0

    import plotly.express as px
    # Row1: affectation pie + monthly bar
    figs_row1 = []
    try:
        if 'Affectation' in df.columns:
            affect_counts = df['Affectation'].value_counts().head(10)
            fig_aff = px.pie(values=affect_counts.values, names=affect_counts.index, title="Répartition par Affectation")
            fig_aff.update_traces(textposition='inside', textinfo='percent+label', textfont=dict(size=14))
            fig_aff.update_layout(title=dict(text="Répartition par Affectation", x=0, xanchor='left', font=TITLE_FONT), height=420, margin=dict(l=20,r=20,t=48,b=12), legend=dict(font=dict(size=14)))
            figs_row1.append(fig_aff)
    except Exception:
        figs_row1.append(None)
    try:
        if date_integration_col in df.columns:
            s = pd.to_datetime(df[date_integration_col], errors='coerce')
            monthly = s.dt.to_period('M').value_counts().sort_index().rename_axis('Mois').reset_index(name='Count')
            monthly['Label'] = monthly['Mois'].astype(str)
            fig_month = px.bar(monthly, x='Label', y='Count', title="Évolution des Intégrations Prévues", text='Count')
            fig_month.update_traces(marker_color='#2ca02c', textposition='outside')
            fig_month.update_layout(height=360, margin=dict(l=20,r=20,t=40,b=10), xaxis_title=None, yaxis_title=None)
            figs_row1.append(fig_month)
    except Exception:
        figs_row1.append(None)

    row1_imgs = [(_plotly_fig_to_pil(f, width=900, height=360) if f else None) for f in figs_row1]
    kpis = [("Intégrations en cours", nb_int, "#1f77b4"), ("Plan d'intégration à préparer", a_preparer, "#ff7f0e"), ("En retard", en_retard, "#d62728")]
    chart_rows = [row1_imgs]
    return _compose_dashboard_image("📊 Intégrations", kpis, chart_rows, 'integrations.png')
def generate_powerpoint_report(df_recrutement, template_path="MASQUE PPT TGCC (2).pptx"):
    """
    Génère un rapport PowerPoint à partir d'un template en remplaçant les placeholders
    par les graphiques et tableaux générés.
    
    Returns: BytesIO contenant le fichier PowerPoint généré
    """
    try:
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN
        from pptx.dml.color import RGBColor
        from pptx.shapes.base import BaseShape
        
        # Charger le template
        prs = Presentation(template_path)
        
        # Préparer les données pour les placeholders
        # Calculer les métriques hebdomadaires
        weekly_metrics = calculate_weekly_metrics(df_recrutement)
        
        # Debug: afficher les informations sur weekly_metrics
        if weekly_metrics:
            st.info(f"✅ Métriques hebdomadaires calculées: {len(weekly_metrics.get('table_data', []))} lignes")
        else:
            st.warning("⚠️ Métriques hebdomadaires vides, le PowerPoint sera généré avec des données limitées.")
        
        # Générer les images HTML avec logos (utilise html2image + Chromium)
        st.info("📊 Génération des images avec les visualisations Streamlit...")
        table_image_path = generate_table_html_image(weekly_metrics) if weekly_metrics else None
        
        # Obtenir la date de reporting pour le filtre Clôture/Désistement
        from datetime import datetime, timedelta
        reporting_date = st.session_state.get('reporting_date', datetime.now().date())
        if isinstance(reporting_date, str):
            reporting_date = pd.to_datetime(reporting_date).date()
        
        # Calculer la fenêtre de dates (semaine précédente + courante)
        current_week_monday = datetime(year=reporting_date.year, month=reporting_date.month, day=reporting_date.day) - timedelta(days=reporting_date.weekday())
        start_filter = current_week_monday - timedelta(days=7)  # Lundi précédent
        end_filter = current_week_monday + timedelta(days=6)    # Dimanche courant
        
        st.info(f"📅 Période de filtrage Clôture/Désistement: {start_filter.date()} au {end_filter.date()}")
        
        # Générer une image par statut Kanban
        statuts = ["Nouvelle demande", "Sourcing", "Shortlisté", "Signature DRH", "Clôture", "Désistement"]
        kanban_images = {}
        
        for statut in statuts:
            st.info(f"🔄 Génération de l'image pour {statut}...")
            
            # Pour Clôture et Désistement, appliquer le filtre de date
            if statut == "Clôture":
                # Trouver la colonne de date d'acceptation
                accept_col = None
                for col_name in df_recrutement.columns:
                    if "date" in col_name.lower() and ("accept" in col_name.lower() or "promesse" in col_name.lower()):
                        accept_col = col_name
                        break
                
                if accept_col:
                    df_filtered = df_recrutement[
                        (df_recrutement['Colonne TG Hire'] == 'Clôture')
                    ].copy()
                    # Filtrer par date
                    df_filtered[accept_col] = pd.to_datetime(df_filtered[accept_col], errors='coerce')
                    df_filtered = df_filtered[
                        (df_filtered[accept_col].dt.date >= start_filter.date()) &
                        (df_filtered[accept_col].dt.date <= end_filter.date())
                    ]
                    st.info(f"📅 Clôture filtrée: {len(df_filtered)} postes (colonne: {accept_col})")
                else:
                    df_filtered = df_recrutement[df_recrutement['Colonne TG Hire'] == 'Clôture'].copy()
                    st.warning(f"⚠️ Colonne date d'acceptation non trouvée, affichage de tous les {len(df_filtered)} postes clôturés")
                
                kanban_images[statut] = generate_kanban_statut_image(df_filtered, statut, max_cards=100)
                
            elif statut == "Désistement":
                # Trouver la colonne de date de désistement
                desist_col = None
                for col_name in df_recrutement.columns:
                    if "date" in col_name.lower() and "désist" in col_name.lower():
                        desist_col = col_name
                        break
                
                if desist_col:
                    df_filtered = df_recrutement[
                        (df_recrutement['Colonne TG Hire'] == 'Désistement')
                    ].copy()
                    # Filtrer par date
                    df_filtered[desist_col] = pd.to_datetime(df_filtered[desist_col], errors='coerce')
                    df_filtered = df_filtered[
                        (df_filtered[desist_col].dt.date >= start_filter.date()) &
                        (df_filtered[desist_col].dt.date <= end_filter.date())
                    ]
                    st.info(f"📅 Désistement filtré: {len(df_filtered)} postes (colonne: {desist_col})")
                else:
                    df_filtered = df_recrutement[df_recrutement['Colonne TG Hire'] == 'Désistement'].copy()
                    st.warning(f"⚠️ Colonne date désistement non trouvée, affichage de tous les {len(df_filtered)} désistements")
                # Toujours générer une image même si aucune donnée
                img_path_desist = generate_kanban_statut_image(df_filtered, statut, max_cards=100)
                if not img_path_desist or not os.path.exists(img_path_desist):
                    # Placeholder image
                    try:
                        from PIL import Image, ImageDraw, ImageFont
                        import tempfile
                        width, height = 1920, 1080
                        img = Image.new('RGB',(width,height),'white')
                        d = ImageDraw.Draw(img)
                        try:
                            f_title = ImageFont.truetype('/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf', 42)
                            f_text = ImageFont.truetype('/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf', 20)
                        except:
                            f_title = ImageFont.load_default(); f_text = ImageFont.load_default()
                        d.text((width//2, 80), 'Désistement (0 poste)', fill='#9C182F', font=f_title, anchor='mm')
                        d.text((width//2, height//2), 'Aucune donnée pour cette période', fill='#666666', font=f_text, anchor='mm')
                        ph_path = os.path.join(tempfile.gettempdir(),'desistement_empty.png')
                        img.save(ph_path)
                        kanban_images[statut] = ph_path
                    except Exception:
                        kanban_images[statut] = img_path_desist
                else:
                    kanban_images[statut] = img_path_desist
            else:
                # Pour les autres statuts, pas de filtre de date
                kanban_images[statut] = generate_kanban_statut_image(df_recrutement, statut, max_cards=100)
        
        # Debug: Vérifier les chemins des images
        if table_image_path and os.path.exists(table_image_path):
            st.success(f"✅ Image tableau générée: {os.path.basename(table_image_path)}")
        else:
            st.warning("⚠️ Échec génération image tableau")
        
        # Vérifier les images Kanban
        for statut, img_path in kanban_images.items():
            if img_path and os.path.exists(img_path):
                st.success(f"✅ Image {statut} générée: {os.path.basename(img_path)}")
            else:
                st.warning(f"⚠️ Échec génération image {statut}")

        # Générer les deux slides supplémentaires: Demandes de Recrutement et Recrutements Clôturés
        st.info("🧩 Génération des images pour DEMANDES DE RECRUTEMENT et RECRUTEMENTS CLÔTURÉS...")
        demandes_img = generate_demandes_recrutement_html_image(df_recrutement)
        clotures_img = generate_recrutements_clotures_html_image(df_recrutement)
        
        # Parcourir chaque slide et remplacer les placeholders
        for slide_idx, slide in enumerate(prs.slides):
            st.info(f"📄 Traitement de la slide {slide_idx + 1}/{len(prs.slides)}")
            shapes_to_remove = []
            
            # Debug: lister tous les shapes avec texte
            text_shapes = []
            for shape in slide.shapes:
                shape_text = getattr(shape, "text", None)
                if shape_text and shape_text.strip():
                    text_shapes.append(shape_text[:100])  # Premier 100 chars
            
            if text_shapes:
                st.code(f"Slide {slide_idx + 1} - Shapes avec texte:\n" + "\n".join([f"  - {t}" for t in text_shapes]))
            
            for shape in slide.shapes:
                # Utiliser getattr pour éviter les erreurs Pylance
                shape_text = getattr(shape, "text", None)
                if shape_text is None:
                    continue
                    
                # Tableau des besoins par entité - Insérer l'image
                if "{{TABLEAU_BESOINS_ENTITES}}" in shape_text:
                    st.info(f"🔍 Placeholder {{{{TABLEAU_BESOINS_ENTITES}}}} trouvé dans slide {slide_idx + 1}")
                    try:
                        # Récupérer la position et taille du shape original
                        left = shape.left
                        top = shape.top
                        width = shape.width
                        height = shape.height
                        
                        st.info(f"📐 Position Tableau: left={left}, top={top}, width={width}, height={height}")
                        
                        # Marquer le shape pour suppression
                        shapes_to_remove.append(shape)
                        
                        # Insérer l'image du tableau si elle existe (en grand format)
                        if table_image_path and os.path.exists(table_image_path):
                            # Utiliser toute la largeur de la slide pour une meilleure visibilité
                            slide_width = prs.slide_width
                            slide_height = prs.slide_height
                            # Positionner l'image en laissant des marges
                            img_left = Inches(0.5)
                            img_top = Inches(1.5)
                            # Calculer dimensions avec vérification
                            if slide_width and slide_height:
                                img_width = Inches((slide_width.inches if hasattr(slide_width, 'inches') else slide_width / 914400) - 1)
                                img_height = Inches((slide_height.inches if hasattr(slide_height, 'inches') else slide_height / 914400) - 2.5)
                                slide.shapes.add_picture(table_image_path, img_left, img_top, width=img_width, height=img_height)
                            else:
                                slide.shapes.add_picture(table_image_path, img_left, img_top)
                            st.success(f"✅ Image tableau insérée dans slide {slide_idx + 1} (dimensions optimisées)")
                        else:
                            st.error(f"❌ Image tableau non trouvée: {table_image_path}")
                            # Fallback: ajouter un message d'erreur
                            txBox = slide.shapes.add_textbox(left, top, width, height)
                            text_frame = txBox.text_frame
                            p = text_frame.paragraphs[0]
                            p.text = "Erreur: Impossible de générer l'image du tableau"
                            p.font.size = Pt(14)
                    except Exception as e:
                        st.error(f"❌ Erreur lors de l'insertion du tableau: {e}")
                        import traceback
                        st.code(traceback.format_exc())
                
                # Métrique total postes - Créer 5 slides Kanban (une par statut) dans l'ordre souhaité
                elif "{{METRIC_TOTAL_POSTES}}" in shape_text:
                    st.info(f"🔍 Placeholder {{{{METRIC_TOTAL_POSTES}}}} trouvé dans slide {slide_idx + 1}")
                    st.info("📋 Création de 5 slides Kanban (une par statut) dans l'ordre: Sourcing, Shortlisté, Signature DRH, Clôture, Désistement...")
                    try:
                        # Marquer le shape pour suppression
                        shapes_to_remove.append(shape)
                        
                        # Pour chaque statut, insérer l'image dans cette slide ou créer de nouvelles slides
                        # On va insérer toutes les images Kanban après cette slide
                        # Pour l'instant, on insère juste la première dans la slide actuelle
                        
                        # Récupérer les dimensions de slide pour positionnement
                        slide_width = prs.slide_width
                        slide_height = prs.slide_height
                        img_left = Inches(0.5)
                        img_top = Inches(1.5)
                        
                        # Calculer dimensions avec vérification
                        if slide_width and slide_height:
                            img_width = Inches((slide_width.inches if hasattr(slide_width, 'inches') else slide_width / 914400) - 1)
                            img_height = Inches((slide_height.inches if hasattr(slide_height, 'inches') else slide_height / 914400) - 2.5)
                        else:
                            img_width = Inches(9)
                            img_height = Inches(6)
                        
                        # Insérer les 5 images Kanban dans des slides séparées selon l'ordre demandé
                        ordered_status = ["Nouvelle demande", "Sourcing", "Shortlisté", "Signature DRH", "Clôture", "Désistement"]
                        for statut in ordered_status:
                            img_path = kanban_images.get(statut)
                            # Créer une nouvelle slide pour chaque statut (ne pas utiliser la slide actuelle)
                            new_slide = prs.slides.add_slide(slide.slide_layout)
                            # Titre de la slide: nom du statut
                            if hasattr(new_slide.shapes, 'title') and new_slide.shapes.title:
                                new_slide.shapes.title.text = f"Kanban - {statut}"
                            # Insérer l'image
                            if img_path and os.path.exists(img_path):
                                new_slide.shapes.add_picture(img_path, img_left, img_top, width=img_width, height=img_height)
                                st.success(f"✅ Slide insérée pour {statut}")
                            else:
                                st.warning(f"⚠️ Image {statut} non trouvée")

                        # Ajouter juste après Désistement: DEMANDES DE RECRUTEMENT puis RECRUTEMENTS CLÔTURÉS, puis INTÉGRATIONS
                        if demandes_img and os.path.exists(demandes_img):
                            slide_dem = prs.slides.add_slide(slide.slide_layout)
                            if hasattr(slide_dem.shapes, 'title') and slide_dem.shapes.title:
                                slide_dem.shapes.title.text = "📋 DEMANDES DE RECRUTEMENT"
                            slide_dem.shapes.add_picture(demandes_img, img_left, img_top, width=img_width, height=img_height)
                            st.success("✅ Slide ajoutée: DEMANDES DE RECRUTEMENT")
                        else:
                            st.warning("⚠️ Image DEMANDES DE RECRUTEMENT non trouvée")

                        if clotures_img and os.path.exists(clotures_img):
                            slide_clo = prs.slides.add_slide(slide.slide_layout)
                            if hasattr(slide_clo.shapes, 'title') and slide_clo.shapes.title:
                                slide_clo.shapes.title.text = "🎯 RECRUTEMENTS CLÔTURÉS"
                            slide_clo.shapes.add_picture(clotures_img, img_left, img_top, width=img_width, height=img_height)
                            st.success("✅ Slide ajoutée: RECRUTEMENTS CLÔTURÉS")
                        else:
                            st.warning("⚠️ Image RECRUTEMENTS CLÔTURÉS non trouvée")

                        # Intégrations
                        integrations_img = generate_integrations_html_image(df_recrutement)
                        if integrations_img and os.path.exists(integrations_img):
                            slide_int = prs.slides.add_slide(slide.slide_layout)
                            if hasattr(slide_int.shapes, 'title') and slide_int.shapes.title:
                                slide_int.shapes.title.text = "📊 Intégrations"
                            slide_int.shapes.add_picture(integrations_img, img_left, img_top, width=img_width, height=img_height)
                            st.success("✅ Slide ajoutée: Intégrations")
                        else:
                            st.warning("⚠️ Image Intégrations non trouvée")
                        
                    except Exception as e:
                        st.error(f"❌ Erreur lors de l'insertion des Kanban: {e}")
                        import traceback
                        st.code(traceback.format_exc())
                
                # Graphiques Plotly - on peut ajouter d'autres placeholders ici
                elif "{{GRAPH_" in shape_text:
                    placeholder = shape_text.strip()
                    text_frame = getattr(shape, "text_frame", None)
                    if text_frame:
                        text_frame.clear()
                        p = text_frame.paragraphs[0] if text_frame.paragraphs else text_frame.add_paragraph()
                        p.text = f"Graphique {placeholder} (à venir)"
                        p.font.size = Pt(12)
            
            # Supprimer les shapes marqués (après la boucle pour éviter les problèmes d'itération)
            for shape in shapes_to_remove:
                sp = shape.element
                sp.getparent().remove(sp)
        
        # Nettoyer les fichiers temporaires
        if table_image_path and os.path.exists(table_image_path):
            try:
                os.remove(table_image_path)
            except:
                pass
        
        # Nettoyer les images Kanban temporaires
        for img_path in kanban_images.values():
            if img_path and os.path.exists(img_path):
                try:
                    os.remove(img_path)
                except:
                    pass

        # Nettoyer les images supplémentaires
        for extra in [demandes_img, clotures_img, integrations_img]:
            try:
                if extra and os.path.exists(extra):
                    os.remove(extra)
            except:
                pass
        
        # Sauvegarder le PowerPoint modifié dans un BytesIO
        ppt_bytes = BytesIO()
        prs.save(ppt_bytes)
        ppt_bytes.seek(0)
        
        return ppt_bytes
    
    except Exception as e:
        st.error(f"Erreur lors de la génération du PowerPoint: {e}")
        return None


def main():
    st.title("📊 Tableau de Bord RH")
    # Afficher le hash du commit actuel pour aider au debug des déploiements
    try:
        commit_hash = get_current_commit_hash()
        st.markdown(f"<div style='font-size:12px;color:#666;margin-top:-8px'>Commit: {commit_hash}</div>", unsafe_allow_html=True)
    except Exception:
        pass
    st.markdown("---")
    # Date de reporting : permet de fixer la date de référence pour tous les calculs
    if 'reporting_date' not in st.session_state:
        # default to today's date
        st.session_state['reporting_date'] = datetime.now().date()

    with st.sidebar:
        st.subheader("🔧 Filtres - Hebdomadaire")
        # Ensure the session value is a plain date object Streamlit expects.
        # Protect against cases where session_state may contain a str or pd.Timestamp
        rd_val = st.session_state.get('reporting_date', None)
        if rd_val is not None:
            try:
                # pandas Timestamp or datetime -> date
                import pandas as _pd
                if isinstance(rd_val, _pd.Timestamp):
                    st.session_state['reporting_date'] = rd_val.date()
                elif hasattr(rd_val, 'date') and not isinstance(rd_val, str):
                    # datetime -> date (leave date as-is)
                    try:
                        st.session_state['reporting_date'] = rd_val.date()
                    except Exception:
                        pass
                elif isinstance(rd_val, str):
                    # try parsing common formats
                    from dateutil import parser as _parser
                    try:
                        parsed = _parser.parse(rd_val)
                        st.session_state['reporting_date'] = parsed.date()
                    except Exception:
                        # fallback: keep original (Streamlit will raise if invalid)
                        pass
            except Exception:
                # If pandas isn't available or any unexpected issue, ignore and let Streamlit handle
                pass
        # Do not pass both `value` and `key` for the same widget to avoid Streamlit's
        # warning: the widget was created with a default value but also had its
        # value set via the Session State API. Passing `key='reporting_date'`
        # is enough for Streamlit to initialize the widget from
        # `st.session_state['reporting_date']`. Avoid providing `value=` here.
        rd = st.date_input(
            "Date de reporting (hebdomadaire)",
            help="Date utilisée comme référence pour le reporting hebdomadaire",
            key='reporting_date',
            format='DD/MM/YYYY'
        )
    
    # Créer les onglets (Demandes et Recrutement regroupés)
    # CSS pour agrandir le texte des onglets
    st.markdown("""
    <style>
    .stTabs [data-baseweb="tab-list"] button {
        font-size: 1.15em !important;
        font-weight: 600 !important;
        padding: 12px 20px !important;
    }
    .stTabs [data-baseweb="tab-list"] button [data-testid="stMarkdownContainer"] p {
        font-size: 1.15em !important;
        font-weight: 600 !important;
    }
    </style>
    """, unsafe_allow_html=True)
    tabs = st.tabs(["📂 Upload & Téléchargement", "🗂️ Demandes & Recrutement", "📅 Hebdomadaire", "🤝 Intégrations", "📖 Méthodologie"])
    
    # Variables pour stocker les fichiers uploadés
    # Use session_state to persist upload/refresh state
    if 'data_updated' not in st.session_state:
        st.session_state.data_updated = False
    if 'uploaded_excel' not in st.session_state:
        st.session_state.uploaded_excel = None
    uploaded_excel = st.session_state.uploaded_excel
    
    with tabs[0]:
        col1, col2 = st.columns(2)

        with col1:
            st.subheader("🔗 Synchroniser depuis Google Sheets")
            # Respecter la date de reporting si fournie dans st.session_state
            reporting_date = st.session_state.get('reporting_date', None)
            if reporting_date is None:
                today = datetime.now()
            else:
                # reporting_date is a date object from st.date_input; convert to datetime
                if isinstance(reporting_date, datetime):
                    today = reporting_date
                else:
                    today = datetime.combine(reporting_date, datetime.min.time())
            start_of_month = today.replace(day=1)
            # valeur par défaut du Google Sheet (identique à celle utilisée précédemment)
            default_sheet = "https://docs.google.com/spreadsheets/d/1hvghSMjcbdY8yNZOWqALBpgMdLWB5CxVJCDwEm6JULI/edit?gid=785271056#gid=785271056"
            gs_url = st.text_input("URL Google Sheet", value=default_sheet, key="gsheet_url")
            
            if 'synced_recrutement_df' not in st.session_state:
                st.session_state.synced_recrutement_df = None
            
            if st.button("🔁 Synchroniser depuis Google Sheets", 
                        help="Synchroniser les données depuis Google Sheets",
                        width="stretch"):
                
                try:
                    # Utiliser la fonction de connexion automatique (comme dans Home.py)
                    df_synced = load_data_from_google_sheets(gs_url)
                    
                    if df_synced is not None and len(df_synced) > 0:
                        st.session_state.synced_recrutement_df = df_synced
                        st.session_state.data_updated = True
                        nb_lignes = len(df_synced)
                        nb_colonnes = len(df_synced.columns)
                        st.success(f"✅ Synchronisation Google Sheets réussie ! Les onglets ont été mis à jour. ({nb_lignes} lignes, {nb_colonnes} colonnes)")
                    else:
                        st.warning("⚠️ Aucune donnée trouvée dans la feuille Google Sheets.")
                        
                except Exception as e:
                    err_str = str(e)
                    st.error(f"Erreur lors de la synchronisation: {err_str}")
                    
                    if '401' in err_str or 'Unauthorized' in err_str or 'HTTP Error 401' in err_str:
                        st.error("❌ **Feuille Google privée** - Vérifiez que:")
                        st.markdown("""
                        1. La feuille est partagée avec: `your-service-account@your-project.iam.gserviceaccount.com`
                        2. Les secrets Streamlit sont correctement configurés
                        3. L'URL de la feuille est correcte
                        """)
                    elif 'secrets' in err_str.lower():
                        st.error("❌ **Configuration des secrets manquante**")
                        st.markdown("""
                        Assurez-vous que les secrets suivants sont configurés:
                        - `GCP_TYPE`, `GCP_PROJECT_ID`, `GCP_PRIVATE_KEY_ID`
                        - `GCP_PRIVATE_KEY`, `GCP_CLIENT_EMAIL`, `GCP_CLIENT_ID`
                        - `GCP_AUTH_URI`, `GCP_TOKEN_URI`, etc.
                        """)
                    else:
                        st.error(f"Erreur technique: {err_str}")

        with col2:
            st.subheader("📊 Fichier Excel - Données de Recrutement")
            uploaded_excel = st.file_uploader(
                "Choisir le fichier Excel de recrutement",
                type=['xlsx', 'xls'],
                help="Fichier Excel contenant les données de recrutement",
                key="excel_uploader"
            )
            
            if uploaded_excel is not None:
                # Aperçu des données
                try:
                    preview_excel = pd.read_excel(uploaded_excel, sheet_name=0)
                    st.success(f"✅ Fichier Excel chargé: {uploaded_excel.name} - {len(preview_excel)} lignes, {len(preview_excel.columns)} colonnes")
                    st.dataframe(preview_excel.head(3), width="stretch")
                    # Reset file pointer for later use
                    uploaded_excel.seek(0)
                    st.session_state.uploaded_excel = uploaded_excel
                except Exception as e:
                    st.error(f"Erreur lors de la lecture de l'Excel: {e}")
        
        # Bouton pour actualiser les données - s'étale sur les deux colonnes
        st.markdown("---")
        if st.button("🔄 Actualiser les Graphiques", type="primary", width="stretch"):
            st.session_state.data_updated = True
            st.success("Données mises à jour ! Consultez les autres onglets.")
    
    # Charger les données ICI (avant l'onglet Upload pour pouvoir les utiliser)
    df_integration, df_recrutement = load_data_from_files(None, uploaded_excel)
    
    # Continuer l'onglet Upload avec la section PowerPoint
    with tabs[0]:
        # Section de génération du PowerPoint
        st.markdown("---")
        st.subheader("📥 Télécharger le Rapport PowerPoint")
        
        col_ppt1, col_ppt2 = st.columns(2)
        
        with col_ppt1:
            st.info("📊 **Génération automatique de votre rapport PowerPoint**")
            st.markdown("""
            Le rapport PowerPoint inclura :
            - 📊 Tableau des besoins par entité (Hebdomadaire)
            - 📈 Métriques et statistiques
            - 📉 Graphiques de Demandes & Recrutement
            - 🤝 Données d'Intégration
            """)
        
        with col_ppt2:
            # Vérifier que des données sont disponibles
            if df_recrutement is not None or st.session_state.get('synced_recrutement_df') is not None:
                # Option de format d'export
                export_format = st.radio(
                    "Format d'export:",
                    ["PowerPoint (.pptx)", "PDF (.pdf)"],
                    horizontal=True,
                    help="PDF recommandé pour une meilleure compatibilité avec Office 2010"
                )
                
                if st.button("📥 Générer et Télécharger le Rapport", type="primary", width="stretch"):
                    with st.spinner("Génération du rapport en cours..."):
                        # Utiliser les données synchronisées si disponibles
                        data_to_use = st.session_state.get('synced_recrutement_df')
                        if data_to_use is None:
                            data_to_use = df_recrutement
                        
                        # Générer le PowerPoint
                        ppt_bytes = generate_powerpoint_report(data_to_use)
                        
                        if ppt_bytes:
                            # Générer un nom de fichier avec la date
                            today_str = datetime.now().strftime("%Y-%m-%d")
                            
                            if export_format == "PDF (.pdf)":
                                # Convertir en PDF
                                try:
                                    from pptx import Presentation
                                    from pptx.enum.shapes import MSO_SHAPE_TYPE
                                    from reportlab.lib.pagesizes import A4, landscape
                                    from reportlab.pdfgen import canvas
                                    from reportlab.lib.utils import ImageReader
                                    from PIL import Image
                                    import io
                                    import tempfile
                                    
                                    st.info("📄 Conversion en PDF...")
                                    
                                    # Charger le PPT depuis les bytes
                                    ppt_bytes.seek(0)
                                    prs = Presentation(ppt_bytes)
                                    
                                    # Créer un PDF
                                    pdf_buffer = io.BytesIO()
                                    c = canvas.Canvas(pdf_buffer, pagesize=landscape(A4))
                                    page_width, page_height = landscape(A4)
                                    
                                    # Pour chaque slide, extraire les images et les ajouter au PDF
                                    for slide_idx, slide in enumerate(prs.slides):
                                        # Ajouter un titre de page
                                        c.setFont("Helvetica-Bold", 16)
                                        c.drawString(50, page_height - 40, f"Slide {slide_idx + 1}")
                                        
                                        # Extraire les images de la slide
                                        img_y = page_height - 80
                                        for shape in slide.shapes:
                                            try:
                                                # Extraire uniquement les images (pictures)
                                                if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.PICTURE:
                                                    img_obj = getattr(shape, "image", None)
                                                    if img_obj is None:
                                                        continue
                                                    img_data = img_obj.blob
                                                    img = Image.open(io.BytesIO(img_data))
                                                    
                                                    # Redimensionner pour tenir dans la page
                                                    max_width = page_width - 100
                                                    max_height = page_height - 150
                                                    img_ratio = img.width / img.height
                                                    
                                                    if img.width > max_width:
                                                        new_width = max_width
                                                        new_height = new_width / img_ratio
                                                    else:
                                                        new_width = img.width
                                                        new_height = img.height
                                                    
                                                    if new_height > max_height:
                                                        new_height = max_height
                                                        new_width = new_height * img_ratio
                                                    
                                                    # Sauvegarder temporairement
                                                    temp_img = io.BytesIO()
                                                    img.save(temp_img, format='PNG')
                                                    temp_img.seek(0)
                                                    
                                                    c.drawImage(ImageReader(temp_img), 50, img_y - new_height, 
                                                               width=new_width, height=new_height)
                                                    img_y -= new_height + 20
                                            except Exception as img_e:
                                                st.warning(f"Image non extraite: {img_e}")
                                        
                                        c.showPage()
                                    
                                    c.save()
                                    pdf_buffer.seek(0)
                                    
                                    st.download_button(
                                        label="💾 Télécharger le Rapport PDF",
                                        data=pdf_buffer,
                                        file_name=f"Rapport_RH_{today_str}.pdf",
                                        mime="application/pdf",
                                        type="primary",
                                        width="stretch"
                                    )
                                    st.success("✅ PDF généré avec succès !")
                                    
                                except ImportError:
                                    st.error("❌ Module 'reportlab' non installé. Installation: pip install reportlab")
                                    st.info("💡 En attendant, voici le PowerPoint:")
                                    ppt_bytes.seek(0)
                                    st.download_button(
                                        label="💾 Télécharger le Rapport PowerPoint",
                                        data=ppt_bytes,
                                        file_name=f"Rapport_RH_{today_str}.pptx",
                                        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                                        type="secondary",
                                        width="stretch"
                                    )
                                except Exception as pdf_e:
                                    st.error(f"❌ Erreur lors de la conversion PDF: {pdf_e}")
                                    st.info("💡 Téléchargez le PowerPoint à la place:")
                                    ppt_bytes.seek(0)
                                    st.download_button(
                                        label="💾 Télécharger le Rapport PowerPoint",
                                        data=ppt_bytes,
                                        file_name=f"Rapport_RH_{today_str}.pptx",
                                        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                                        type="secondary",
                                        width="stretch"
                                    )
                            else:
                                # Export PowerPoint standard
                                filename = f"Rapport_RH_{today_str}.pptx"
                                
                                st.download_button(
                                    label="💾 Télécharger le Rapport PowerPoint",
                                    data=ppt_bytes,
                                    file_name=filename,
                                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                                    type="primary",
                                    width="stretch"
                                )
                                st.success("✅ PowerPoint généré avec succès !")
            else:
                st.warning("⚠️ Veuillez d'abord charger des données (Google Sheets ou Excel) avant de générer le PowerPoint.")
    
    # Message d'information sur les données chargées
    has_uploaded = (st.session_state.uploaded_excel is not None) or (st.session_state.get('synced_recrutement_df') is not None)
    if df_recrutement is None and df_integration is None:
        st.sidebar.warning("⚠️ Aucune donnée disponible. Veuillez uploader vos fichiers dans l'onglet 'Upload Fichiers'.")
    elif df_recrutement is None:
        st.sidebar.warning("⚠️ Données de recrutement non disponibles. Seules les données d'intégration sont chargées.")
    elif df_integration is None:
        st.sidebar.warning("⚠️ Données d'intégration non disponibles. Seules les données de recrutement sont chargées.")

    with tabs[1]:
        if df_recrutement is not None:
            create_demandes_recrutement_combined_tab(df_recrutement)
        else:
            st.warning("📊 Aucune donnée de recrutement disponible. Veuillez uploader un fichier Excel dans l'onglet 'Upload Fichiers'.")
    
    with tabs[2]:
        create_weekly_report_tab(df_recrutement)

    with tabs[3]:
        # Onglet Intégrations basé sur les données Excel
        if df_recrutement is not None:
            # Créer les filtres spécifiques pour les intégrations (sans période)
            st.sidebar.subheader("🔧 Filtres - Intégrations")
            int_filters = create_integration_filters(df_recrutement, "integrations")
            create_integrations_tab(df_recrutement, int_filters)
        else:
            st.warning("📊 Aucune donnée disponible pour les intégrations. Veuillez uploader un fichier Excel dans l'onglet 'Upload Fichiers'.")

    with tabs[4]:
        st.header("📖 Méthodologie & Guide Utilisateur")
        
        st.subheader("Comment générer le reporting ?")
        st.markdown("""
        Avant de générer le reporting, assurez-vous que tous les recrutements sont saisis, que la colonne «Status» de la demande est renseignée et que la colonne «TG Hire» reflète l'état du recrutement (Nouvelle demande, Sourcing, Signature DRH, ...).
        """)

        st.markdown('<div style="color:red; font-weight:700; margin-top:6px;">NB : au fur et à mesure de l\'évolution des recrutements, mettez à jour leur état, sinon le reporting ne sera pas fiable.</div>', unsafe_allow_html=True)

        st.markdown("""
        1.  **Chargement des Données** : 
            *   Allez dans l'onglet **"📂 Upload & Téléchargement"**.
            *   **Option A** : Cliquez sur le bouton rouge **"🔁 Synchroniser depuis Google Sheets"** pour récupérer les données les plus récentes.
            *   **Option B** : Glissez-déposez votre fichier Excel de recrutement dans la zone de chargement.
        2.  **Actualisation** : 
            *   Cliquez impérativement sur le bouton bleu **"🔄 Actualiser les Graphiques"**. Cela déclenche le calcul de tous les indicateurs et la mise à jour des visuels.
        3.  **Configuration & Filtres** :
            *   Utilisez le **menu à gauche (sidebar)** pour sélectionner la **"Date de reporting"**. 
            *   Le reporting hebdomadaire prend en compte la performance de la semaine précédente (période du **Vendredi au Lundi S-1**).
            *   Affinez vos analyses grâce aux filtres par **Entité**, **Direction** ou **Année**.
        4.  **Consultation et Export** :
            *   Naviguez dans les onglets (**Demandes**, **Hebdo**, **Intégrations**) pour visualiser les résultats.
            *   Pour finir, prenez des captures d'écran des graphiques et tableaux pour les insérer dans le PowerPoint.
            *   Le reporting des stagaires se prend pour le moment avec Power BI (Lien : https://drive.google.com/file/d/1BF3JNIq11O9FGNzN428r3JiEl3qqYWID/view?usp=drive_link)

        💡 *Une fonction permettant d'automatiser totalement la génération périodique du reporting est actuellement en cours de développement (les captures d'écran des graphiques et tableaux y seront automatiquement intégrées).*
        """)

        st.markdown("---")
        st.subheader("🔍 Détail des Indicateurs & Méthodologie")
        st.markdown("""
        Cette section explique comment chaque indicateur est calculé pour garantir la fiabilité des données :

        ### 📂 Onglet Demandes & Recrutement
        - **Besoins en cours par entité** : Calculés à partir des demandes validées par la DRH. Un poste est considéré "en cours" s'il a le statut `En cours`. Si les données sont incomplètes, le système utilise la balance `(Stock Initial + Nouveaux - Pourvus)`.
        - **Recrutements en cours par recruteur** : Tableau croisé dynamique basé sur la `Colonne TG Hire`. Le **"Total (hors clôture)"** additionne les étapes *Sourcing*, *Shortlisté* et *Signature DRH* pour mesurer la charge active.
        - **Comparaison par direction / poste** : Histogrammes de volume basés sur les colonnes `Direction concernée` et `Poste demandé`. Les labels sont tronqués pour la lisibilité mais le nom complet est visible au survol.
        - **Évolution mensuelle** : Agrégation temporelle basée sur la `Date de réception de la demande`.

        ### 📅 Onglet Reporting Hebdomadaire
        - **Fenêtre de calcul** : Basée sur la `Date de reporting` sélectionnée dans le menu latéral. Elle analyse la performance du **Lundi au Vendredi** de la semaine précédente.
        - **Postes Avant** : Demandes ouvertes avant le lundi de la semaine S-1.
        - **Nouveaux** : Demandes reçues entre le lundi et le vendredi de la semaine S-1.
        - **Pourvus** : Recrutements dont la date d'acceptation du candidat tombe dans la semaine S-1.
        - **Stock (En cours)** : Postes restant ouverts (Statut 'En cours' + Pas de candidat retenu) à la fin de la période.

        ### 🤝 Onglet Intégrations
        - **Critères d'Inclusion** : Seules les demandes avec le statut `En cours` ET un `Nom de candidat` renseigné sont affichées ici.
        - **Signal "⚠️ En retard"** : Une alerte rouge s'affiche si la `Date d'entrée prévisionnelle` est antérieure à la date du jour (ou la date de reporting choisie).
        - **Répartition par Entité** : Visualisation en "Donut Chart" pour identifier les filiales avec le plus gros volume d'arrivées imminentes.

        ### 🎯 Indicateurs de Performance (KPIs)
        - **Nombre de candidats présélectionnés** : Somme cumulative de la colonne `Nb de candidats pré-selectionnés`.
        - **Délai de recrutement (Duree de recrutement)** : Calculé selon la formule : `DATEDIFF('Date de réception de la demande','Date du 1er retour equipe RH', day)`. Cet indicateur se base **uniquement** sur les postes ayant le statut **"Clôture"**.
        - **Taux de refus** : Ratio entre le nombre de refus et le nombre de promesses d'embauche réalisées.
        """, unsafe_allow_html=True)

if __name__ == "__main__":
    main()